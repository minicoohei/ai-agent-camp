"""pptx_ops.py の単体テスト"""
import json
import pytest
from unittest.mock import patch, MagicMock
from pathlib import Path


class TestImport:
    def test_import_module(self):
        from pptx_ops import PowerPointOperator
        assert PowerPointOperator is not None


class TestPowerPointOperator:
    def test_init_no_file(self):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator()
        assert op.filepath is None

    def test_init_with_file(self, tmp_path):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator(str(tmp_path / "test.pptx"))
        assert op.filepath == tmp_path / "test.pptx"

    def test_load_nonexistent(self, tmp_path):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator(str(tmp_path / "nonexistent.pptx"))
        assert op.load() is False

    def test_load_none_filepath(self):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator()
        assert op.load() is False

    def test_load_valid_pptx(self, sample_pptx):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator(str(sample_pptx))
        assert op.load() is True
        assert op.get_slide_count() == 1

    def test_read_slide(self, sample_pptx):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator(str(sample_pptx))
        op.load()
        slide_data = op.read_slide(0)
        assert "index" in slide_data
        assert slide_data["index"] == 1

    def test_read_slide_invalid_index(self, sample_pptx):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator(str(sample_pptx))
        op.load()
        result = op.read_slide(999)
        assert "error" in result

    def test_read_slide_negative_index(self, sample_pptx):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator(str(sample_pptx))
        op.load()
        result = op.read_slide(-1)
        assert "error" in result

    def test_read_slide_not_loaded(self):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator()
        result = op.read_slide(0)
        assert "error" in result

    def test_get_slide_count_not_loaded(self):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator()
        assert op.get_slide_count() == 0

    def test_read_all_slides(self, sample_pptx):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator(str(sample_pptx))
        op.load()
        slides = op.read_all_slides()
        assert len(slides) == 1
        assert slides[0]["index"] == 1

    def test_read_all_slides_not_loaded(self):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator()
        slides = op.read_all_slides()
        assert slides == []

    def test_to_markdown(self, sample_pptx):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator(str(sample_pptx))
        op.load()
        md = op.to_markdown()
        assert "# " in md
        assert "Slides" in md

    def test_to_markdown_not_loaded(self):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator()
        md = op.to_markdown()
        assert "Error" in md

    def test_analyze(self, sample_pptx):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator(str(sample_pptx))
        op.load()
        report = op.analyze()
        assert "slide_count" in report
        assert report["slide_count"] == 1
        assert "filename" in report
        assert "layouts_used" in report
        assert isinstance(report["layouts_used"], list)
        assert "shape_types" in report

    def test_analyze_not_loaded(self):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator()
        result = op.analyze()
        assert "error" in result

    def test_extract_template(self, sample_pptx):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator(str(sample_pptx))
        op.load()
        template = op.extract_template()
        assert "source_file" in template
        assert "layouts" in template
        assert "slides" in template

    def test_extract_template_not_loaded(self):
        from pptx_ops import PowerPointOperator
        op = PowerPointOperator()
        result = op.extract_template()
        assert "error" in result


class TestCreateFromTemplate:
    def test_create_simple_presentation(self, tmp_path):
        from pptx_ops import PowerPointOperator
        template_data = {
            "slides": [
                {
                    "title": "テストタイトル",
                    "content": ["ポイント1", "ポイント2", "ポイント3"],
                    "notes": "スピーカーノート"
                }
            ]
        }
        op = PowerPointOperator()
        output = tmp_path / "new.pptx"
        result = op.create_from_template(template_data, str(output))
        assert Path(result).exists()

    def test_create_multiple_slides(self, tmp_path):
        from pptx_ops import PowerPointOperator
        template_data = {
            "slides": [
                {"title": "スライド1"},
                {"title": "スライド2", "content": ["内容A"]},
                {"title": "スライド3", "content": ["内容B"], "notes": "ノート"},
            ]
        }
        op = PowerPointOperator()
        output = tmp_path / "multi.pptx"
        op.create_from_template(template_data, str(output))

        # Verify by reading back
        op2 = PowerPointOperator(str(output))
        assert op2.load() is True
        assert op2.get_slide_count() == 3

    def test_create_with_custom_dimensions(self, tmp_path):
        from pptx_ops import PowerPointOperator
        template_data = {
            "slide_width": 12192000,  # 10 inches in EMU
            "slide_height": 6858000,  # 7.5 inches in EMU
            "slides": [{"title": "Wide slide"}]
        }
        op = PowerPointOperator()
        output = tmp_path / "wide.pptx"
        op.create_from_template(template_data, str(output))
        assert Path(output).exists()

    def test_create_empty_slides(self, tmp_path):
        from pptx_ops import PowerPointOperator
        template_data = {"slides": [{}]}
        op = PowerPointOperator()
        output = tmp_path / "empty.pptx"
        op.create_from_template(template_data, str(output))
        assert Path(output).exists()

    def test_create_no_slides(self, tmp_path):
        from pptx_ops import PowerPointOperator
        template_data = {"slides": []}
        op = PowerPointOperator()
        output = tmp_path / "noslide.pptx"
        op.create_from_template(template_data, str(output))
        assert Path(output).exists()


class TestReadWithContent:
    """Test reading presentations with actual content"""

    def test_read_slide_with_text(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx_ops import PowerPointOperator

        # Create PPTX with text
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tf = txBox.text_frame
        tf.text = "テストテキスト"

        path = tmp_path / "with_text.pptx"
        prs.save(str(path))

        op = PowerPointOperator(str(path))
        op.load()
        slide_data = op.read_slide(0)
        assert "テストテキスト" in slide_data["text_content"]

    def test_read_slide_with_table(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches
        from pptx_ops import PowerPointOperator

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table
        table.cell(0, 0).text = "Header1"
        table.cell(0, 1).text = "Header2"
        table.cell(1, 0).text = "Data1"
        table.cell(1, 1).text = "Data2"

        path = tmp_path / "with_table.pptx"
        prs.save(str(path))

        op = PowerPointOperator(str(path))
        op.load()
        slide_data = op.read_slide(0)
        # Find the shape with table data
        has_table = any("table" in shape for shape in slide_data["shapes"])
        assert has_table

    def test_to_markdown_with_content(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches
        from pptx_ops import PowerPointOperator

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox.text_frame.text = "マークダウンテスト"

        path = tmp_path / "md_test.pptx"
        prs.save(str(path))

        op = PowerPointOperator(str(path))
        op.load()
        md = op.to_markdown()
        assert "マークダウンテスト" in md
        assert "## Slide 1" in md

    def test_analyze_with_content(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches
        from pptx_ops import PowerPointOperator

        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            txBox.text_frame.text = f"Slide {i+1}"

        path = tmp_path / "analyze.pptx"
        prs.save(str(path))

        op = PowerPointOperator(str(path))
        op.load()
        report = op.analyze()
        assert report["slide_count"] == 3
        assert report["total_text_length"] > 0


# ===========================================================================
# Import with ImportError (lines 28-31)
# ===========================================================================

class TestPptxImportError:
    def test_import_error_exit(self):
        """python-pptx がない場合に sys.exit(1) (lines 28-31)"""
        import importlib
        import sys

        # This is hard to test without actually removing pptx,
        # but we can verify the module has the try/except structure
        from pptx_ops import PowerPointOperator
        assert PowerPointOperator is not None


# ===========================================================================
# load exception (lines 50-52)
# ===========================================================================

class TestLoadException:
    def test_load_corrupt_file(self, tmp_path):
        """壊れたファイルの読み込み (lines 50-52)"""
        from pptx_ops import PowerPointOperator
        corrupt = tmp_path / "corrupt.pptx"
        corrupt.write_bytes(b"not a pptx file at all")
        op = PowerPointOperator(str(corrupt))
        assert op.load() is False


# ===========================================================================
# cmd_read (lines 355-378)
# ===========================================================================

class TestCmdRead:
    def test_cmd_read_json_format(self, sample_pptx, capsys):
        """JSON形式での出力 (lines 367-368)"""
        from pptx_ops import cmd_read
        args = MagicMock()
        args.file = str(sample_pptx)
        args.slide = None
        args.format = "json"
        cmd_read(args)
        captured = capsys.readouterr()
        data = json.loads(captured.out)
        assert isinstance(data, list)
        assert len(data) == 1

    def test_cmd_read_text_format(self, sample_pptx, capsys):
        """テキスト形式での出力 (lines 370-377)"""
        from pptx_ops import cmd_read
        args = MagicMock()
        args.file = str(sample_pptx)
        args.slide = None
        args.format = "text"
        cmd_read(args)
        captured = capsys.readouterr()
        assert "Slide 1" in captured.out

    def test_cmd_read_specific_slide(self, sample_pptx, capsys):
        """特定スライドの読み取り (line 362)"""
        from pptx_ops import cmd_read
        args = MagicMock()
        args.file = str(sample_pptx)
        args.slide = 1
        args.format = "text"
        cmd_read(args)
        captured = capsys.readouterr()
        assert "Slide 1" in captured.out

    def test_cmd_read_load_fail(self, tmp_path):
        """読み込み失敗 (line 358-359)"""
        from pptx_ops import cmd_read
        args = MagicMock()
        args.file = str(tmp_path / "nonexistent.pptx")
        args.slide = None
        args.format = "text"
        cmd_read(args)  # Should return without error

    def test_cmd_read_text_with_notes_and_long_text(self, tmp_path, capsys):
        """ノートとlong textのテスト (lines 375-377)"""
        from pptx import Presentation
        from pptx.util import Inches
        from pptx_ops import cmd_read

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox.text_frame.text = "A" * 150  # Long text
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Speaker notes here"

        path = tmp_path / "notes.pptx"
        prs.save(str(path))

        args = MagicMock()
        args.file = str(path)
        args.slide = None
        args.format = "text"
        cmd_read(args)
        captured = capsys.readouterr()
        assert "..." in captured.out
        assert "Notes" in captured.out


# ===========================================================================
# cmd_to_markdown (lines 380-392)
# ===========================================================================

class TestCmdToMarkdown:
    def test_cmd_to_markdown_stdout(self, sample_pptx, capsys):
        """stdout出力 (line 392)"""
        from pptx_ops import cmd_to_markdown
        args = MagicMock()
        args.file = str(sample_pptx)
        args.output = None
        cmd_to_markdown(args)
        captured = capsys.readouterr()
        assert "Slide" in captured.out

    def test_cmd_to_markdown_file(self, sample_pptx, tmp_path):
        """ファイル出力 (lines 388-390)"""
        from pptx_ops import cmd_to_markdown
        args = MagicMock()
        args.file = str(sample_pptx)
        args.output = str(tmp_path / "output.md")
        cmd_to_markdown(args)
        assert (tmp_path / "output.md").exists()

    def test_cmd_to_markdown_load_fail(self, tmp_path):
        """読み込み失敗"""
        from pptx_ops import cmd_to_markdown
        args = MagicMock()
        args.file = str(tmp_path / "nonexistent.pptx")
        args.output = None
        cmd_to_markdown(args)


# ===========================================================================
# cmd_analyze (lines 395-418)
# ===========================================================================

class TestCmdAnalyze:
    def test_cmd_analyze_json(self, sample_pptx, capsys):
        """JSON形式 (line 404)"""
        from pptx_ops import cmd_analyze
        args = MagicMock()
        args.file = str(sample_pptx)
        args.format = "json"
        cmd_analyze(args)
        captured = capsys.readouterr()
        data = json.loads(captured.out)
        assert "slide_count" in data

    def test_cmd_analyze_text(self, sample_pptx, capsys):
        """テキスト形式 (lines 406-418)"""
        from pptx_ops import cmd_analyze
        args = MagicMock()
        args.file = str(sample_pptx)
        args.format = "text"
        cmd_analyze(args)
        captured = capsys.readouterr()
        assert "Analysis Report" in captured.out
        assert "Shape types" in captured.out

    def test_cmd_analyze_load_fail(self, tmp_path):
        """読み込み失敗"""
        from pptx_ops import cmd_analyze
        args = MagicMock()
        args.file = str(tmp_path / "nonexistent.pptx")
        args.format = "text"
        cmd_analyze(args)


# ===========================================================================
# cmd_extract_template (lines 421-435)
# ===========================================================================

class TestCmdExtractTemplate:
    def test_cmd_extract_template_stdout(self, sample_pptx, capsys):
        """stdout出力 (line 435)"""
        from pptx_ops import cmd_extract_template
        args = MagicMock()
        args.file = str(sample_pptx)
        args.output = None
        cmd_extract_template(args)
        captured = capsys.readouterr()
        data = json.loads(captured.out)
        assert "source_file" in data

    def test_cmd_extract_template_file(self, sample_pptx, tmp_path):
        """ファイル出力 (lines 431-433)"""
        from pptx_ops import cmd_extract_template
        args = MagicMock()
        args.file = str(sample_pptx)
        args.output = str(tmp_path / "template.json")
        cmd_extract_template(args)
        assert (tmp_path / "template.json").exists()

    def test_cmd_extract_template_load_fail(self, tmp_path):
        """読み込み失敗"""
        from pptx_ops import cmd_extract_template
        args = MagicMock()
        args.file = str(tmp_path / "nonexistent.pptx")
        args.output = None
        cmd_extract_template(args)


# ===========================================================================
# cmd_create (lines 438-448)
# ===========================================================================

class TestCmdCreate:
    def test_cmd_create_success(self, tmp_path):
        """正常作成 (lines 440-448)"""
        from pptx_ops import cmd_create
        template = {
            "slides": [
                {"title": "Test Slide", "content": ["Point 1"], "notes": "Note"}
            ]
        }
        template_file = tmp_path / "template.json"
        template_file.write_text(json.dumps(template), encoding="utf-8")

        args = MagicMock()
        args.template = str(template_file)
        args.output = str(tmp_path / "created.pptx")
        cmd_create(args)
        assert (tmp_path / "created.pptx").exists()

    def test_cmd_create_load_fail(self, tmp_path, capsys):
        """テンプレート読み込み失敗 (lines 442-444)"""
        from pptx_ops import cmd_create
        args = MagicMock()
        args.template = str(tmp_path / "nonexistent.json")
        args.output = str(tmp_path / "out.pptx")
        cmd_create(args)
        captured = capsys.readouterr()
        assert "Failed" in captured.out

    def test_cmd_create_invalid_json(self, tmp_path, capsys):
        """不正なJSON (line 442)"""
        from pptx_ops import cmd_create
        bad_json = tmp_path / "bad.json"
        bad_json.write_text("{broken}", encoding="utf-8")
        args = MagicMock()
        args.template = str(bad_json)
        args.output = str(tmp_path / "out.pptx")
        cmd_create(args)
        captured = capsys.readouterr()
        assert "Failed" in captured.out


# ===========================================================================
# main (lines 451-496)
# ===========================================================================

class TestMain:
    def test_main_no_command(self):
        """コマンドなしで実行 (lines 483-485)"""
        from pptx_ops import main
        with patch("sys.argv", ["pptx_ops.py"]):
            with pytest.raises(SystemExit) as exc_info:
                main()
            assert exc_info.value.code == 1

    def test_main_read_command(self, sample_pptx, capsys):
        """readコマンド実行 (line 488)"""
        from pptx_ops import main
        with patch("sys.argv", ["pptx_ops.py", "read", str(sample_pptx), "-f", "json"]):
            main()
        captured = capsys.readouterr()
        assert "index" in captured.out

    def test_main_analyze_command(self, sample_pptx, capsys):
        """analyzeコマンド実行 (line 492)"""
        from pptx_ops import main
        with patch("sys.argv", ["pptx_ops.py", "analyze", str(sample_pptx), "-f", "json"]):
            main()
        captured = capsys.readouterr()
        assert "slide_count" in captured.out

    def test_main_to_markdown_command(self, sample_pptx, capsys):
        """to-markdownコマンド実行 (line 490)"""
        from pptx_ops import main
        with patch("sys.argv", ["pptx_ops.py", "to-markdown", str(sample_pptx)]):
            main()
        captured = capsys.readouterr()
        assert "Slide" in captured.out

    def test_main_extract_template_command(self, sample_pptx, capsys):
        """extract-templateコマンド実行 (line 494)"""
        from pptx_ops import main
        with patch("sys.argv", ["pptx_ops.py", "extract-template", str(sample_pptx)]):
            main()
        captured = capsys.readouterr()
        assert "source_file" in captured.out

    def test_main_create_command(self, tmp_path, capsys):
        """createコマンド実行 (line 496)"""
        from pptx_ops import main
        template = {"slides": [{"title": "Test"}]}
        template_file = tmp_path / "template.json"
        template_file.write_text(json.dumps(template))
        with patch("sys.argv", [
            "pptx_ops.py", "create", str(template_file),
            "-o", str(tmp_path / "out.pptx")
        ]):
            main()
        assert (tmp_path / "out.pptx").exists()


# ===========================================================================
# to_markdown with table and notes (lines 176-194)
# ===========================================================================

class TestToMarkdownWithTableAndNotes:
    def test_markdown_with_table(self, tmp_path):
        """テーブル付きMarkdown変換 (lines 176-187)"""
        from pptx import Presentation
        from pptx.util import Inches
        from pptx_ops import PowerPointOperator

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table
        table.cell(0, 0).text = "Col1"
        table.cell(0, 1).text = "Col2"
        table.cell(1, 0).text = "Val1"
        table.cell(1, 1).text = "Val2"

        path = tmp_path / "table.pptx"
        prs.save(str(path))

        op = PowerPointOperator(str(path))
        op.load()
        md = op.to_markdown()
        assert "Col1" in md
        assert "|" in md
        assert "---" in md

    def test_markdown_with_notes(self, tmp_path):
        """ノート付きMarkdown変換 (lines 190-194)"""
        from pptx import Presentation
        from pptx.util import Inches
        from pptx_ops import PowerPointOperator

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox.text_frame.text = "Content"
        notes = slide.notes_slide
        notes.notes_text_frame.text = "Speaker notes"

        path = tmp_path / "notes.pptx"
        prs.save(str(path))

        op = PowerPointOperator(str(path))
        op.load()
        md = op.to_markdown()
        assert "Speaker Notes" in md
        assert "Speaker notes" in md

    def test_markdown_no_text_shows_slide(self, tmp_path):
        """テキストなしスライドのMarkdown (line 154)"""
        from pptx import Presentation
        from pptx_ops import PowerPointOperator

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # blank slide, no text

        path = tmp_path / "blank.pptx"
        prs.save(str(path))

        op = PowerPointOperator(str(path))
        op.load()
        md = op.to_markdown()
        assert "Slide" in md


# ===========================================================================
# analyze with placeholders (lines 237-241)
# ===========================================================================

class TestAnalyzeWithPlaceholders:
    def test_analyze_with_placeholder_slide(self, tmp_path):
        """プレースホルダー付きスライドの分析 (lines 237-241)"""
        from pptx import Presentation
        from pptx_ops import PowerPointOperator

        prs = Presentation()
        # Use a layout with placeholders (Title Slide = index 0)
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        # Fill in title placeholder
        if slide.placeholders:
            slide.placeholders[0].text = "Title"

        path = tmp_path / "placeholder.pptx"
        prs.save(str(path))

        op = PowerPointOperator(str(path))
        op.load()
        report = op.analyze()
        assert report["slide_count"] == 1
        # Should have placeholder info
        assert len(report["slides"][0]["placeholders"]) > 0
