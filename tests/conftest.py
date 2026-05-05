"""
pytest fixtures for ai-agent-camp tests.
決定論的テスト用の共通フィクスチャ。
"""
import pytest
import os
import sys
import json
import importlib.util
from pathlib import Path
from unittest.mock import MagicMock


# ---------------------------------------------------------------------------
# Path helpers
# ---------------------------------------------------------------------------

@pytest.fixture
def project_root():
    """プロジェクトルートパスを返す"""
    return Path(__file__).parent.parent


def import_module_from_repo(module_name: str, relative_path: str | Path):
    """リポジトリ内の指定ファイルを module_name として毎回再ロードする。"""
    project_root = Path(__file__).parent.parent
    module_path = project_root / relative_path

    sys.modules.pop(module_name, None)
    spec = importlib.util.spec_from_file_location(module_name, module_path)
    if spec is None or spec.loader is None:
        raise ImportError(f"Could not load module spec for {module_path}")

    module = importlib.util.module_from_spec(spec)
    sys.modules[module_name] = module
    spec.loader.exec_module(module)
    return module


@pytest.fixture(autouse=True)
def _add_project_to_path(project_root):
    """tools/ と skills/ をsys.pathに追加（インポート可能にする）"""
    paths_to_add = [
        str(project_root),
        str(project_root / "tools"),
        str(project_root / "skills" / "check-inbox" / "scripts"),
        str(project_root / "skills" / "x-research" / "scripts"),
        str(project_root / "skills" / "pptx-converter" / "scripts"),
        str(project_root / "skills" / "pptx-creator" / "scripts"),
        str(project_root / "skills" / "monitoring-dashboard" / "scripts"),
    ]
    original = sys.path.copy()
    for p in paths_to_add:
        if p not in sys.path:
            sys.path.insert(0, p)
    yield
    sys.path[:] = original


# ---------------------------------------------------------------------------
# Sample data
# ---------------------------------------------------------------------------

@pytest.fixture
def sample_slack_data(project_root):
    """Slackサンプルデータを読み込む"""
    sample_path = project_root / "tests/fixtures/slack-sample-1.json"
    if sample_path.exists():
        with open(sample_path) as f:
            return json.load(f)
    return None


@pytest.fixture
def sample_gmail_data(project_root):
    """Gmailサンプルデータを読み込む"""
    sample_path = project_root / "tests/fixtures/gmail-sample-1.json"
    if sample_path.exists():
        with open(sample_path) as f:
            return json.load(f)
    return None


# ---------------------------------------------------------------------------
# Output / temp directories
# ---------------------------------------------------------------------------

@pytest.fixture
def output_dir(tmp_path):
    """テスト用出力ディレクトリ（自動作成）"""
    d = tmp_path / "output"
    d.mkdir()
    return d


# ---------------------------------------------------------------------------
# Environment variable mocks
# ---------------------------------------------------------------------------

@pytest.fixture
def mock_env(monkeypatch):
    """テスト用環境変数（全API キーをダミー値に設定）"""
    env_vars = {
        "GEMINI_API_KEY": "test_gemini_key_12345",
        "GOOGLE_CLOUD_PROJECT": "test-project",
        "GOOGLE_API_KEY": "test_google_api_key",
        "OPENAI_API_KEY": "test_openai_key",
        "ANTHROPIC_API_KEY": "test_anthropic_key",
        "SLACK_BOT_TOKEN": "test-slack-bot-token-dummy",
        "SLACK_USER_TOKEN": "test-slack-user-token-dummy",
        "X_BEARER_TOKEN": "test_x_bearer_token",
        "RESEND_API_KEY": "test_resend_key",
    }
    for key, value in env_vars.items():
        monkeypatch.setenv(key, value)
    return env_vars


@pytest.fixture
def clean_env(monkeypatch):
    """API キーが未設定の環境をシミュレート"""
    keys_to_remove = [
        "GEMINI_API_KEY", "GOOGLE_CLOUD_PROJECT", "GOOGLE_API_KEY",
        "OPENAI_API_KEY", "ANTHROPIC_API_KEY", "SLACK_BOT_TOKEN",
        "X_BEARER_TOKEN",
    ]
    for key in keys_to_remove:
        monkeypatch.delenv(key, raising=False)


# ---------------------------------------------------------------------------
# API mock helpers
# ---------------------------------------------------------------------------

@pytest.fixture
def mock_gemini_response():
    """Gemini API のモックレスポンスを返すファクトリ"""
    def _make(text="mock response", finish_reason="STOP"):
        resp = MagicMock()
        resp.text = text
        candidate = MagicMock()
        candidate.finish_reason = finish_reason
        candidate.content.parts = [MagicMock(text=text)]
        resp.candidates = [candidate]
        return resp
    return _make


@pytest.fixture
def mock_gemini_image_response():
    """Gemini Image Generation API のモックレスポンスを返すファクトリ"""
    def _make(image_bytes=b"\x89PNG\r\n\x1a\n" + b"\x00" * 100):
        resp = MagicMock()
        part = MagicMock()
        part.inline_data.data = image_bytes
        part.inline_data.mime_type = "image/png"
        resp.candidates = [MagicMock()]
        resp.candidates[0].content.parts = [part]
        return resp
    return _make


@pytest.fixture
def mock_slack_client():
    """Slack WebClient のモック"""
    client = MagicMock()
    client.conversations_list.return_value = {
        "ok": True,
        "channels": [{"id": "C123", "name": "general"}],
    }
    client.conversations_history.return_value = {
        "ok": True,
        "messages": [{"text": "hello", "ts": "1234567890.000000"}],
    }
    client.chat_postMessage.return_value = {"ok": True, "ts": "1234567890.000001"}
    return client


# ---------------------------------------------------------------------------
# File helpers
# ---------------------------------------------------------------------------

@pytest.fixture
def sample_image(tmp_path):
    """1x1 PNG テスト画像を作成"""
    try:
        from PIL import Image
        img = Image.new("RGB", (100, 100), color="red")
        path = tmp_path / "test_image.png"
        img.save(path)
        return path
    except ImportError:
        pytest.skip("Pillow not installed")


@pytest.fixture
def sample_pptx(tmp_path):
    """空のPPTXファイルを作成"""
    try:
        from pptx import Presentation
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = tmp_path / "test.pptx"
        prs.save(path)
        return path
    except ImportError:
        pytest.skip("python-pptx not installed")


@pytest.fixture
def sample_csv(tmp_path):
    """テスト用CSVファイルを作成"""
    path = tmp_path / "test_data.csv"
    path.write_text("name,age,city\nAlice,30,Tokyo\nBob,25,Osaka\nCharlie,35,Nagoya\n", encoding="utf-8")
    return path
