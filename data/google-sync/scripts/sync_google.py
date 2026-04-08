#!/usr/bin/env python3
"""
Google Sync
gogcli を使って Gmail/Drive/Docs/Sheets/Slides/Calendar を同期し、
Markdown/CSV 形式で保存する。
"""

from __future__ import annotations

import argparse
import base64
import csv
import json
import os
import re
import subprocess
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional

import openpyxl
from pptx import Presentation


BASE_DIR = Path(__file__).resolve().parents[1]
DATA_DIR = BASE_DIR / "data"


def safe_name(name: str, max_len: int = 80) -> str:
    cleaned = name.strip()
    cleaned = re.sub(r"[\\/:*?\"<>|]", "_", cleaned)
    cleaned = re.sub(r"\s+", "_", cleaned)
    cleaned = re.sub(r"_+", "_", cleaned)
    if not cleaned:
        cleaned = "untitled"
    return cleaned[:max_len]


def run_cmd(cmd: List[str]) -> tuple[bool, str]:
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode == 0:
        return True, result.stdout
    return False, (result.stderr or result.stdout)


def run_gog_json(account: str, *args: str) -> tuple[bool, Dict[str, Any]]:
    cmd = ["gog", "--no-input", "--account", account, "--json", *args]
    ok, output = run_cmd(cmd)
    if not ok:
        return False, {"error": output.strip()}
    try:
        return True, json.loads(output)
    except json.JSONDecodeError as exc:
        return False, {"error": f"json_decode_failed: {exc}"}


def run_gog(account: str, *args: str) -> tuple[bool, str]:
    cmd = ["gog", "--no-input", "--account", account, *args]
    return run_cmd(cmd)


def decode_base64url(data: str) -> str:
    if not data:
        return ""
    padding = "=" * (-len(data) % 4)
    try:
        raw = base64.urlsafe_b64decode(data + padding)
        return raw.decode("utf-8", errors="replace")
    except Exception:
        return ""


def load_accounts() -> List[Dict[str, str]]:
    config = os.getenv("GOG_ACCOUNTS_CONFIG", "").strip()
    accounts: List[Dict[str, str]] = []
    if config:
        data = json.loads(config)
        for entry in data:
            account = entry.get("account")
            if not account:
                continue
            label = entry.get("label") or safe_name(account.split("@")[0])
            accounts.append(
                {
                    "label": label,
                    "account": account,
                    "drive_folder": entry.get("drive_folder", ""),
                }
            )
        return accounts
    account = os.getenv("GOG_ACCOUNT", "").strip()
    if account:
        accounts.append(
            {
                "label": "default",
                "account": account,
                "drive_folder": os.getenv("GOG_DRIVE_FOLDER_ID", "").strip(),
            }
        )
    return accounts


def extract_plain_text(payload: Dict[str, Any]) -> str:
    body = payload.get("body", {})
    if body.get("data"):
        return decode_base64url(body["data"])
    for part in payload.get("parts", []) or []:
        mime = part.get("mimeType")
        if mime == "text/plain":
            return decode_base64url(part.get("body", {}).get("data", ""))
        text = extract_plain_text(part)
        if text:
            return text
    return ""


def write_text(path: Path, lines: Iterable[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))


def write_markdown_from_text(title: str, content: str, out_path: Path) -> None:
    lines = [
        f"# {title}",
        "",
        f"_Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}_",
        "",
        content,
    ]
    write_text(out_path, lines)


def yaml_escape(value: str) -> str:
    """
    YAML フロントマター用に値をエスケープ
    特殊文字を含む場合はダブルクォートで囲む
    """
    if not value:
        return '""'
    # 特殊文字チェック: ", ', :, #, [, ], {, }, >, |, *, &, !, %, @, `, 改行
    special_chars = ['"', "'", ":", "#", "[", "]", "{", "}", ">", "|", "*", "&", "!", "%", "@", "`", "\n"]
    needs_quote = any(c in value for c in special_chars) or value.startswith(" ") or value.endswith(" ")
    
    if needs_quote:
        # ダブルクォート内のダブルクォートはエスケープ
        escaped = value.replace("\\", "\\\\").replace('"', '\\"')
        return f'"{escaped}"'
    return value


def gmail_thread_to_markdown(thread: Dict[str, Any], debug: bool = False) -> tuple[List[str], str]:
    """
    Gmail スレッドを YAML フロントマター形式の Markdown に変換
    
    email_parser.py と互換性のある形式で出力:
    ---
    id: thread_id
    subject: ...
    from: ...
    date: ...
    ---
    本文
    
    Returns:
        (lines, filename_base): lines は Markdown 行のリスト、
        filename_base は "2026-02-05_メールの件名" 形式
    """
    thread_id = thread.get("id", "unknown")
    
    # gogcli の出力構造に対応: messages は直接または thread キー内にある可能性
    messages = thread.get("messages", []) or []
    if not messages and "thread" in thread:
        messages = thread["thread"].get("messages", []) or []
    
    if debug:
        print(f"    DEBUG: thread_data keys = {list(thread.keys())}")
        print(f"    DEBUG: messages count = {len(messages)}")
        if messages:
            print(f"    DEBUG: first msg keys = {list(messages[0].keys())}")
    
    if not messages:
        return [], ""
    
    # 最新メッセージからメタデータ取得
    latest = messages[-1]
    headers = {
        h.get("name", "").lower(): h.get("value", "")
        for h in latest.get("payload", {}).get("headers", [])
    }
    
    subject = headers.get("subject", "No subject")
    sender = headers.get("from", "Unknown")
    date_str = headers.get("date", "")
    
    # internalDate があればそれを使用（より正確）
    internal = latest.get("internalDate")
    date_prefix = "unknown"
    if internal:
        try:
            dt = datetime.fromtimestamp(int(internal) / 1000, tz=timezone.utc)
            date_str = dt.strftime("%a, %d %b %Y %H:%M:%S %z")
            date_prefix = dt.strftime("%Y-%m-%d")
        except Exception:
            pass
    
    # YAML フロントマター形式（特殊文字をエスケープ）
    lines = [
        "---",
        f"id: {thread_id}",
        f"subject: {yaml_escape(subject)}",
        f"from: {yaml_escape(sender)}",
        f"date: {yaml_escape(date_str)}",
        "---",
        "",
    ]
    
    # 本文（全メッセージを結合）
    for msg in messages:
        body = extract_plain_text(msg.get("payload", {})).strip()
        snippet = msg.get("snippet", "")
        content = body or snippet or ""
        if content:
            lines.append(content)
            lines.append("")
    
    # ファイル名: 日付_件名.md 形式
    safe_subject = safe_name(subject, max_len=50)
    filename_base = f"{date_prefix}_{safe_subject}"
    
    return lines, filename_base


def sync_gmail(account: str, label: str, days: int, max_threads: int) -> None:
    query = f"newer_than:{days}d"
    ok, data = run_gog_json(account, "gmail", "search", query, "--max", str(max_threads))
    if not ok:
        print(f"  ⚠️ Gmail search failed: {data.get('error')}")
        return
    threads = data.get("threads", []) or []
    gmail_dir = DATA_DIR / label / "gmail"
    
    # デバッグ: 最初のスレッドだけ詳細ログ
    debug_first = True
    success_count = 0
    
    for thread in threads:
        thread_id = thread.get("id")
        if not thread_id:
            continue
        ok, thread_data = run_gog_json(account, "gmail", "thread", "get", thread_id)
        if not ok:
            print(f"  ⚠️ Gmail thread get failed: {thread_id}")
            continue
        
        lines, filename_base = gmail_thread_to_markdown(thread_data, debug=debug_first)
        debug_first = False
        
        if lines and filename_base:
            write_text(gmail_dir / f"{filename_base}.md", lines)
            success_count += 1
        else:
            # フォールバック: thread_id をファイル名に使用（空でも記録）
            print(f"  ⚠️ Empty content for thread: {thread_id}")
            write_text(gmail_dir / f"{thread_id}_empty.md", [f"# Thread {thread_id}", "", "Content could not be extracted."])
    
    print(f"  ✅ Gmail threads: {success_count}/{len(threads)}")


def sync_calendar(account: str, label: str, days: int) -> None:
    ok, data = run_gog_json(account, "calendar", "events", "primary", "--days", str(days))
    if not ok:
        print(f"  ⚠️ Calendar fetch failed: {data.get('error')}")
        return
    events = data.get("events", data.get("items", [])) or []
    by_date: Dict[str, List[Dict[str, Any]]] = {}
    for event in events:
        start = event.get("start", {})
        date_key = ""
        if "dateTime" in start:
            date_key = start["dateTime"][:10]
        elif "date" in start:
            date_key = start["date"]
        if not date_key:
            continue
        by_date.setdefault(date_key, []).append(event)

    cal_dir = DATA_DIR / label / "calendar"
    for date_key, items in by_date.items():
        lines = [f"# {date_key}", ""]
        lines.append(f"_Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}_")
        lines.append("")
        for event in items:
            summary = event.get("summary", "(タイトルなし)")
            location = event.get("location", "")
            start = event.get("start", {})
            end = event.get("end", {})
            start_str = start.get("dateTime", start.get("date", ""))
            end_str = end.get("dateTime", end.get("date", ""))
            lines.append(f"- {summary}")
            if start_str:
                lines.append(f"  - start: {start_str}")
            if end_str:
                lines.append(f"  - end: {end_str}")
            if location:
                lines.append(f"  - location: {location}")
            description = event.get("description", "")
            if description:
                lines.append(f"  - description: {description}")
        write_text(cal_dir / f"{date_key}.md", lines)
    print(f"  ✅ Calendar days: {len(by_date)}")


def list_drive_files(account: str, drive_folder: str, max_files: int) -> List[Dict[str, Any]]:
    args = ["drive", "ls", "--max", str(max_files)]
    if drive_folder:
        args.extend(["--parent", drive_folder])
    ok, data = run_gog_json(account, *args)
    if not ok:
        print(f"  ⚠️ Drive list failed: {data.get('error')}")
        return []
    return data.get("files", []) or []


def export_docs(account: str, label: str, files: List[Dict[str, Any]]) -> None:
    docs_dir = DATA_DIR / label / "docs"
    for f in files:
        if f.get("mimeType") != "application/vnd.google-apps.document":
            continue
        doc_id = f.get("id")
        title = f.get("name", "untitled")
        if not doc_id:
            continue
        safe_title = safe_name(title)
        tmp_path = docs_dir / f"{safe_title}.txt"
        ok, _ = run_gog(
            account, "docs", "export", doc_id, "--format", "txt", "--out", str(tmp_path)
        )
        if not ok:
            print(f"  ⚠️ Docs export failed: {title}")
            continue
        md_path = docs_dir / f"{safe_title}.md"
        with open(tmp_path, "r", encoding="utf-8") as ftxt:
            content = ftxt.read()
        write_markdown_from_text(title, content, md_path)
        tmp_path.unlink(missing_ok=True)
    print("  ✅ Docs export complete")


def export_slides(account: str, label: str, files: List[Dict[str, Any]]) -> None:
    slides_dir = DATA_DIR / label / "slides"
    for f in files:
        if f.get("mimeType") != "application/vnd.google-apps.presentation":
            continue
        file_id = f.get("id")
        title = f.get("name", "untitled")
        if not file_id:
            continue
        safe_title = safe_name(title)
        tmp_pptx = slides_dir / f"{safe_title}.pptx"
        ok, _ = run_gog(
            account, "drive", "export", file_id, "--format", "pptx", "--out", str(tmp_pptx)
        )
        if not ok:
            print(f"  ⚠️ Slides export failed: {title}")
            continue
        md_path = slides_dir / f"{safe_title}.md"
        pptx_to_markdown(tmp_pptx, md_path, title)
        tmp_pptx.unlink(missing_ok=True)
    print("  ✅ Slides export complete")


def pptx_to_markdown(pptx_path: Path, out_path: Path, title: str) -> None:
    prs = Presentation(str(pptx_path))
    lines = [f"# {title}", "", f"_Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}_", ""]
    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"## Slide {idx}")
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = (shape.text or "").strip()
            if text:
                lines.append(text)
        lines.append("")
    write_text(out_path, lines)


def export_sheets(account: str, label: str, files: List[Dict[str, Any]]) -> None:
    sheets_root = DATA_DIR / label / "sheets"
    for f in files:
        if f.get("mimeType") != "application/vnd.google-apps.spreadsheet":
            continue
        sheet_id = f.get("id")
        title = f.get("name", "untitled")
        if not sheet_id:
            continue
        safe_title = safe_name(title)
        out_dir = sheets_root / safe_title
        out_dir.mkdir(parents=True, exist_ok=True)
        tmp_xlsx = out_dir / f"{safe_title}.xlsx"

        ok, _ = run_gog(
            account, "sheets", "export", sheet_id, "--format", "xlsx", "--out", str(tmp_xlsx)
        )
        if not ok:
            ok, _ = run_gog(
                account, "drive", "export", sheet_id, "--format", "xlsx", "--out", str(tmp_xlsx)
            )
        if ok:
            xlsx_to_csvs(tmp_xlsx, out_dir)
            tmp_xlsx.unlink(missing_ok=True)
            continue

        # Fallback: single CSV export
        tmp_csv = out_dir / f"{safe_title}.csv"
        ok, _ = run_gog(
            account, "sheets", "export", sheet_id, "--format", "csv", "--out", str(tmp_csv)
        )
        if not ok:
            print(f"  ⚠️ Sheets export failed: {title}")
    print("  ✅ Sheets export complete")


def xlsx_to_csvs(xlsx_path: Path, out_dir: Path) -> None:
    wb = openpyxl.load_workbook(str(xlsx_path), data_only=True)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        safe_sheet = safe_name(sheet_name)
        out_path = out_dir / f"{safe_sheet}.csv"
        with open(out_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            for row in ws.iter_rows(values_only=True):
                writer.writerow(["" if v is None else v for v in row])


def sync_account(account_cfg: Dict[str, str], days: int, max_threads: int, max_files: int) -> None:
    label = account_cfg["label"]
    account = account_cfg["account"]
    drive_folder = account_cfg.get("drive_folder", "")
    print(f"\n=== Sync: {label} ({account}) ===")

    DATA_DIR.joinpath(label).mkdir(parents=True, exist_ok=True)
    sync_gmail(account, label, days, max_threads)
    sync_calendar(account, label, days)

    files = list_drive_files(account, drive_folder, max_files)
    export_docs(account, label, files)
    export_slides(account, label, files)
    export_sheets(account, label, files)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--days", type=int, default=int(os.getenv("DAYS", "1")))
    parser.add_argument("--max-threads", type=int, default=int(os.getenv("MAX_GMAIL_THREADS", "50")))
    parser.add_argument("--max-files", type=int, default=int(os.getenv("MAX_DRIVE_FILES", "100")))
    args = parser.parse_args()

    accounts = load_accounts()
    if not accounts:
        print("❌ No accounts configured (GOG_ACCOUNT or GOG_ACCOUNTS_CONFIG)")
        sys.exit(1)

    for account_cfg in accounts:
        sync_account(account_cfg, args.days, args.max_threads, args.max_files)


if __name__ == "__main__":
    main()
