#!/usr/bin/env python3
"""
PPTX ディープパーサー - 全要素抽出モジュール

PowerPointファイルの全要素（テキスト・チャート・テーブル・画像・図形・グループ・SmartArt）
を詳細に抽出し、構造化データとして返す。

使用例:
    from extractor import PPTXExtractor
    ext = PPTXExtractor("source.pptx", assets_dir="extracted_assets")
    data = ext.extract_all()
"""

import json
import os
import re
import sys
import warnings
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

try:
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
except ImportError:
    XL_CHART_TYPE = None
    XL_LEGEND_POSITION = None

from lxml import etree


# ===== ユーティリティ関数 =====

def _safe_str(value) -> str:
    """None安全な文字列変換"""
    if value is None:
        return ""
    return str(value)


def _safe_enum_name(enum_val) -> str:
    """Enum値の安全な文字列取得"""
    if enum_val is None:
        return "unknown"
    try:
        return enum_val.name if hasattr(enum_val, 'name') else str(enum_val)
    except Exception:
        return str(enum_val)


def _extract_rgb_color(color_obj) -> str:
    """RGB色を '#RRGGBB' 形式で抽出。テーマ色は 'theme:X' 形式"""
    if color_obj is None:
        return ""
    try:
        rgb = color_obj.rgb
        if rgb is not None:
            return str(rgb)
    except (AttributeError, TypeError):
        pass
    try:
        if hasattr(color_obj, 'theme_color') and color_obj.theme_color is not None:
            return f"theme:{_safe_enum_name(color_obj.theme_color)}"
    except Exception:
        pass
    return ""


def _extract_color_from_font(font) -> str:
    """フォントオブジェクトからRGB色を抽出"""
    if font is None or font.color is None:
        return ""
    return _extract_rgb_color(font.color)


def _alignment_name(alignment) -> str:
    """PP_ALIGN列挙値を文字列に変換"""
    if alignment is None:
        return ""
    mapping = {
        PP_ALIGN.LEFT: "left",
        PP_ALIGN.CENTER: "center",
        PP_ALIGN.RIGHT: "right",
        PP_ALIGN.JUSTIFY: "justify",
        PP_ALIGN.DISTRIBUTE: "distribute",
    }
    return mapping.get(alignment, _safe_enum_name(alignment))


def _position_dict(shape) -> Dict[str, int]:
    """シェイプの位置・サイズをEMU値の辞書で返す"""
    return {
        "left": shape.left or 0,
        "top": shape.top or 0,
        "width": shape.width or 0,
        "height": shape.height or 0,
    }


def _emu_to_pt(emu_value) -> float:
    """EMU値をポイントに変換"""
    if emu_value is None:
        return 0
    return round(emu_value / 12700, 1)


# ===== PPTXExtractor クラス =====

class PPTXExtractor:
    """
    PowerPointファイルの全要素をディープ抽出するパーサー。

    テキスト・チャート・テーブル・画像・図形・グループ・SmartArtに対応。
    """

    MAX_GROUP_DEPTH = 3  # グループの再帰展開の最大深さ

    def __init__(self, pptx_path: str, assets_dir: Optional[str] = None):
        """
        Args:
            pptx_path: PPTXファイルパス
            assets_dir: 画像抽出先ディレクトリ（Noneの場合は画像を保存しない）
        """
        self.pptx_path = Path(pptx_path)
        self.assets_dir = Path(assets_dir) if assets_dir else None
        self._image_counter = 0  # 画像ファイル連番

        if not self.pptx_path.exists():
            raise FileNotFoundError(f"PPTXファイルが見つかりません: {pptx_path}")

        self.prs = Presentation(str(self.pptx_path))

        if self.assets_dir:
            self.assets_dir.mkdir(parents=True, exist_ok=True)

    def extract_all(self) -> Dict[str, Any]:
        """
        全スライドの全要素を抽出して構造化辞書を返す。

        Returns:
            {
                "source": "filename.pptx",
                "slide_width": EMU,
                "slide_height": EMU,
                "slides": [ { "slide_number": 1, "layout": "...", "elements": [...] }, ... ]
            }
        """
        result = {
            "source": self.pptx_path.name,
            "slide_width": self.prs.slide_width,
            "slide_height": self.prs.slide_height,
            "slides": [],
        }

        for idx, slide in enumerate(self.prs.slides):
            slide_data = self._extract_slide(slide, idx + 1)
            result["slides"].append(slide_data)

        return result

    def extract_to_json(self, indent: int = 2) -> str:
        """extract_all()の結果をJSON文字列で返す"""
        data = self.extract_all()
        return json.dumps(data, ensure_ascii=False, indent=indent, default=str)

    def _extract_slide(self, slide, slide_number: int) -> Dict[str, Any]:
        """1スライドの全要素を抽出"""
        layout_name = "Unknown"
        try:
            if slide.slide_layout:
                layout_name = slide.slide_layout.name
        except Exception:
            pass

        slide_data = {
            "slide_number": slide_number,
            "layout": layout_name,
            "elements": [],
        }

        for shape in slide.shapes:
            try:
                element = self._extract_shape(shape, slide_number)
                if element:
                    slide_data["elements"].append(element)
            except Exception as e:
                # 抽出失敗してもスキップして続行
                slide_data["elements"].append({
                    "id": getattr(shape, 'shape_id', 0),
                    "name": getattr(shape, 'name', 'unknown'),
                    "type": "error",
                    "error": str(e),
                })

        return slide_data

    def _extract_shape(self, shape, slide_number: int, depth: int = 0) -> Optional[Dict[str, Any]]:
        """シェイプを判別して適切な抽出関数を呼ぶ"""
        shape_type = shape.shape_type

        # チャート（GraphicFrame内）
        if shape.has_chart:
            return self._extract_chart(shape, slide_number)

        # テーブル（GraphicFrame内）
        if shape.has_table:
            return self._extract_table(shape, slide_number)

        # SmartArt（GraphicFrame内）
        if self._is_smartart(shape):
            return self._extract_smartart(shape, slide_number)

        # グループシェイプ
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            return self._extract_group(shape, slide_number, depth)

        # 画像
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            return self._extract_image(shape, slide_number)

        # プレースホルダー
        if shape.is_placeholder:
            return self._extract_text(shape, slide_number, is_placeholder=True)

        # テキストボックス
        if shape.has_text_frame:
            # オートシェイプ（テキスト付き）
            if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                return self._extract_auto_shape(shape, slide_number)
            # フリーフォーム等
            if shape_type == MSO_SHAPE_TYPE.FREEFORM:
                return self._extract_auto_shape(shape, slide_number)
            # 通常テキストボックス
            return self._extract_text(shape, slide_number)

        # テキストなしのオートシェイプ
        if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return self._extract_auto_shape(shape, slide_number)

        # フリーフォーム
        if shape_type == MSO_SHAPE_TYPE.FREEFORM:
            return self._extract_auto_shape(shape, slide_number)

        # その他（コネクタ等）
        return self._extract_generic(shape, slide_number)

    # ----- テキスト抽出 -----

    def _extract_text(self, shape, slide_number: int, is_placeholder: bool = False) -> Dict[str, Any]:
        """テキスト要素の詳細抽出"""
        element = {
            "id": shape.shape_id,
            "name": shape.name,
            "type": "text",
            "position": _position_dict(shape),
        }

        # プレースホルダー情報
        if is_placeholder and shape.is_placeholder:
            ph_format = shape.placeholder_format
            element["placeholder_type"] = _safe_enum_name(ph_format.type) if ph_format else ""
            element["placeholder_idx"] = ph_format.idx if ph_format else -1

        # テキスト全文
        full_text = self._get_shape_text(shape)
        element["value"] = full_text

        # スタイル情報（最初の有効なrunから取得）
        style = self._extract_text_style(shape)
        element["style"] = style

        # 段落数
        if shape.has_text_frame:
            element["style"]["paragraphs"] = len(shape.text_frame.paragraphs)
            element["style"]["word_wrap"] = shape.text_frame.word_wrap

        return element

    def _extract_text_style(self, shape) -> Dict[str, Any]:
        """テキストスタイル情報を抽出（最初の有効なrunから）"""
        style = {
            "font": "",
            "size": 0,
            "bold": False,
            "italic": False,
            "color": "",
            "align": "",
            "line_spacing": None,
        }

        if not shape.has_text_frame:
            return style

        for para in shape.text_frame.paragraphs:
            # アライメント
            if para.alignment is not None and not style["align"]:
                style["align"] = _alignment_name(para.alignment)

            # 行間
            if para.line_spacing is not None and style["line_spacing"] is None:
                style["line_spacing"] = para.line_spacing

            for run in para.runs:
                font = run.font
                if font is None:
                    continue

                if font.name and not style["font"]:
                    style["font"] = font.name

                if font.size and not style["size"]:
                    style["size"] = _emu_to_pt(font.size)

                if font.bold is not None:
                    style["bold"] = font.bold

                if font.italic is not None:
                    style["italic"] = font.italic

                color = _extract_color_from_font(font)
                if color and not style["color"]:
                    style["color"] = color

                # 最初の有効なrunで十分
                if style["font"] and style["size"]:
                    break

            if style["font"] and style["size"]:
                break

        return style

    def _get_shape_text(self, shape) -> str:
        """シェイプからテキスト全文を取得"""
        if not shape.has_text_frame:
            return ""
        lines = []
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs)
            if text.strip():
                lines.append(text)
        return "\n".join(lines)

    # ----- チャート抽出 -----

    def _extract_chart(self, shape, slide_number: int) -> Dict[str, Any]:
        """チャートデータの完全抽出"""
        chart = shape.chart

        element = {
            "id": shape.shape_id,
            "name": shape.name,
            "type": "chart",
            "position": _position_dict(shape),
        }

        # チャートタイプ
        try:
            element["chart_type"] = _safe_enum_name(chart.chart_type)
        except Exception:
            element["chart_type"] = "unknown"

        # チャート設定
        chart_config = {
            "has_legend": False,
            "legend_position": "",
            "has_data_labels": False,
            "value_axis_title": None,
            "category_axis_title": None,
        }

        try:
            chart_config["has_legend"] = chart.has_legend
            if chart.has_legend and chart.legend:
                chart_config["legend_position"] = _safe_enum_name(chart.legend.position)
        except Exception:
            pass

        try:
            for plot in chart.plots:
                if plot.has_data_labels:
                    chart_config["has_data_labels"] = True
                    break
        except Exception:
            pass

        try:
            if chart.value_axis and chart.value_axis.has_title:
                chart_config["value_axis_title"] = chart.value_axis.axis_title.text_frame.text
        except Exception:
            pass

        try:
            if chart.category_axis and chart.category_axis.has_title:
                chart_config["category_axis_title"] = chart.category_axis.axis_title.text_frame.text
        except Exception:
            pass

        element["chart_config"] = chart_config

        # データ抽出
        value = {"categories": [], "series": []}

        try:
            # カテゴリ取得
            if chart.plots and len(chart.plots) > 0:
                plot = chart.plots[0]
                cats = plot.categories
                if cats:
                    value["categories"] = [str(c) for c in cats]
        except Exception:
            pass

        try:
            # シリーズ取得
            for series in chart.series:
                series_data = {
                    "name": "",
                    "values": [],
                }
                try:
                    series_data["name"] = str(series.name) if series.name else ""
                except Exception:
                    pass
                try:
                    series_data["values"] = list(series.values) if series.values else []
                except Exception:
                    pass
                value["series"].append(series_data)
        except Exception:
            pass

        element["value"] = value
        return element

    # ----- テーブル抽出 -----

    def _extract_table(self, shape, slide_number: int) -> Dict[str, Any]:
        """テーブルの完全抽出"""
        table = shape.table

        rows = len(table.rows)
        cols = len(table.columns)

        element = {
            "id": shape.shape_id,
            "name": shape.name,
            "type": "table",
            "position": _position_dict(shape),
        }

        # テーブル設定
        table_config = {
            "rows": rows,
            "cols": cols,
            "header_row": self._detect_header_row(table),
            "merged_cells": [],
        }

        # マージセル検出
        for row_idx in range(rows):
            for col_idx in range(cols):
                try:
                    cell = table.cell(row_idx, col_idx)
                    if cell.is_merge_origin:
                        table_config["merged_cells"].append({
                            "row": row_idx,
                            "col": col_idx,
                            "span_width": cell.span_width,
                            "span_height": cell.span_height,
                        })
                except Exception:
                    pass

        element["table_config"] = table_config

        # セルスタイル（ヘッダーとボディで分けて記録）
        cell_styles = {"header": {}, "body": {}}
        if rows > 0 and cols > 0:
            try:
                header_cell = table.cell(0, 0)
                cell_styles["header"] = self._extract_cell_style(header_cell)
            except Exception:
                pass
            if rows > 1:
                try:
                    body_cell = table.cell(1, 0)
                    cell_styles["body"] = self._extract_cell_style(body_cell)
                except Exception:
                    pass

        element["cell_styles"] = cell_styles

        # 全セルのテキスト値
        value = []
        for row_idx in range(rows):
            row_data = []
            for col_idx in range(cols):
                try:
                    cell = table.cell(row_idx, col_idx)
                    row_data.append(cell.text or "")
                except Exception:
                    row_data.append("")
            value.append(row_data)

        element["value"] = value
        return element

    def _detect_header_row(self, table) -> bool:
        """ヘッダー行の検出（1行目と2行目のスタイル比較）"""
        if len(table.rows) < 2 or len(table.columns) < 1:
            return False

        try:
            header_cell = table.cell(0, 0)
            body_cell = table.cell(1, 0)

            # ヘッダーが太字かチェック
            header_bold = False
            body_bold = False

            for para in header_cell.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.bold:
                        header_bold = True

            for para in body_cell.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.bold:
                        body_bold = True

            # ヘッダーが太字でボディが太字でなければヘッダー行あり
            if header_bold and not body_bold:
                return True

            # 背景色が違えばヘッダー行あり
            header_fill = self._get_fill_color(header_cell)
            body_fill = self._get_fill_color(body_cell)
            if header_fill and body_fill and header_fill != body_fill:
                return True

        except Exception:
            pass

        return False

    def _extract_cell_style(self, cell) -> Dict[str, Any]:
        """セルのスタイルを抽出"""
        style = {
            "fill": "",
            "text_color": "",
            "font": "",
            "size": 0,
            "bold": False,
        }

        try:
            style["fill"] = self._get_fill_color(cell)
        except Exception:
            pass

        try:
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    font = run.font
                    if font.name:
                        style["font"] = font.name
                    if font.size:
                        style["size"] = _emu_to_pt(font.size)
                    if font.bold is not None:
                        style["bold"] = font.bold
                    color = _extract_color_from_font(font)
                    if color:
                        style["text_color"] = color
                    if style["font"]:
                        break
                if style["font"]:
                    break
        except Exception:
            pass

        return style

    def _get_fill_color(self, obj) -> str:
        """オブジェクトの塗りつぶし色を取得"""
        try:
            fill = obj.fill
            if fill is None:
                return ""
            if fill.type is not None:
                try:
                    return _extract_rgb_color(fill.fore_color)
                except Exception:
                    pass
        except Exception:
            pass
        return ""

    # ----- 画像抽出 -----

    def _extract_image(self, shape, slide_number: int) -> Dict[str, Any]:
        """画像の抽出・保存"""
        element = {
            "id": shape.shape_id,
            "name": shape.name,
            "type": "image",
            "position": _position_dict(shape),
        }

        image_info = {
            "content_type": "",
            "original_size": [0, 0],
            "extracted_path": "",
        }

        try:
            image = shape.image
            image_info["content_type"] = image.content_type or ""

            # 画像サイズ（ピクセル）
            try:
                from PIL import Image as PILImage
                import io
                img = PILImage.open(io.BytesIO(image.blob))
                image_info["original_size"] = list(img.size)
            except Exception:
                pass

            # 画像を保存
            if self.assets_dir:
                self._image_counter += 1
                ext = self._content_type_to_ext(image.content_type)
                filename = f"slide{slide_number}_pic{self._image_counter}{ext}"
                save_path = self.assets_dir / filename
                save_path.write_bytes(image.blob)
                image_info["extracted_path"] = str(save_path)

        except Exception as e:
            image_info["error"] = str(e)

        element["image_info"] = image_info
        element["replace_mode"] = "keep"
        return element

    def _content_type_to_ext(self, content_type: str) -> str:
        """MIMEタイプからファイル拡張子を返す"""
        mapping = {
            "image/png": ".png",
            "image/jpeg": ".jpg",
            "image/gif": ".gif",
            "image/bmp": ".bmp",
            "image/tiff": ".tiff",
            "image/svg+xml": ".svg",
            "image/x-emf": ".emf",
            "image/x-wmf": ".wmf",
        }
        return mapping.get(content_type, ".png")

    # ----- オートシェイプ抽出 -----

    def _extract_auto_shape(self, shape, slide_number: int) -> Dict[str, Any]:
        """図形の抽出"""
        element = {
            "id": shape.shape_id,
            "name": shape.name,
            "type": "shape",
            "position": _position_dict(shape),
        }

        # 形状タイプ
        try:
            element["shape_type"] = _safe_enum_name(shape.auto_shape_type)
        except Exception:
            element["shape_type"] = "unknown"

        # スタイル
        style = {
            "fill": "",
            "line_color": "",
            "line_width": 0,
            "rotation": 0,
        }

        try:
            style["fill"] = self._get_fill_color(shape)
        except Exception:
            pass

        try:
            if shape.line and shape.line.color:
                style["line_color"] = _extract_rgb_color(shape.line.color)
            if shape.line and shape.line.width:
                style["line_width"] = shape.line.width or 0
        except Exception:
            pass

        try:
            style["rotation"] = shape.rotation or 0
        except Exception:
            pass

        element["style"] = style

        # テキスト（あれば）
        text = self._get_shape_text(shape)
        if text:
            element["value"] = text
            # テキストスタイルも追加
            text_style = self._extract_text_style(shape)
            element["text_style"] = text_style

        return element

    # ----- グループ抽出 -----

    def _extract_group(self, shape, slide_number: int, depth: int = 0) -> Dict[str, Any]:
        """グループの再帰展開（深さ制限あり）"""
        element = {
            "id": shape.shape_id,
            "name": shape.name,
            "type": "group",
            "position": _position_dict(shape),
            "children": [],
        }

        if depth >= self.MAX_GROUP_DEPTH:
            element["_warning"] = f"最大深度({self.MAX_GROUP_DEPTH})に達したため子要素を省略"
            return element

        try:
            for child_shape in shape.shapes:
                child = self._extract_shape(child_shape, slide_number, depth + 1)
                if child:
                    element["children"].append(child)
        except Exception as e:
            element["_warning"] = f"子要素の抽出に失敗: {e}"

        return element

    # ----- SmartArt抽出 -----

    def _is_smartart(self, shape) -> bool:
        """SmartArtかどうかを判定"""
        try:
            # GraphicFrame で dgm (diagram) namespace があればSmartArt
            el = shape._element
            tag = el.tag
            if 'graphicFrame' not in tag and 'GraphicFrame' not in tag:
                return False

            # SmartArt の namespace を確認
            xml_str = etree.tostring(el, encoding='unicode')

            # dgm namespace は SmartArt の指標
            if 'schemas.openxmlformats.org/drawingml/2006/diagram' in xml_str:
                return True
            # relIds も確認
            if 'dgm:' in xml_str:
                return True

            # has_chart/has_table でないGraphicFrameはSmartArtの可能性
            if not shape.has_chart and not shape.has_table:
                if 'graphicFrame' in tag or 'GraphicFrame' in tag:
                    # OLE等の他の可能性もあるが、テキスト抽出を試みる
                    return True

        except Exception:
            pass
        return False

    def _extract_smartart(self, shape, slide_number: int) -> Dict[str, Any]:
        """SmartArtのテキストノードをXMLから抽出"""
        element = {
            "id": shape.shape_id,
            "name": shape.name,
            "type": "smartart",
            "position": _position_dict(shape),
            "supported": "partial",
            "text_nodes": [],
        }

        try:
            el = shape._element
            # SmartArt内のテキストノードを検索
            # 名前空間定義
            nsmap = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
            }

            # テキストノード (a:t) を全て取得
            text_elements = el.findall('.//a:t', namespaces=nsmap)

            for idx, text_el in enumerate(text_elements):
                text = text_el.text or ""
                if text.strip():
                    # xpathを生成（簡易版）
                    xpath = f".//a:t[{idx + 1}]"
                    element["text_nodes"].append({
                        "xpath": xpath,
                        "value": text.strip(),
                        "index": idx,
                    })

        except Exception as e:
            element["_warning"] = f"SmartArtテキスト抽出に失敗: {e}"

        return element

    # ----- 汎用抽出 -----

    def _extract_generic(self, shape, slide_number: int) -> Dict[str, Any]:
        """その他のシェイプの汎用抽出"""
        element = {
            "id": shape.shape_id,
            "name": shape.name,
            "type": "other",
            "shape_type": _safe_enum_name(shape.shape_type),
            "position": _position_dict(shape),
        }

        # テキストがあれば取得
        text = self._get_shape_text(shape)
        if text:
            element["value"] = text

        return element


# ===== CLI エントリーポイント =====

def main():
    """コマンドラインから直接実行"""
    import argparse

    parser = argparse.ArgumentParser(
        description="PPTX ディープパーサー - 全要素抽出ツール"
    )
    parser.add_argument("pptx_file", help="入力PPTXファイル")
    parser.add_argument(
        "--output", "-o",
        help="出力JSONファイルパス（省略時は標準出力）"
    )
    parser.add_argument(
        "--assets-dir", "-a",
        help="画像抽出先ディレクトリ"
    )
    parser.add_argument(
        "--pretty", action="store_true", default=True,
        help="JSON整形出力（デフォルト: True）"
    )

    args = parser.parse_args()

    extractor = PPTXExtractor(args.pptx_file, assets_dir=args.assets_dir)
    data = extractor.extract_all()

    indent = 2 if args.pretty else None
    json_str = json.dumps(data, ensure_ascii=False, indent=indent, default=str)

    if args.output:
        Path(args.output).write_text(json_str, encoding="utf-8")
        print(f"抽出完了: {args.output}")
        slide_count = len(data["slides"])
        element_count = sum(len(s["elements"]) for s in data["slides"])
        print(f"  スライド数: {slide_count}")
        print(f"  要素数: {element_count}")
    else:
        print(json_str)


if __name__ == "__main__":
    main()
