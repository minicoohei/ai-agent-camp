"""
PPTX Generator - コピー&書き換えエンジン

元PPTXをコピーして、マッピングYAMLに基づいて各要素を書き換える。
スライドマスター・テーマ・アニメーション・SmartArtを完全に保持。

使い方:
    from generator import PPTXGenerator
    gen = PPTXGenerator("source.pptx", "output.pptx")
    gen.apply_replacements(mapping, new_data)
"""

import json
import os
import shutil
import sys
import warnings
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches, Pt

try:
    from lxml import etree
except ImportError:
    etree = None


# ===== PPTXGenerator クラス =====

class PPTXGenerator:
    """
    元PPTXをコピーして要素を書き換えるエンジン。

    コピー方式により、以下が完全に保持される:
    - スライドマスター・テーマカラー・フォント定義
    - アニメーション・トランジション
    - SmartArt・OLEオブジェクト
    - レイアウト・配置
    """

    def __init__(self, source_pptx: str, output_path: str):
        """
        Args:
            source_pptx: 元PPTXファイルパス
            output_path: 出力PPTXファイルパス
        """
        self.source_path = Path(source_pptx)
        self.output_path = Path(output_path)

        if not self.source_path.exists():
            raise FileNotFoundError(f"元PPTXが見つかりません: {source_pptx}")

        # 出力ディレクトリを作成
        self.output_path.parent.mkdir(parents=True, exist_ok=True)

        # 元PPTXをコピー
        shutil.copy2(str(self.source_path), str(self.output_path))

        # コピーを開く
        self.prs = Presentation(str(self.output_path))
        self._stats = {"replaced": 0, "skipped": 0, "errors": []}

    def apply_replacements(self, mapping: dict, new_data: dict) -> dict:
        """マッピングに基づいて全要素を書き換え、結果を保存する。

        Args:
            mapping: mapper.py が生成したマッピング辞書。
                {
                    "slides": [
                        {
                            "slide_number": 1,
                            "elements": [
                                { "id": N, "name": "...", "type": "text",
                                  "placeholder": "{{slide_1_title}}", ... },
                                ...
                            ]
                        }, ...
                    ]
                }
            new_data: プレースホルダーキー → 新しい値 のマッピング。
                {
                    "{{slide_1_title}}": "新しいタイトル",
                    "{{slide_3_chart}}": {
                        "categories": [...],
                        "series": [{"name": "...", "values": [...]}]
                    },
                    "{{slide_5_table}}": [["A","B"], ["1","2"]],
                    ...
                }

        Returns:
            統計情報の辞書。
        """
        for slide_map in mapping.get("slides", []):
            slide_num = slide_map.get("slide_number", 0)
            if slide_num < 1 or slide_num > len(self.prs.slides):
                self._stats["errors"].append(
                    f"スライド番号 {slide_num} が範囲外です"
                )
                continue

            slide = self.prs.slides[slide_num - 1]

            for element in slide_map.get("elements", []):
                self._process_element(slide, element, new_data)

        # 保存
        self.prs.save(str(self.output_path))
        return self._stats

    def _process_element(self, slide, element: dict, new_data: dict):
        """1要素の書き換え処理。"""
        etype = element.get("type", "")
        placeholder = element.get("placeholder")

        # グループ内子要素の処理
        if etype == "group" and "children" in element:
            self._process_group_element(slide, element, new_data)
            return

        # SmartArt テキストノードの処理
        if etype == "smartart" and "text_nodes" in element:
            self._process_smartart_element(slide, element, new_data)
            return

        # プレースホルダーがない or データがなければスキップ
        if not placeholder or placeholder not in new_data:
            self._stats["skipped"] += 1
            return

        new_value = new_data[placeholder]
        self._replace_element(slide, element, new_value)

    def _replace_element(self, slide, element: dict, new_value):
        """要素タイプに応じた書き換え。"""
        shape_id = element.get("id", 0)
        shape_name = element.get("name", "")
        etype = element.get("type", "")

        shape = self._find_shape(slide, shape_id, shape_name)
        if not shape:
            self._stats["errors"].append(
                f"シェイプが見つかりません: id={shape_id}, name={shape_name}"
            )
            return

        try:
            if etype == "text":
                self._replace_text(shape, new_value, element.get("style"))
            elif etype == "chart":
                self._replace_chart(shape, new_value, element)
            elif etype == "table":
                self._replace_table(shape, new_value, element.get("table_config"))
            elif etype == "image":
                self._replace_image(shape, new_value, element)
            elif etype == "shape":
                self._replace_shape_text(shape, new_value)
            else:
                # 未対応タイプ: テキストがあればテキストとして書き換え
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    self._replace_shape_text(shape, str(new_value))
                else:
                    self._stats["skipped"] += 1
                    return

            self._stats["replaced"] += 1

        except Exception as e:
            self._stats["errors"].append(
                f"書き換え失敗 (id={shape_id}, type={etype}): {e}"
            )

    # ================================================================
    # シェイプ検索
    # ================================================================

    def _find_shape(self, slide, shape_id: int, shape_name: str = ""):
        """shape_idでシェイプを検索（グループ内も再帰探索）。"""
        for shape in slide.shapes:
            if shape.shape_id == shape_id:
                return shape
            # グループ内の再帰検索
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                found = self._find_shape_recursive(shape, shape_id, shape_name)
                if found:
                    return found

        # shape_id で見つからなければ name で検索（フォールバック）
        if shape_name:
            for shape in slide.shapes:
                if shape.name == shape_name:
                    return shape
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    found = self._find_by_name_recursive(shape, shape_name)
                    if found:
                        return found

        return None

    def _find_shape_recursive(self, group_shape, shape_id: int, shape_name: str = "", depth: int = 0):
        """グループ内をshape_idで再帰検索。"""
        if depth > 5:
            return None
        try:
            for child in group_shape.shapes:
                if child.shape_id == shape_id:
                    return child
                if child.shape_type == MSO_SHAPE_TYPE.GROUP:
                    found = self._find_shape_recursive(child, shape_id, shape_name, depth + 1)
                    if found:
                        return found
        except Exception:
            pass
        return None

    def _find_by_name_recursive(self, group_shape, shape_name: str, depth: int = 0):
        """グループ内をnameで再帰検索。"""
        if depth > 5:
            return None
        try:
            for child in group_shape.shapes:
                if child.name == shape_name:
                    return child
                if child.shape_type == MSO_SHAPE_TYPE.GROUP:
                    found = self._find_by_name_recursive(child, shape_name, depth + 1)
                    if found:
                        return found
        except Exception:
            pass
        return None

    # ================================================================
    # テキスト書き換え
    # ================================================================

    def _replace_text(self, shape, new_text, style_info: dict = None):
        """テキストを書き換え（run単位でスタイル保持）。

        戦略:
        1. 既存のrunsのスタイル情報を記録
        2. 段落構造に合わせてテキストを配分
        3. スタイルを復元
        """
        if not shape.has_text_frame:
            return

        if not isinstance(new_text, str):
            new_text = str(new_text)

        tf = shape.text_frame
        new_lines = new_text.split("\n") if new_text else [""]

        # 既存段落のスタイルを保存
        saved_styles = self._save_paragraph_styles(tf)

        # テキストを書き換え
        for i, para in enumerate(tf.paragraphs):
            if i < len(new_lines):
                self._replace_paragraph_text(para, new_lines[i], saved_styles.get(i))
            else:
                # 余った段落はテキストをクリア
                self._clear_paragraph(para)

        # 新しい行が段落数より多い場合、最後の段落に追記
        if len(new_lines) > len(tf.paragraphs):
            last_para = tf.paragraphs[-1]
            remaining = "\n".join(new_lines[len(tf.paragraphs):])
            current = last_para.text
            combined = current + "\n" + remaining if current else remaining
            style = saved_styles.get(len(tf.paragraphs) - 1)
            self._replace_paragraph_text(last_para, combined, style)

    def _save_paragraph_styles(self, text_frame) -> Dict[int, dict]:
        """各段落の最初のrunのスタイルを保存。"""
        styles = {}
        for idx, para in enumerate(text_frame.paragraphs):
            style = {
                "alignment": para.alignment,
                "line_spacing": para.line_spacing,
                "space_before": para.space_before,
                "space_after": para.space_after,
                "level": para.level,
                "runs": [],
            }
            for run in para.runs:
                run_style = self._save_run_style(run)
                style["runs"].append(run_style)
            styles[idx] = style
        return styles

    def _save_run_style(self, run) -> dict:
        """runのスタイル情報を保存。"""
        font = run.font
        style = {}

        try:
            style["name"] = font.name
        except Exception:
            style["name"] = None

        try:
            style["size"] = font.size
        except Exception:
            style["size"] = None

        try:
            style["bold"] = font.bold
        except Exception:
            style["bold"] = None

        try:
            style["italic"] = font.italic
        except Exception:
            style["italic"] = None

        try:
            style["underline"] = font.underline
        except Exception:
            style["underline"] = None

        # 色情報はXML要素レベルで保存（最も確実）
        try:
            rPr = run._r.get_or_add_rPr()
            solidFill = rPr.find(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill'
            )
            if solidFill is not None:
                style["color_xml"] = etree.tostring(solidFill, encoding='unicode') if etree else None
            else:
                style["color_xml"] = None
        except Exception:
            style["color_xml"] = None

        try:
            if font.color and font.color.rgb:
                style["color_rgb"] = str(font.color.rgb)
            else:
                style["color_rgb"] = None
        except Exception:
            style["color_rgb"] = None

        return style

    def _apply_run_style(self, run, style: dict):
        """保存したスタイルをrunに適用。"""
        font = run.font

        if style.get("name"):
            font.name = style["name"]
        if style.get("size") is not None:
            font.size = style["size"]
        if style.get("bold") is not None:
            font.bold = style["bold"]
        if style.get("italic") is not None:
            font.italic = style["italic"]
        if style.get("underline") is not None:
            font.underline = style["underline"]

        # 色復元（XML直接操作が最も確実）
        if style.get("color_xml") and etree is not None:
            try:
                rPr = run._r.get_or_add_rPr()
                ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                old_fill = rPr.find(f'{ns}solidFill')
                if old_fill is not None:
                    rPr.remove(old_fill)
                new_fill = etree.fromstring(style["color_xml"])
                rPr.insert(0, new_fill)
            except Exception:
                # XML復元失敗時はRGB fallback
                if style.get("color_rgb"):
                    try:
                        from pptx.dml.color import RGBColor
                        font.color.rgb = RGBColor.from_string(style["color_rgb"])
                    except Exception:
                        pass
        elif style.get("color_rgb"):
            try:
                from pptx.dml.color import RGBColor
                font.color.rgb = RGBColor.from_string(style["color_rgb"])
            except Exception:
                pass

    def _replace_paragraph_text(self, para, new_text: str, saved_style: dict = None):
        """段落のテキストをスタイル保持で書き換え。"""
        runs = list(para.runs)

        if not runs:
            # runがなければ新規作成
            run = para.add_run()
            run.text = new_text
            if saved_style and saved_style.get("runs"):
                self._apply_run_style(run, saved_style["runs"][0])
            return

        # 最初のrunに全テキストを設定、残りのrunをクリア
        first_run_style = saved_style["runs"][0] if saved_style and saved_style.get("runs") else None

        runs[0].text = new_text
        if first_run_style:
            self._apply_run_style(runs[0], first_run_style)

        # 残りのrunのテキストをクリア
        for run in runs[1:]:
            run.text = ""

    def _clear_paragraph(self, para):
        """段落のテキストをクリア（構造は保持）。"""
        for run in para.runs:
            run.text = ""

    # ================================================================
    # チャート書き換え
    # ================================================================

    def _replace_chart(self, shape, new_data, element: dict = None):
        """チャートデータを差し替え（チャートタイプ・色保持）。

        Args:
            shape: チャートを含むシェイプ
            new_data: 新しいデータ
                {
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": [
                        {"name": "売上", "values": [120, 145, 160, 180]},
                        {"name": "利益", "values": [30, 38, 42, 50]}
                    ]
                }
            element: マッピング要素（chart_type等参照用）
        """
        if not shape.has_chart:
            return

        chart = shape.chart

        if not isinstance(new_data, dict):
            return

        categories = new_data.get("categories", [])
        series_list = new_data.get("series", [])

        if not categories and not series_list:
            return

        # チャートタイプに応じてデータオブジェクトを選択
        chart_type_name = ""
        if element and element.get("chart_type"):
            chart_type_name = element["chart_type"].upper()
        else:
            try:
                chart_type_name = chart.chart_type.name if hasattr(chart.chart_type, 'name') else str(chart.chart_type)
            except Exception:
                chart_type_name = ""

        # XY (散布図) の場合は XyChartData
        if "XY" in chart_type_name or "SCATTER" in chart_type_name:
            self._replace_xy_chart(chart, new_data)
        else:
            self._replace_category_chart(chart, categories, series_list)

    def _replace_category_chart(self, chart, categories: list, series_list: list):
        """カテゴリ型チャート（棒、折れ線、円等）のデータ差し替え。"""
        chart_data = CategoryChartData()

        if categories:
            chart_data.categories = categories

        for s in series_list:
            name = s.get("name", "")
            values = s.get("values", [])
            # 値を数値に変換
            numeric_values = []
            for v in values:
                try:
                    numeric_values.append(float(v) if v is not None else 0)
                except (ValueError, TypeError):
                    numeric_values.append(0)
            chart_data.add_series(name, numeric_values)

        chart.replace_data(chart_data)

    def _replace_xy_chart(self, chart, new_data: dict):
        """散布図（XY）チャートのデータ差し替え。"""
        chart_data = XyChartData()

        for s in new_data.get("series", []):
            name = s.get("name", "")
            series = chart_data.add_series(name)
            xs = s.get("x_values", s.get("categories", []))
            ys = s.get("values", s.get("y_values", []))
            for x, y in zip(xs, ys):
                try:
                    series.add_data_point(float(x), float(y))
                except (ValueError, TypeError):
                    pass

        chart.replace_data(chart_data)

    # ================================================================
    # テーブル書き換え
    # ================================================================

    def _replace_table(self, shape, new_data, table_config: dict = None):
        """テーブルのセル内容を差し替え（セルスタイル保持）。

        Args:
            shape: テーブルを含むシェイプ
            new_data: 2次元配列
                [["A","B","C"], ["1","2","3"], ...]
            table_config: テーブル設定（rows, cols, header_row等）
        """
        if not shape.has_table:
            return

        if not isinstance(new_data, list):
            return

        table = shape.table
        rows = len(table.rows)
        cols = len(table.columns)

        for row_idx, row_data in enumerate(new_data):
            if row_idx >= rows:
                break
            if not isinstance(row_data, list):
                continue
            for col_idx, cell_value in enumerate(row_data):
                if col_idx >= cols:
                    break
                try:
                    cell = table.cell(row_idx, col_idx)
                    self._replace_cell_text(cell, str(cell_value) if cell_value is not None else "")
                except Exception:
                    pass

    def _replace_cell_text(self, cell, new_text: str):
        """テーブルセルのテキストをスタイル保持で差し替え。"""
        if not cell.text_frame:
            return

        paras = cell.text_frame.paragraphs
        if not paras:
            return

        # 最初の段落の最初のrunにテキストを設定
        first_para = paras[0]
        runs = list(first_para.runs)

        if runs:
            # 既存runのスタイルを保持してテキストだけ差し替え
            runs[0].text = new_text
            for run in runs[1:]:
                run.text = ""
        else:
            # runがなければ新規追加
            run = first_para.add_run()
            run.text = new_text

        # 残りの段落をクリア
        for para in paras[1:]:
            for run in para.runs:
                run.text = ""

    # ================================================================
    # 画像書き換え
    # ================================================================

    def _replace_image(self, shape, new_value, element: dict):
        """画像を差し替え（replace_modeに応じて）。

        replace_mode:
            "keep": 何もしない（デフォルト）
            "file": new_value をファイルパスとして画像差し替え
            "generate": 将来的にAI画像生成対応
        """
        replace_mode = element.get("replace_mode", "keep")

        if replace_mode == "keep":
            return

        if replace_mode == "file":
            if isinstance(new_value, str) and Path(new_value).exists():
                self._swap_image(shape, new_value)
            else:
                self._stats["errors"].append(
                    f"画像ファイルが見つかりません: {new_value}"
                )
            return

        if replace_mode == "generate":
            # AI画像生成（将来実装）
            if isinstance(new_value, str):
                self._stats["errors"].append(
                    f"画像生成モードは未実装です: {new_value}"
                )
            return

    def _swap_image(self, shape, image_path: str):
        """シェイプの画像を差し替え。

        python-pptxは直接的な画像差し替えAPIがないため、
        XML操作で画像リレーションを更新する。
        """
        try:
            from pptx.opc.constants import RELATIONSHIP_TYPE as RT

            image_path = Path(image_path)
            if not image_path.exists():
                return

            # 画像のパーツを取得して差し替え
            image_part = shape.image
            blob = image_path.read_bytes()

            # content_typeを拡張子から推定
            ext_map = {
                ".png": "image/png",
                ".jpg": "image/jpeg",
                ".jpeg": "image/jpeg",
                ".gif": "image/gif",
                ".bmp": "image/bmp",
                ".tiff": "image/tiff",
                ".svg": "image/svg+xml",
            }
            # 画像パーツの blob を直接更新
            image_part._blob = blob

        except Exception as e:
            self._stats["errors"].append(f"画像差し替え失敗: {e}")

    # ================================================================
    # シェイプテキスト書き換え
    # ================================================================

    def _replace_shape_text(self, shape, new_text):
        """図形内テキストを差し替え。"""
        if not isinstance(new_text, str):
            new_text = str(new_text)

        if not shape.has_text_frame:
            return

        self._replace_text(shape, new_text)

    # ================================================================
    # SmartArt 書き換え
    # ================================================================

    def _process_smartart_element(self, slide, element: dict, new_data: dict):
        """SmartArtのテキストノードを個別に書き換え。"""
        shape_id = element.get("id", 0)
        shape_name = element.get("name", "")
        text_nodes = element.get("text_nodes", [])

        shape = self._find_shape(slide, shape_id, shape_name)
        if not shape:
            self._stats["errors"].append(
                f"SmartArtシェイプが見つかりません: id={shape_id}"
            )
            return

        for node in text_nodes:
            placeholder = node.get("placeholder")
            if not placeholder or placeholder not in new_data:
                continue

            new_text = str(new_data[placeholder])
            self._replace_smartart_node(shape, node, new_text)

    def _replace_smartart_node(self, shape, node_info: dict, new_text: str):
        """SmartArtの1テキストノードをXML直接操作で書き換え。"""
        if etree is None:
            self._stats["errors"].append("lxml が必要です: pip install lxml")
            return

        try:
            el = shape._element
            nsmap = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            }

            # テキストノードを検索
            text_elements = el.findall('.//a:t', namespaces=nsmap)
            index = node_info.get("index", -1)

            if 0 <= index < len(text_elements):
                text_elements[index].text = new_text
                self._stats["replaced"] += 1
            else:
                # xpath で検索（フォールバック）
                xpath = node_info.get("xpath", "")
                if xpath:
                    found = el.xpath(xpath, namespaces=nsmap)
                    if found:
                        found[0].text = new_text
                        self._stats["replaced"] += 1
                    else:
                        self._stats["errors"].append(
                            f"SmartArt xpath '{xpath}' でノードが見つかりません"
                        )

        except Exception as e:
            self._stats["errors"].append(f"SmartArtノード書き換え失敗: {e}")

    # ================================================================
    # グループ書き換え
    # ================================================================

    def _process_group_element(self, slide, element: dict, new_data: dict):
        """グループ内の子要素を書き換え。"""
        shape_id = element.get("id", 0)
        shape_name = element.get("name", "")
        children = element.get("children", [])

        shape = self._find_shape(slide, shape_id, shape_name)
        if not shape:
            self._stats["errors"].append(
                f"グループシェイプが見つかりません: id={shape_id}"
            )
            return

        for child in children:
            child_placeholder = child.get("placeholder")
            if not child_placeholder or child_placeholder not in new_data:
                continue

            child_id = child.get("id", 0)
            child_name = child.get("name", "")
            child_type = child.get("type", "")
            new_value = new_data[child_placeholder]

            # グループ内から子シェイプを検索
            child_shape = self._find_shape_recursive(shape, child_id, child_name)
            if not child_shape:
                # フォールバック: name検索
                child_shape = self._find_by_name_recursive(shape, child_name)

            if not child_shape:
                self._stats["errors"].append(
                    f"グループ子要素が見つかりません: id={child_id}, name={child_name}"
                )
                continue

            try:
                if child_type == "text":
                    self._replace_text(child_shape, str(new_value), child.get("style"))
                elif child_type == "shape":
                    self._replace_shape_text(child_shape, str(new_value))
                elif child_type == "chart":
                    self._replace_chart(child_shape, new_value, child)
                elif child_type == "table":
                    self._replace_table(child_shape, new_value, child.get("table_config"))
                elif child_type == "image":
                    self._replace_image(child_shape, new_value, child)
                else:
                    if hasattr(child_shape, 'has_text_frame') and child_shape.has_text_frame:
                        self._replace_shape_text(child_shape, str(new_value))
                self._stats["replaced"] += 1
            except Exception as e:
                self._stats["errors"].append(
                    f"グループ子要素書き換え失敗 (id={child_id}): {e}"
                )

    # ================================================================
    # ユーティリティ
    # ================================================================

    def get_stats(self) -> dict:
        """書き換え統計を取得。"""
        return dict(self._stats)


# ===== CLI エントリーポイント =====

def main():
    """コマンドラインから直接実行（テスト・デバッグ用）。

    使い方:
        python generator.py source.pptx mapping.yaml data.yaml -o output.pptx
    """
    import argparse

    import yaml

    parser = argparse.ArgumentParser(
        description="PPTX Generator - コピー&書き換えエンジン"
    )
    parser.add_argument("source", help="元PPTXファイル")
    parser.add_argument("mapping", help="マッピングYAMLファイル")
    parser.add_argument("--data", "-d", required=True, help="データYAMLファイル")
    parser.add_argument("--output", "-o", required=True, help="出力PPTXファイル")

    args = parser.parse_args()

    # マッピングYAML読み込み
    with open(args.mapping, "r", encoding="utf-8") as f:
        mapping = yaml.safe_load(f)

    # データYAML読み込み
    with open(args.data, "r", encoding="utf-8") as f:
        new_data = yaml.safe_load(f)

    print(f"[generator] 元PPTX: {args.source}")
    print(f"[generator] マッピング: {args.mapping}")
    print(f"[generator] データ: {args.data}")
    print(f"[generator] 出力先: {args.output}")

    # 生成実行
    gen = PPTXGenerator(args.source, args.output)
    stats = gen.apply_replacements(mapping, new_data)

    print(f"\n[generator] === 書き換え完了 ===")
    print(f"[generator]   書き換え済み: {stats['replaced']}要素")
    print(f"[generator]   スキップ: {stats['skipped']}要素")
    if stats["errors"]:
        print(f"[generator]   エラー: {len(stats['errors'])}件")
        for err in stats["errors"][:10]:
            print(f"[generator]     - {err}")

    print(f"[generator]   出力: {Path(args.output).resolve()}")


if __name__ == "__main__":
    main()
