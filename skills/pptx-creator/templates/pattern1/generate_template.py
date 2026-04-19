#!/usr/bin/env python3
"""Generate McKinsey-style template PPTX (pattern1).

Creates a sample deck demonstrating all available slide layouts
with placeholder content. This serves as both a reusable template
and a visual reference for the design system.
"""

import io
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# ===========================================================================
# Design System — McKinsey-Inspired
# ===========================================================================

# Canvas
SW = Inches(13.333)
SH = Inches(7.5)

# Color Palette
NAVY       = RGBColor(0x05, 0x1C, 0x2C)
BLUE       = RGBColor(0x00, 0x5E, 0xB8)
SLATE      = RGBColor(0x6B, 0x7B, 0x8D)
CHARCOAL   = RGBColor(0x22, 0x22, 0x22)
GREY_DARK  = RGBColor(0x4A, 0x4A, 0x4A)
GREY_MED   = RGBColor(0x66, 0x66, 0x66)
GREY_LIGHT = RGBColor(0xA2, 0xAA, 0xAD)
GREY_LINE  = RGBColor(0xD0, 0xD0, 0xD0)
GREY_BG    = RGBColor(0xF5, 0xF5, 0xF5)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

# Margins
ML = Inches(0.6)
MR = Inches(0.6)
CW = SW - ML - MR

# Vertical Grid
TITLE_Y  = Inches(0.25)
TITLE_H  = Inches(0.55)
SUB_Y    = Inches(0.80)
SEP_Y    = Inches(1.12)
BODY_Y   = Inches(1.33)
BODY_H   = SH - BODY_Y - Inches(0.3)
FOOTER_Y = Inches(7.18)

# Section Labels
LABEL_H   = Inches(0.3)
LABEL_GAP = Inches(0.06)

# Fonts
FONT_TITLE = "Georgia"
FONT_BODY  = "Arial"
FONT_NUM   = "Futura"


# ===========================================================================
# Primitive Drawing Helpers
# ===========================================================================

def add_textbox(slide, left, top, width, height, text,
                font_name=FONT_BODY, font_size=16, bold=False,
                color=CHARCOAL, align=PP_ALIGN.LEFT, line_spacing=None,
                v_anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    txBox.text_frame.paragraphs[0].alignment = align
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    if line_spacing:
        p.line_spacing = line_spacing
    # Vertical anchor
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", {
            MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b"
        }.get(v_anchor, "t"))
    return txBox


def add_rect(slide, left, top, width, height,
             fill=None, line_color=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_section_label(slide, x, y, w, text, color=BLUE):
    """Section label: thin colored bar with centered text legend."""
    line_y = y + Inches(0.14)
    add_rect(slide, x, line_y, w, Pt(2.5), fill=color)
    # Text width estimation (CJK-aware)
    cjk_count = sum(1 for c in text if ord(c) > 0x2E7F)
    ascii_count = len(text) - cjk_count
    text_w = max(int((cjk_count * Pt(15) + ascii_count * Pt(9)) * 1.15), Inches(1.5))
    text_w = min(text_w, w - Inches(0.2))
    text_x = x + (w - text_w) // 2
    add_rect(slide, text_x - Inches(0.05), y + Inches(0.02),
             text_w + Inches(0.1), Inches(0.26), fill=WHITE)
    add_textbox(slide, text_x, y, text_w, LABEL_H, text,
                font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def add_action_title(slide, title, subtitle=None):
    """McKinsey action title with thin separator."""
    add_textbox(slide, ML, TITLE_Y, CW, TITLE_H, title,
                font_name=FONT_TITLE, font_size=24, bold=True, color=NAVY)
    if subtitle:
        add_textbox(slide, ML, SUB_Y, CW, Inches(0.3), subtitle,
                    font_size=14, color=GREY_MED)
    # Thin separator
    add_rect(slide, ML, SEP_Y, CW, Pt(1.5), fill=GREY_LINE)


def add_footer(slide, page_num, total=10):
    add_textbox(slide, SW - Inches(1.2), FOOTER_Y, Inches(0.8), Inches(0.25),
                f"{page_num} / {total}",
                font_name=FONT_NUM, font_size=9, color=GREY_LIGHT,
                align=PP_ALIGN.RIGHT)


# ===========================================================================
# Slide Builders — Each demonstrates a layout pattern
# ===========================================================================

def build_title(prs):
    """1. Title slide — dark navy background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    # Accent line
    add_rect(slide, Inches(2.5), Inches(4.4), Inches(8.3), Pt(2.5), fill=BLUE)
    # Title
    add_textbox(slide, Inches(2.5), Inches(2.0), Inches(8.3), Inches(1.8),
                "Presentation Title",
                font_name=FONT_TITLE, font_size=36, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    # Subtitle
    add_textbox(slide, Inches(2.5), Inches(4.7), Inches(8.3), Inches(0.5),
                "Subtitle  |  Month Year",
                font_size=16, color=GREY_LIGHT, align=PP_ALIGN.CENTER)
    add_footer(slide, 1)


def build_exec_summary(prs):
    """2. Executive Summary — KPI dots + descriptions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "Executive Summary",
                     "Key takeaways and decision points")

    kpis = [
        ("KPI Label 1", "Description of the first key metric or finding"),
        ("KPI Label 2", "Description of the second key metric or finding"),
        ("KPI Label 3", "Description of the third key metric or finding"),
        ("KPI Label 4", "Description of the fourth key metric or finding"),
    ]
    dot_colors = [BLUE, NAVY, SLATE, GREY_DARK]
    y = BODY_Y + Inches(0.1)
    for i, (label, desc) in enumerate(kpis):
        # Dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      ML + Inches(0.1), y + Inches(0.08),
                                      Inches(0.26), Inches(0.26))
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_colors[i]
        dot.line.fill.background()
        # Label
        add_textbox(slide, ML + Inches(0.55), y, Inches(3.0), Inches(0.4),
                    label, font_size=16, bold=True, color=NAVY)
        # Description
        add_textbox(slide, ML + Inches(3.8), y, Inches(8.2), Inches(0.4),
                    desc, font_size=16, color=CHARCOAL)
        # Hairline
        add_rect(slide, ML + Inches(0.55), y + Inches(0.55),
                 CW - Inches(0.85), Pt(0.5), fill=GREY_LINE)
        y += Inches(0.7)

    add_footer(slide, 2)


def build_content(prs):
    """3. Content slide — bullets in grey box."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "Content Slide Title",
                     "Action-oriented subtitle summarizing the key insight")

    add_section_label(slide, ML, BODY_Y, CW, "Section Heading")
    box_y = BODY_Y + LABEL_H + LABEL_GAP
    box_h = BODY_H - Inches(0.15) - LABEL_H - LABEL_GAP
    add_rect(slide, ML, box_y, CW, box_h,
             fill=GREY_BG, line_color=GREY_LINE, line_width=Pt(0.5))

    bullets = "•  First bullet point with key information\n•  Second bullet point elaborating on the topic\n•  Third bullet point with supporting data\n•  Fourth bullet point with recommendation"
    add_textbox(slide, ML + Inches(0.2), box_y + Inches(0.15),
                CW - Inches(0.4), box_h - Inches(0.3), bullets,
                font_size=16, color=CHARCOAL, line_spacing=Pt(30))

    add_footer(slide, 3)


def build_content_with_image(prs):
    """4. Content + Image — left bullets, right image area."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "Content + Visual Layout",
                     "Data or diagram supports the narrative")

    gap_col = Inches(0.3)
    main_w = CW * 0.58
    side_w = CW - main_w - gap_col
    side_x = ML + main_w + gap_col

    add_section_label(slide, ML, BODY_Y, main_w, "Key Points")
    add_section_label(slide, side_x, BODY_Y, side_w, "Visual / Data", color=SLATE)

    box_y = BODY_Y + LABEL_H + LABEL_GAP
    box_h = BODY_H - Inches(0.15) - LABEL_H - LABEL_GAP

    # Left: bullets
    add_rect(slide, ML, box_y, main_w, box_h,
             fill=GREY_BG, line_color=GREY_LINE, line_width=Pt(0.5))
    bullets = "•  Primary insight or finding\n•  Supporting data point\n•  Strategic recommendation\n•  Next steps"
    add_textbox(slide, ML + Inches(0.2), box_y + Inches(0.15),
                main_w - Inches(0.4), box_h - Inches(0.3), bullets,
                font_size=16, color=CHARCOAL, line_spacing=Pt(30))

    # Right: image placeholder
    add_rect(slide, side_x, box_y, side_w, box_h,
             fill=RGBColor(0xE8, 0xEE, 0xF4), line_color=GREY_LINE, line_width=Pt(0.5))
    add_textbox(slide, side_x, box_y + box_h // 2 - Inches(0.2), side_w, Inches(0.4),
                "[Image / Chart / Diagram]",
                font_size=14, color=GREY_MED, align=PP_ALIGN.CENTER)

    add_footer(slide, 4)


def build_comparison(prs):
    """5. Comparison — two equal panels."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "Comparison Layout",
                     "Side-by-side analysis of two approaches")

    gap_col = Inches(0.3)
    half_w = (CW - gap_col) / 2
    right_x = ML + half_w + gap_col

    add_section_label(slide, ML, BODY_Y, half_w, "Option A")
    add_section_label(slide, right_x, BODY_Y, half_w, "Option B", color=SLATE)

    box_y = BODY_Y + LABEL_H + LABEL_GAP
    box_h = BODY_H - Inches(0.7) - LABEL_H - LABEL_GAP

    # Left panel
    add_rect(slide, ML, box_y, half_w, box_h,
             fill=GREY_BG, line_color=GREY_LINE, line_width=Pt(0.5))
    add_textbox(slide, ML + Inches(0.2), box_y + Inches(0.15),
                half_w - Inches(0.4), box_h - Inches(0.3),
                "•  Advantage point 1\n•  Advantage point 2\n•  Advantage point 3\n•  Advantage point 4",
                font_size=16, color=CHARCOAL, line_spacing=Pt(30))

    # Right panel
    add_rect(slide, right_x, box_y, half_w, box_h,
             fill=GREY_BG, line_color=GREY_LINE, line_width=Pt(0.5))
    add_textbox(slide, right_x + Inches(0.2), box_y + Inches(0.15),
                half_w - Inches(0.4), box_h - Inches(0.3),
                "•  Advantage point 1\n•  Advantage point 2\n•  Advantage point 3\n•  Advantage point 4",
                font_size=16, color=CHARCOAL, line_spacing=Pt(30))

    # Key message
    km_y = box_y + box_h + Inches(0.1)
    add_textbox(slide, ML + Inches(1.5), km_y, CW - Inches(3), Inches(0.5),
                "Key takeaway message summarizing the comparison",
                font_size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_footer(slide, 5)


def build_quadrant(prs):
    """6. Quadrant (2x2 Matrix) — four boxes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "2x2 Framework",
                     "Four-quadrant analysis")

    gap = Inches(0.25)
    half_w = (CW - gap) / 2
    half_h = (BODY_H - Inches(0.15) - gap) / 2
    colors = [
        RGBColor(0xE8, 0xEE, 0xF4), RGBColor(0xF0, 0xF4, 0xF8),
        RGBColor(0xF0, 0xF4, 0xF8), RGBColor(0xE8, 0xEE, 0xF4),
    ]
    titles = ["Quadrant 1", "Quadrant 2", "Quadrant 3", "Quadrant 4"]
    bodies = [
        "Description of the first\nquadrant analysis",
        "Description of the second\nquadrant analysis",
        "Description of the third\nquadrant analysis",
        "Description of the fourth\nquadrant analysis",
    ]
    positions = [
        (ML, BODY_Y), (ML + half_w + gap, BODY_Y),
        (ML, BODY_Y + half_h + gap), (ML + half_w + gap, BODY_Y + half_h + gap),
    ]
    t_colors = [NAVY, BLUE, BLUE, NAVY]

    for i, ((x, y), title, body) in enumerate(zip(positions, titles, bodies)):
        add_rect(slide, x, y, half_w, half_h,
                 fill=colors[i], line_color=GREY_LINE, line_width=Pt(0.5))
        add_textbox(slide, x + Inches(0.2), y + Inches(0.15),
                    half_w - Inches(0.4), Inches(0.35),
                    title, font_size=16, bold=True, color=t_colors[i])
        add_textbox(slide, x + Inches(0.2), y + Inches(0.6),
                    half_w - Inches(0.4), half_h - Inches(0.8),
                    body, font_size=16, color=CHARCOAL, line_spacing=Pt(26))

    add_footer(slide, 6)


def build_process_flow(prs):
    """7. Process Flow — horizontal numbered steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "Process Flow",
                     "Step-by-step visualization")

    steps = [
        ("Step 1", "Description of\nthe first step"),
        ("Step 2", "Description of\nthe second step"),
        ("Step 3", "Description of\nthe third step"),
        ("Step 4", "Description of\nthe fourth step"),
    ]
    n = len(steps)
    gap = Inches(0.45)
    step_w = (CW - gap * (n - 1)) / n

    for i, (title, desc) in enumerate(steps):
        x = ML + i * (step_w + gap)
        y = BODY_Y + Inches(0.3)
        # Number circle
        cx = x + step_w // 2 - Inches(0.28)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, y,
                                         Inches(0.56), Inches(0.56))
        circle.fill.solid()
        circle.fill.fore_color.rgb = NAVY
        circle.line.fill.background()
        add_textbox(slide, cx, y + Inches(0.08), Inches(0.56), Inches(0.4),
                    str(i + 1), font_name=FONT_NUM, font_size=18, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        # Title
        add_textbox(slide, x, y + Inches(0.65), step_w, Inches(0.55),
                    title, font_size=16, bold=True, color=NAVY,
                    align=PP_ALIGN.CENTER)
        # Description
        add_textbox(slide, x, y + Inches(1.3), step_w, Inches(0.55),
                    desc, font_size=14, color=GREY_MED, align=PP_ALIGN.CENTER)
        # Arrow
        if i < n - 1:
            arrow_x = x + step_w + Inches(0.05)
            add_rect(slide, arrow_x, y + Inches(0.22), Inches(0.35), Pt(1.5), fill=BLUE)

    add_footer(slide, 7)


def build_table(prs):
    """8. Table slide — McKinsey-style table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "Data Table",
                     "Structured data presentation")

    headers = ["Category", "Metric", "Status", "Notes"]
    rows = [
        ["Item A", "Value 1", "Active", "Supporting detail"],
        ["Item B", "Value 2", "Planned", "Supporting detail"],
        ["Item C", "Value 3", "Target", "Supporting detail"],
        ["Item D", "Value 4", "Active", "Supporting detail"],
    ]

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, ML, BODY_Y + Inches(0.1), CW, Inches(0.01))
    table = table_shape.table
    col_w = CW // n_cols
    for i in range(n_cols):
        table.columns[i].width = col_w

    # Header
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_BODY
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE

    # Data
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = GREY_BG if ri % 2 == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_BODY
                p.font.size = Pt(16)
                p.font.color.rgb = CHARCOAL

    add_footer(slide, 8)


def build_full_image(prs):
    """9. Full Image slide — image placeholder fills body."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "Full Image / Diagram Slide",
                     "Visual-first layout for architecture or data visualization")

    add_rect(slide, ML, BODY_Y, CW, BODY_H - Inches(0.15),
             fill=RGBColor(0xE8, 0xEE, 0xF4), line_color=GREY_LINE, line_width=Pt(0.5))
    add_textbox(slide, ML, BODY_Y + BODY_H // 2 - Inches(0.3), CW, Inches(0.6),
                "[Full-width Image / Diagram / Architecture]",
                font_size=18, color=GREY_MED, align=PP_ALIGN.CENTER)

    add_footer(slide, 9)


def build_closing(prs):
    """10. Closing slide — dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    add_rect(slide, Inches(2.5), Inches(4.0), Inches(8.3), Pt(2.5), fill=BLUE)

    add_textbox(slide, Inches(2.5), Inches(2.0), Inches(8.3), Inches(1.5),
                "Next Steps",
                font_name=FONT_TITLE, font_size=36, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(1.2),
                "Action Item 1  |  Action Item 2  |  Action Item 3",
                font_size=16, color=GREY_LIGHT, align=PP_ALIGN.CENTER)

    add_footer(slide, 10)


# ===========================================================================
# Main — Generate template
# ===========================================================================

def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    build_title(prs)
    build_exec_summary(prs)
    build_content(prs)
    build_content_with_image(prs)
    build_comparison(prs)
    build_quadrant(prs)
    build_process_flow(prs)
    build_table(prs)
    build_full_image(prs)
    build_closing(prs)

    out = Path(__file__).parent / "template.pptx"
    prs.save(str(out))
    print(f"[OK] Generated template: {out} ({len(prs.slides)} slides)")
    return out


if __name__ == "__main__":
    main()
