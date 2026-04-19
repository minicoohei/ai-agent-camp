#!/usr/bin/env python3
"""
PPTX構造解析スクリプト

PowerPointファイルの構造（スライド、図形、プレースホルダー、テキスト）を
解析し、JSON形式で出力します。

Usage:
    python analyze_pptx.py <pptx_path> [--output-dir PATH] [--with-images] [--with-gemini]

Examples:
    # 基本的な構造解析
    python analyze_pptx.py template.pptx
    
    # 画像生成付き
    python analyze_pptx.py template.pptx --with-images
    
    # Geminiによる意味解析付き
    python analyze_pptx.py template.pptx --with-gemini
"""

import argparse
import json
import os
import platform
import shutil
import subprocess
from pathlib import Path
from dataclasses import dataclass, asdict
from typing import Optional, List

from pptx import Presentation
from pptx.util import Emu

# Gemini API（オプション）
try:
    from google import genai
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False

# PIL（画像処理用、オプション）
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False


@dataclass
class ShapeInfo:
    """図形情報"""
    shape_id: int
    name: str
    shape_type: str
    left: int
    top: int
    width: int
    height: int
    text: str
    has_text_frame: bool
    is_placeholder: bool
    placeholder_type: Optional[str] = None


@dataclass
class SlideInfo:
    """スライド情報"""
    slide_index: int
    layout_name: str
    shapes: List[ShapeInfo]
    image_path: Optional[str] = None


def parse_args():
    parser = argparse.ArgumentParser(description="PPTX構造解析ツール")
    parser.add_argument("pptx_path", type=Path, help="解析対象のPPTXファイル")
    parser.add_argument(
        "--output-dir", "-o", type=Path, default=None,
        help="出力ディレクトリ（デフォルト: PPTXと同じディレクトリ）"
    )
    parser.add_argument(
        "--with-images", action="store_true",
        help="スライド画像を生成（LibreOffice/qlmanage必要）"
    )
    parser.add_argument(
        "--with-gemini", action="store_true",
        help="Geminiによる意味解析を実行"
    )
    parser.add_argument(
        "--format", choices=["json", "txt", "both"], default="both",
        help="出力形式（デフォルト: both）"
    )
    return parser.parse_args()


def emu_to_inches(emu: int) -> float:
    """EMU単位をインチに変換"""
    return emu / 914400


def extract_shape_info(shape) -> ShapeInfo:
    """図形から情報を抽出"""
    text = ""
    has_text = False
    try:
        has_text = shape.has_text_frame
        if has_text:
            texts = []
            for paragraph in shape.text_frame.paragraphs:
                para_text = "".join(run.text for run in paragraph.runs)
                if para_text.strip():
                    texts.append(para_text.strip())
            text = "\n".join(texts)
    except:
        has_text = False
    
    is_placeholder = False
    placeholder_type = None
    try:
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            is_placeholder = True
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                placeholder_type = str(shape.placeholder_format.type)
    except:
        pass
    
    shape_type = type(shape).__name__
    left = getattr(shape, 'left', 0) or 0
    top = getattr(shape, 'top', 0) or 0
    width = getattr(shape, 'width', 0) or 0
    height = getattr(shape, 'height', 0) or 0
    
    return ShapeInfo(
        shape_id=shape.shape_id,
        name=shape.name,
        shape_type=shape_type,
        left=left,
        top=top,
        width=width,
        height=height,
        text=text[:200] if text else "",
        has_text_frame=has_text,
        is_placeholder=is_placeholder,
        placeholder_type=placeholder_type
    )


def extract_slide_structure(pptx_path: Path) -> List[SlideInfo]:
    """PPTXからスライド構造を抽出"""
    prs = Presentation(pptx_path)
    slides_info = []
    
    for idx, slide in enumerate(prs.slides):
        layout_name = "Unknown"
        try:
            layout_name = slide.slide_layout.name
        except:
            pass
        
        shapes_info = []
        for shape in slide.shapes:
            try:
                shape_info = extract_shape_info(shape)
                shapes_info.append(shape_info)
            except Exception as e:
                print(f"  ⚠️ 図形抽出エラー (slide {idx+1}): {e}")
        
        slides_info.append(SlideInfo(
            slide_index=idx,
            layout_name=layout_name,
            shapes=shapes_info
        ))
    
    return slides_info


def _find_soffice() -> Optional[str]:
    """LibreOffice (soffice) の実行パスを探す"""
    soffice_cmd = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice_cmd:
        return soffice_cmd
    # Windows のデフォルトインストールパスを確認
    if platform.system() == "Windows":
        for candidate in [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]:
            if Path(candidate).exists():
                return candidate
    return None


def pptx_to_images(pptx_path: Path, output_dir: Path) -> List[Path]:
    """PPTXをPDF経由で画像に変換"""
    pptx_path = pptx_path.resolve()
    output_dir = output_dir.resolve()

    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTXファイルが見つかりません: {pptx_path}")
    if not pptx_path.is_file():
        raise ValueError(f"PPTXファイルではありません: {pptx_path}")

    output_dir.mkdir(parents=True, exist_ok=True)
    images = []
    pdf_path = output_dir / f"{pptx_path.stem}.pdf"

    # LibreOfficeでPDF変換
    soffice_cmd = _find_soffice()
    if soffice_cmd:
        try:
            # 古いPDFが残っていると再利用されるため削除
            if pdf_path.exists():
                pdf_path.unlink()

            result = subprocess.run([
                soffice_cmd, "--headless", "--convert-to", "pdf",
                "--outdir", str(output_dir), str(pptx_path)
            ], capture_output=True, text=True, timeout=120)

            if result.returncode != 0:
                print(f"  ⚠️ LibreOffice PDF変換失敗: {result.stderr.strip()}")
            elif pdf_path.exists():
                print(f"  ✅ PDF変換成功: {pdf_path}")
                try:
                    from pdf2image import convert_from_path
                    pdf_images = convert_from_path(pdf_path, dpi=150)
                    for i, img in enumerate(pdf_images):
                        img_path = output_dir / f"slide_{i+1:02d}.png"
                        img.save(img_path, "PNG")
                        images.append(img_path)
                    return images
                except ImportError:
                    print("  ⚠️ pdf2imageがインストールされていません")
        except (subprocess.TimeoutExpired, FileNotFoundError):
            pass

    # macOS Quick Look (macOS専用)
    if platform.system() == "Darwin":
        try:
            # 古いQuick Look出力を掃除
            for old in output_dir.glob(f"{pptx_path.stem}*.png"):
                old.unlink()

            result = subprocess.run([
                "qlmanage", "-t", "-s", "1200", "-o", str(output_dir), str(pptx_path)
            ], capture_output=True, text=True, timeout=30)

            if result.returncode != 0:
                print(f"  ⚠️ Quick Look失敗: {result.stderr.strip()}")
            else:
                ql_files = list(output_dir.glob(f"{pptx_path.stem}*.png"))
                if ql_files:
                    for i, ql_file in enumerate(sorted(ql_files)):
                        new_path = output_dir / f"slide_{i+1:02d}.png"
                        if ql_file != new_path:
                            ql_file.rename(new_path)
                        images.append(new_path)
                    print(f"  ✅ Quick Look画像生成成功: {len(images)}枚")
                    return images
        except (subprocess.TimeoutExpired, FileNotFoundError):
            pass

    print("  ⚠️ 画像変換できませんでした")
    return images


def format_structure_text(slides_info: List[SlideInfo]) -> str:
    """LLM/人間可読なテキスト形式に変換"""
    output = []
    
    for slide in slides_info:
        output.append(f"\n=== Slide {slide.slide_index + 1} (Layout: {slide.layout_name}) ===")
        
        for shape in slide.shapes:
            pos = f"({emu_to_inches(shape.left):.1f}\", {emu_to_inches(shape.top):.1f}\")"
            size = f"{emu_to_inches(shape.width):.1f}\" x {emu_to_inches(shape.height):.1f}\""
            
            output.append(f"  [{shape.shape_id}] {shape.name}")
            output.append(f"      Type: {shape.shape_type}, Pos: {pos}, Size: {size}")
            
            if shape.is_placeholder:
                output.append(f"      Placeholder: {shape.placeholder_type}")
            
            if shape.text:
                text_preview = shape.text[:50].replace('\n', ' ')
                suffix = "..." if len(shape.text) > 50 else ""
                output.append(f"      Text: \"{text_preview}{suffix}\"")
    
    return "\n".join(output)


def analyze_with_gemini(slides_info: List[SlideInfo], image_paths: List[Path]) -> dict:
    """Geminiを使用してスライドの意味を解析"""
    if not GEMINI_AVAILABLE:
        print("  ⚠️ google-genai がインストールされていません")
        return {}

    api_key = os.environ.get("GEMINI_API_KEY") or os.environ.get("GOOGLE_API_KEY")
    if not api_key:
        print("  ⚠️ GEMINI_API_KEY が設定されていません")
        return {}

    client = genai.Client(api_key=api_key)
    
    print("\n🤖 Geminiによるスライド意味解析")
    
    results = {"slides": [], "recommended_usage": {}}
    
    for slide in slides_info:
        slide_idx = slide.slide_index
        print(f"   Slide {slide_idx + 1} を解析中...")
        
        structure_text = f"Layout: {slide.layout_name}\nShapes ({len(slide.shapes)}):\n"
        for shape in slide.shapes:
            pos = f"({emu_to_inches(shape.left):.1f}\", {emu_to_inches(shape.top):.1f}\")"
            size = f"{emu_to_inches(shape.width):.1f}\" x {emu_to_inches(shape.height):.1f}\""
            structure_text += f"  - [{shape.shape_id}] {shape.name}: {shape.shape_type}\n"
            structure_text += f"    Position: {pos}, Size: {size}\n"
            if shape.text:
                structure_text += f"    Text: \"{shape.text[:80]}...\"\n" if len(shape.text) > 80 else f"    Text: \"{shape.text}\"\n"
        
        prompt = f"""このPowerPointスライドを分析してください。

【構造情報】
{structure_text}

以下の形式でJSON出力してください（説明文不要）:
{{
    "slide_role": "タイトル/目次/データ表示/箇条書き/まとめ/その他",
    "recommended_for": "タイトルスライド/グラフ表示/表/箇条書き/まとめ/汎用",
    "elements": [
        {{
            "shape_id": 123,
            "purpose": "メインタイトル/サブタイトル/グラフエリア/テーブル/本文/フッター",
            "editable": true
        }}
    ]
}}
"""
        
        try:
            content = [prompt]
            if PIL_AVAILABLE and slide_idx < len(image_paths) and image_paths[slide_idx].exists():
                img = Image.open(image_paths[slide_idx])
                content.append(img)

            response = client.models.generate_content(
                model="gemini-3-flash-preview",
                contents=content
            )
            result_text = response.text.strip()
            
            if result_text.startswith("```"):
                lines = result_text.split("\n")
                result_text = "\n".join(lines[1:-1])
            
            slide_analysis = json.loads(result_text)
            slide_analysis["slide_index"] = slide_idx
            results["slides"].append(slide_analysis)
            
            print(f"      → 役割: {slide_analysis.get('slide_role', 'N/A')}")
            
        except Exception as e:
            print(f"      ⚠️ 解析エラー: {e}")
            results["slides"].append({"slide_index": slide_idx, "error": str(e)})
    
    for slide_result in results["slides"]:
        if "recommended_for" in slide_result and "error" not in slide_result:
            rec = slide_result["recommended_for"]
            if rec not in results["recommended_usage"]:
                results["recommended_usage"][rec] = []
            results["recommended_usage"][rec].append(slide_result["slide_index"])
    
    return results


def save_results(
    slides_info: List[SlideInfo],
    output_dir: Path,
    pptx_name: str,
    output_format: str,
    image_paths: List[Path] = None,
    gemini_results: dict = None
):
    """解析結果を保存"""
    output_dir.mkdir(parents=True, exist_ok=True)
    
    if image_paths:
        for i, slide in enumerate(slides_info):
            if i < len(image_paths):
                slide.image_path = str(image_paths[i].name)
    
    result = {
        "source_file": pptx_name,
        "total_slides": len(slides_info),
        "slides": []
    }
    
    for slide in slides_info:
        slide_dict = {
            "slide_index": slide.slide_index,
            "layout_name": slide.layout_name,
            "image_path": slide.image_path,
            "shapes": [asdict(s) for s in slide.shapes]
        }
        result["slides"].append(slide_dict)
    
    if gemini_results:
        result["gemini_analysis"] = gemini_results
    
    # JSON保存
    if output_format in ["json", "both"]:
        json_path = output_dir / f"{pptx_name}_structure.json"
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
        print(f"   📄 JSON: {json_path}")
    
    # テキスト保存
    if output_format in ["txt", "both"]:
        txt_path = output_dir / f"{pptx_name}_structure.txt"
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write(format_structure_text(slides_info))
        print(f"   📄 TXT:  {txt_path}")


def main():
    args = parse_args()
    
    print("=" * 60)
    print("PPTX構造解析ツール")
    print("=" * 60)
    
    if not args.pptx_path.exists():
        print(f"❌ ファイルが見つかりません: {args.pptx_path}")
        return 1
    
    output_dir = args.output_dir or args.pptx_path.parent / "pptx_analysis"
    
    print(f"\n📂 対象: {args.pptx_path}")
    print(f"📁 出力: {output_dir}")
    
    # 画像生成
    image_paths = []
    if args.with_images:
        print("\n📸 画像生成中...")
        images_dir = output_dir / "images"
        image_paths = pptx_to_images(args.pptx_path, images_dir)
    
    # 構造抽出
    print("\n🔍 構造解析中...")
    slides_info = extract_slide_structure(args.pptx_path)
    print(f"   スライド数: {len(slides_info)}")
    print(f"   総図形数: {sum(len(s.shapes) for s in slides_info)}")
    
    # Gemini解析
    gemini_results = None
    if args.with_gemini:
        gemini_results = analyze_with_gemini(slides_info, image_paths)
    
    # 結果保存
    print("\n💾 結果保存中...")
    save_results(
        slides_info, output_dir, args.pptx_path.stem,
        args.format, image_paths, gemini_results
    )
    
    # サマリー
    print("\n" + "=" * 60)
    print("✅ 解析完了")
    print("=" * 60)
    
    for slide in slides_info:
        text_shapes = sum(1 for s in slide.shapes if s.has_text_frame and s.text)
        placeholders = sum(1 for s in slide.shapes if s.is_placeholder)
        print(f"   Slide {slide.slide_index + 1}: {slide.layout_name}")
        print(f"      図形: {len(slide.shapes)}, テキスト有: {text_shapes}, PH: {placeholders}")
    
    return 0


if __name__ == "__main__":
    exit(main())
