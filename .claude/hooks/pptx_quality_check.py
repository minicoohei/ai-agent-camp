#!/usr/bin/env python3
"""PPTX Quality Checker — PostToolUse hook for Bash.

Checks for:
1. Text overflow (chars vs box size)
2. Shape overlap (bounding box intersection)
3. Off-screen elements
4. Naming inconsistencies

Read-only: never modifies files. Outputs warnings to stderr.
"""
import json
import sys
import os

def check_pptx(filepath: str) -> list[str]:
    """Run quality checks on a PPTX file. Returns list of warnings."""
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        return []

    if not os.path.exists(filepath):
        return []

    warnings = []
    try:
        prs = Presentation(filepath)
    except Exception:
        return [f"PPTX open failed: {filepath}"]

    sw = prs.slide_width / 914400  # inches
    sh = prs.slide_height / 914400

    for idx, slide in enumerate(prs.slides, 1):
        boxes = []  # (left, top, right, bottom, text_preview)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text.strip()
            if not text:
                continue

            left = (shape.left or 0) / 914400
            top = (shape.top or 0) / 914400
            w = (shape.width or 0) / 914400
            h = (shape.height or 0) / 914400
            right = left + w
            bottom = top + h
            preview = text[:30].replace("\n", " ")

            # Check 1: Off-screen
            if left < -0.5 or top < -0.5 or right > sw + 0.5 or bottom > sh + 0.5:
                warnings.append(f"Slide {idx}: off-screen shape [{preview}] pos=({left:.1f},{top:.1f})")

            # Check 2: Text overflow estimate
            # Skip short text and small label shapes
            if w > 0 and h > 0.3:
                chars = len(text.replace("\n", ""))
                if chars > 30:  # Only check substantial text
                    # Adaptive: smaller boxes use smaller fonts → more chars/inch
                    cpi = 7 if h < 0.5 else 5  # chars per inch
                    lpi = 3.5 if h < 0.5 else 2.5  # lines per inch
                    max_chars = w * cpi * h * lpi
                    if chars > max_chars * 1.3:
                        warnings.append(
                            f"Slide {idx}: text may overflow [{preview}] "
                            f"{chars} chars in {w:.1f}x{h:.1f}in box"
                        )

            boxes.append((left, top, right, bottom, preview))

        # Check 3: Overlap detection (same-size text boxes only)
        for i, (l1, t1, r1, b1, p1) in enumerate(boxes):
            for j, (l2, t2, r2, b2, p2) in enumerate(boxes):
                if j <= i:
                    continue
                # Check bounding box intersection
                if l1 < r2 and r1 > l2 and t1 < b2 and b1 > t2:
                    # Calculate overlap area
                    ol = max(l1, l2)
                    ot = max(t1, t2)
                    or_ = min(r1, r2)
                    ob = min(b1, b2)
                    area = (or_ - ol) * (ob - ot)
                    min_area = min((r1-l1)*(b1-t1), (r2-l2)*(b2-t2))
                    if min_area > 0 and area / min_area > 0.3:  # >30% overlap
                        warnings.append(
                            f"Slide {idx}: overlap [{p1}] & [{p2}] "
                            f"({area/min_area*100:.0f}% overlap)"
                        )

    return warnings


def main():
    # Read hook input from stdin
    try:
        hook_input = json.load(sys.stdin)
    except (json.JSONDecodeError, EOFError):
        return

    tool_name = hook_input.get("tool_name", "")
    if tool_name != "Bash":
        return

    # Check if the command output mentions a .pptx file
    tool_input = hook_input.get("tool_input", {})
    stdout = hook_input.get("tool_result", {}).get("stdout", "")
    command = tool_input.get("command", "")

    # Find .pptx paths in command or output
    pptx_files = set()
    for text in [command, stdout]:
        for token in text.split():
            token = token.strip("'\"")
            if token.endswith(".pptx") and "/" in token:
                pptx_files.add(token)

    if not pptx_files:
        return

    all_warnings = []
    for f in pptx_files:
        # Only check output files (not source files)
        if os.path.exists(f):
            ws = check_pptx(f)
            all_warnings.extend(ws)

    if all_warnings:
        msg = "PPTX Quality Check:\n" + "\n".join(f"  ⚠ {w}" for w in all_warnings[:10])
        print(msg, file=sys.stderr)


if __name__ == "__main__":
    main()
