#!/usr/bin/env python3
"""
PPTX Converter - メインCLI

4つのサブコマンド:
  convert: 元PPTXをコピーして新トピックで書き換え（1コマンド完結）
  extract: マッピングYAML生成のみ
  build:   マッピングYAML + データYAML → PPTX書き換え
  deck:    ゼロからPPTXデッキ生成

使用例:
  python pptx_converter.py convert source.pptx --topic "2026年Q1報告" -o output.pptx
  python pptx_converter.py extract source.pptx -o mapping.yaml
  python pptx_converter.py build source.pptx mapping.yaml --data data.yaml -o output.pptx
  python pptx_converter.py deck --topic "AI活用提案" --type proposal -o output.pptx
"""

import argparse
import json
import os
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional

import yaml

# ─── パス設定 ─────────────────────────────────────────────────
_SCRIPT_DIR = Path(__file__).resolve().parent
_TOOLS_DIR = _SCRIPT_DIR.parent.parent.parent / "tools"
sys.path.insert(0, str(_SCRIPT_DIR))
sys.path.insert(0, str(_TOOLS_DIR))

from extractor import PPTXExtractor
from mapper import PPTXMapper
from generator import PPTXGenerator

try:
    from bootcamp_utils import get_client, get_flash_model
except ImportError:
    def get_client():
        return None
    def get_flash_model() -> str:
        return "gemini-3-flash-preview"


# ─── デッキ生成用スタイル定義 ─────────────────────────────────
try:
    from pptx.util import Pt, Inches, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx import Presentation
except ImportError:
    RGBColor = None
    PP_ALIGN = None

STYLES = {
    "corporate": {
        "primary": (0x25, 0x63, 0xEB),
        "secondary": (0x1E, 0x40, 0xAF),
        "accent": (0xFB, 0xBF, 0x24),
        "bg": (0xFF, 0xFF, 0xFF),
        "text": (0x1F, 0x2D, 0x3D),
        "light_bg": (0xF0, 0xF4, 0xF8),
        "title_font": "Meiryo UI",
        "body_font": "Meiryo",
        "title_size": 36,
        "body_size": 16,
    },
    "minimal": {
        "primary": (0x11, 0x18, 0x27),
        "secondary": (0x37, 0x41, 0x51),
        "accent": (0xEF, 0x44, 0x44),
        "bg": (0xFF, 0xFF, 0xFF),
        "text": (0x11, 0x18, 0x27),
        "light_bg": (0xF9, 0xFA, 0xFB),
        "title_font": "Meiryo UI",
        "body_font": "Meiryo",
        "title_size": 32,
        "body_size": 14,
    },
    "creative": {
        "primary": (0x7C, 0x3A, 0xED),
        "secondary": (0xEC, 0x48, 0x99),
        "accent": (0x06, 0xB6, 0xD4),
        "bg": (0xFF, 0xFF, 0xFF),
        "text": (0x1E, 0x1B, 0x4B),
        "light_bg": (0xF5, 0xF3, 0xFF),
        "title_font": "Meiryo UI",
        "body_font": "Meiryo",
        "title_size": 36,
        "body_size": 16,
    },
    "academic": {
        "primary": (0x1E, 0x3A, 0x5F),
        "secondary": (0x2C, 0x5F, 0x7C),
        "accent": (0xD4, 0x8B, 0x2C),
        "bg": (0xFF, 0xFF, 0xFF),
        "text": (0x2D, 0x33, 0x36),
        "light_bg": (0xF8, 0xF9, 0xFA),
        "title_font": "Meiryo UI",
        "body_font": "Meiryo",
        "title_size": 32,
        "body_size": 14,
    },
}


# =====================================================================
# convert コマンド
# =====================================================================
def cmd_convert(args):
    """元PPTXをコピー → 全要素抽出 → Gemini意味解析 → 新コンテンツ生成 → 書き換え → 保存"""
    source = args.source
    topic = args.topic
    output = args.output

    print(f"[convert] 元PPTX: {source}")
    print(f"[convert] トピック: {topic}")
    print(f"[convert] 出力先: {output}")

    # Step 1: 抽出
    print("\n[convert] Step 1/4: 要素抽出中...")
    assets_dir = str(Path(output).parent / "extracted_assets")
    extractor = PPTXExtractor(source, assets_dir=assets_dir)
    extracted = extractor.extract_all()

    slide_count = len(extracted["slides"])
    element_count = sum(len(s["elements"]) for s in extracted["slides"])
    print(f"[convert]   {slide_count}スライド, {element_count}要素を抽出")

    # Step 2: マッピング生成
    print("\n[convert] Step 2/4: マッピング生成中...")
    mapper = PPTXMapper(extracted)
    mapping = mapper.generate_mapping()

    placeholder_count = len(mapping.get("placeholders", []))
    print(f"[convert]   {placeholder_count}個のプレースホルダーを生成")

    # Step 3: Geminiで新コンテンツ生成
    print("\n[convert] Step 3/4: 新コンテンツ生成中...")
    new_data = _generate_new_content(mapping, topic)

    if new_data:
        print(f"[convert]   {len(new_data)}個のプレースホルダーを生成")
    else:
        print("[convert]   コンテンツ生成に失敗しました")
        return

    # Step 4: 書き換え
    print("\n[convert] Step 4/4: PPTX書き換え中...")
    gen = PPTXGenerator(source, output)
    stats = gen.apply_replacements(mapping, new_data)

    print(f"\n[convert] === 完了 ===")
    print(f"[convert]   書き換え: {stats['replaced']}要素")
    print(f"[convert]   スキップ: {stats['skipped']}要素")
    if stats["errors"]:
        print(f"[convert]   エラー: {len(stats['errors'])}件")
        for err in stats["errors"][:5]:
            print(f"[convert]     - {err}")
    print(f"[convert]   出力: {Path(output).resolve()}")


def _generate_new_content(mapping: dict, topic: str) -> Optional[dict]:
    """Geminiで新コンテンツを生成する。"""
    client = get_client()
    if not client:
        print("[convert] Gemini APIキーが設定されていません。手動でデータYAMLを作成してください。")
        return None

    placeholders = mapping.get("placeholders", [])
    if not placeholders:
        return {}

    # プロンプト構築
    placeholder_list = []
    for ph in placeholders:
        entry = {
            "key": ph["key"],
            "type": ph["type"],
            "role": ph.get("role", ""),
            "current": ph.get("current", ""),
        }
        # チャート・テーブルの構造情報を追加
        if ph["type"] == "chart":
            # マッピングから構造情報を検索
            for slide in mapping.get("slides", []):
                for elem in slide.get("elements", []):
                    if elem.get("placeholder") == ph["key"]:
                        value = elem.get("value", {})
                        if isinstance(value, dict):
                            cats = value.get("categories", [])
                            series = value.get("series", [])
                            entry["structure"] = {
                                "categories_count": len(cats),
                                "series": [s.get("name", "") for s in series] if isinstance(series, list) else [],
                            }
                        break
        elif ph["type"] == "table":
            for slide in mapping.get("slides", []):
                for elem in slide.get("elements", []):
                    if elem.get("placeholder") == ph["key"]:
                        value = elem.get("value", [])
                        if isinstance(value, list):
                            entry["structure"] = {
                                "rows": len(value),
                                "cols": len(value[0]) if value else 0,
                            }
                        break

        placeholder_list.append(entry)

    prompt = _build_content_generation_prompt(placeholder_list, topic)

    try:
        response = client.models.generate_content(
            model=get_flash_model(),
            contents=[prompt],
        )
        raw_text = response.text if hasattr(response, "text") else str(response)
        return _parse_content_response(raw_text)
    except Exception as e:
        print(f"[convert] Gemini コンテンツ生成エラー: {e}")
        return None


def _build_content_generation_prompt(placeholders: list, topic: str) -> str:
    """コンテンツ生成用のGeminiプロンプトを構築。"""
    ph_json = json.dumps(placeholders, ensure_ascii=False, indent=2)

    return (
        f"以下のプレゼンテーションテンプレートを「{topic}」というトピックで書き換えます。\n"
        f"各プレースホルダーに対して、適切な新しい値を生成してください。\n\n"
        f"プレースホルダー一覧:\n{ph_json}\n\n"
        f"重要なルール:\n"
        f"- テキストは role と current の文字数を参考に適切な長さで生成\n"
        f"- チャートは categories + series の構造を維持（同じ数のカテゴリ・シリーズ）\n"
        f"- テーブルは行列数を維持\n"
        f"- 画像プレースホルダーは無視（値なしでOK）\n"
        f"- 日本語で生成\n"
        f"- 内容は「{topic}」に合わせて統一感を持たせる\n\n"
        f"回答は以下のJSON形式のみで返してください（説明文不要）:\n"
        f"```json\n"
        f'{{\n'
        f'  "{{{{placeholder_key}}}}": "新しいテキスト",\n'
        f'  "{{{{chart_key}}}}": {{\n'
        f'    "categories": ["A","B","C"],\n'
        f'    "series": [{{"name":"系列1","values":[1,2,3]}}]\n'
        f'  }},\n'
        f'  "{{{{table_key}}}}": [["ヘッダ1","ヘッダ2"],["値1","値2"]]\n'
        f'}}\n'
        f"```"
    )


def _parse_content_response(text: str) -> Optional[dict]:
    """Geminiレスポンスからコンテンツ辞書を抽出。"""
    import re

    if not text:
        return None

    # ```json ... ``` ブロック抽出
    match = re.search(r"```(?:json)?\s*\n?(.*?)\n?\s*```", text, re.DOTALL)
    json_str = match.group(1).strip() if match else text.strip()

    # JSON辞書部分を抽出
    brace_match = re.search(r"\{.*\}", json_str, re.DOTALL)
    if brace_match:
        json_str = brace_match.group(0)

    try:
        return json.loads(json_str)
    except json.JSONDecodeError:
        return None


# =====================================================================
# extract コマンド
# =====================================================================
def cmd_extract(args):
    """マッピングYAML生成のみ。"""
    source = args.source
    output = args.output

    print(f"[extract] 元PPTX: {source}")
    print(f"[extract] 出力先: {output}")

    # 画像抽出先
    assets_dir = str(Path(output).parent / "extracted_assets")

    # 抽出
    print("\n[extract] Step 1/2: 要素抽出中...")
    extractor = PPTXExtractor(source, assets_dir=assets_dir)
    extracted = extractor.extract_all()

    slide_count = len(extracted["slides"])
    element_count = sum(len(s["elements"]) for s in extracted["slides"])
    print(f"[extract]   {slide_count}スライド, {element_count}要素を抽出")

    # マッピング生成
    print("\n[extract] Step 2/2: マッピング生成中...")
    mapper = PPTXMapper(extracted)

    if args.no_gemini:
        mapper._gemini_available = False
        mapper._client = None
        print("[extract]   Gemini解析を無効にしました")

    mapping = mapper.generate_mapping()
    mapper.save_yaml(mapping, output)

    # アセット保存
    mapper.save_assets(extracted, assets_dir)

    # サマリー
    placeholder_count = len(mapping.get("placeholders", []))
    warning_count = len(mapping.get("warnings", []))
    print(f"\n[extract] === 完了 ===")
    print(f"[extract]   プレースホルダー: {placeholder_count}個")
    if warning_count:
        print(f"[extract]   警告: {warning_count}件")
    print(f"[extract]   マッピング: {Path(output).resolve()}")
    print(f"[extract]   アセット: {Path(assets_dir).resolve()}")


# =====================================================================
# build コマンド
# =====================================================================
def cmd_build(args):
    """マッピングYAML + データYAML → PPTX書き換え。"""
    source = args.source
    mapping_path = args.mapping
    data_path = args.data
    output = args.output

    print(f"[build] 元PPTX: {source}")
    print(f"[build] マッピング: {mapping_path}")
    print(f"[build] データ: {data_path}")
    print(f"[build] 出力先: {output}")

    # ファイル読み込み
    with open(mapping_path, "r", encoding="utf-8") as f:
        mapping = yaml.safe_load(f)

    with open(data_path, "r", encoding="utf-8") as f:
        new_data = yaml.safe_load(f)

    # 書き換え
    gen = PPTXGenerator(source, output)
    stats = gen.apply_replacements(mapping, new_data)

    print(f"\n[build] === 完了 ===")
    print(f"[build]   書き換え: {stats['replaced']}要素")
    print(f"[build]   スキップ: {stats['skipped']}要素")
    if stats["errors"]:
        print(f"[build]   エラー: {len(stats['errors'])}件")
        for err in stats["errors"][:5]:
            print(f"[build]     - {err}")
    print(f"[build]   出力: {Path(output).resolve()}")


# =====================================================================
# deck コマンド
# =====================================================================
def cmd_deck(args):
    """ゼロからPPTXデッキ生成。"""
    topic = args.topic
    deck_type = args.type
    style_name = args.style
    slides_count = args.slides
    audience = args.audience
    output = args.output
    outline_only = args.outline_only

    print(f"[deck] トピック: {topic}")
    print(f"[deck] タイプ: {deck_type}")
    print(f"[deck] スタイル: {style_name}")
    print(f"[deck] 対象: {audience}")
    print(f"[deck] 出力先: {output}")

    # Step 1: Geminiでアウトライン生成
    print("\n[deck] Step 1: アウトライン生成中...")
    outline = _generate_deck_outline(topic, deck_type, slides_count, audience)

    if not outline:
        print("[deck] アウトライン生成に失敗しました")
        return

    total_slides = len(outline.get("slides", []))
    print(f"[deck]   {total_slides}スライドのアウトラインを生成")

    if outline_only:
        # アウトラインのみ出力
        outline_path = Path(output).with_suffix(".yaml")
        with open(outline_path, "w", encoding="utf-8") as f:
            yaml.dump(outline, f, allow_unicode=True, default_flow_style=False, sort_keys=False)
        print(f"[deck] アウトラインを保存: {outline_path.resolve()}")
        return

    # Step 2: PPTX生成
    print("\n[deck] Step 2: PPTX生成中...")
    style = STYLES.get(style_name, STYLES["corporate"])
    _build_deck_pptx(outline, style, output)

    print(f"\n[deck] === 完了 ===")
    print(f"[deck]   スライド数: {total_slides}")
    print(f"[deck]   出力: {Path(output).resolve()}")


def _generate_deck_outline(
    topic: str, deck_type: str, slides_count: int, audience: str
) -> Optional[dict]:
    """Geminiでデッキアウトラインを生成。"""
    client = get_client()
    if not client:
        print("[deck] Gemini APIキーが設定されていません")
        return None

    slides_instruction = f"{slides_count}枚" if slides_count > 0 else "AIが最適な枚数を決定（8-15枚推奨）"

    prompt = (
        f"以下のトピックについて、{deck_type}用のプレゼンテーションデッキを設計してください。\n"
        f"対象: {audience}\n\n"
        f"トピック: {topic}\n"
        f"スライド枚数: {slides_instruction}\n\n"
        f"各スライドについて以下を出力してください（JSONのみ、説明文不要）:\n"
        f"```json\n"
        f'{{\n'
        f'  "title": "デッキタイトル",\n'
        f'  "total_slides": 10,\n'
        f'  "slides": [\n'
        f'    {{\n'
        f'      "slide_number": 1,\n'
        f'      "slide_type": "title",\n'
        f'      "title": "...",\n'
        f'      "subtitle": "...",\n'
        f'      "bullets": ["..."],\n'
        f'      "chart_type": "COLUMN_CLUSTERED",\n'
        f'      "chart_data": {{"categories":["A","B"],"series":[{{"name":"X","values":[1,2]}}]}},\n'
        f'      "table_data": [["col1","col2"],["val1","val2"]],\n'
        f'      "speaker_notes": "..."\n'
        f'    }}\n'
        f'  ]\n'
        f'}}\n'
        f"```\n\n"
        f"slide_typeの選択肢: title, section, content, two_column, table, chart, summary, closing\n"
        f"- title: タイトルスライド（title, subtitle必須）\n"
        f"- section: セクション区切り（title, subtitle）\n"
        f"- content: 本文（title, bullets必須、3-5項目）\n"
        f"- two_column: 2列レイアウト（title, bullets を左右に分割）\n"
        f"- table: テーブルスライド（title, table_data必須）\n"
        f"- chart: チャートスライド（title, chart_type, chart_data必須）\n"
        f"- summary: まとめ（title, bullets）\n"
        f"- closing: 締めスライド（title, subtitle）\n\n"
        f"chart_type: COLUMN_CLUSTERED, LINE, PIE, BAR_CLUSTERED から選択\n"
        f"日本語で生成してください。\n"
    )

    try:
        response = client.models.generate_content(
            model=get_flash_model(),
            contents=[prompt],
        )
        raw_text = response.text if hasattr(response, "text") else str(response)
        return _parse_content_response(raw_text)
    except Exception as e:
        print(f"[deck] Gemini アウトライン生成エラー: {e}")
        return None


def _build_deck_pptx(outline: dict, style: dict, output_path: str):
    """アウトラインからPPTXを生成。"""
    if Presentation is None:
        print("[deck] python-pptx が必要です: pip install python-pptx")
        return

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_info in outline.get("slides", []):
        slide_type = slide_info.get("slide_type", "content")
        _add_deck_slide(prs, slide_info, slide_type, style)

    # 出力ディレクトリ作成
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def _add_deck_slide(prs, slide_info: dict, slide_type: str, style: dict):
    """1スライドを追加。"""
    from pptx.util import Pt, Inches, Emu

    # 空白レイアウトを使用
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    title_text = slide_info.get("title", "")
    subtitle_text = slide_info.get("subtitle", "")
    bullets = slide_info.get("bullets", [])
    notes = slide_info.get("speaker_notes", "")

    primary = RGBColor(*style["primary"]) if RGBColor else None
    text_color = RGBColor(*style["text"]) if RGBColor else None
    accent = RGBColor(*style["accent"]) if RGBColor else None
    bg_color = RGBColor(*style.get("light_bg", (0xF0, 0xF4, 0xF8))) if RGBColor else None

    title_font = style.get("title_font", "Meiryo UI")
    body_font = style.get("body_font", "Meiryo")
    title_size = style.get("title_size", 36)
    body_size = style.get("body_size", 16)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    if slide_type in ("title", "closing"):
        # タイトルスライド: 中央配置
        _add_bg_rect(slide, slide_w, slide_h, primary)
        _add_textbox(
            slide, title_text,
            left=Inches(1), top=Inches(2.2), width=Inches(11.333), height=Inches(1.5),
            font_name=title_font, font_size=Pt(44),
            color=RGBColor(0xFF, 0xFF, 0xFF) if RGBColor else None,
            alignment=PP_ALIGN.CENTER if PP_ALIGN else None,
            bold=True,
        )
        if subtitle_text:
            _add_textbox(
                slide, subtitle_text,
                left=Inches(1.5), top=Inches(4), width=Inches(10.333), height=Inches(1),
                font_name=body_font, font_size=Pt(20),
                color=RGBColor(0xE0, 0xE0, 0xE0) if RGBColor else None,
                alignment=PP_ALIGN.CENTER if PP_ALIGN else None,
            )

    elif slide_type == "section":
        # セクション区切り
        _add_bg_rect(slide, slide_w, slide_h, primary)
        _add_textbox(
            slide, title_text,
            left=Inches(1.5), top=Inches(2.5), width=Inches(10.333), height=Inches(1.2),
            font_name=title_font, font_size=Pt(40),
            color=RGBColor(0xFF, 0xFF, 0xFF) if RGBColor else None,
            alignment=PP_ALIGN.CENTER if PP_ALIGN else None,
            bold=True,
        )
        if subtitle_text:
            _add_textbox(
                slide, subtitle_text,
                left=Inches(2), top=Inches(4), width=Inches(9.333), height=Inches(0.8),
                font_name=body_font, font_size=Pt(18),
                color=RGBColor(0xCC, 0xCC, 0xCC) if RGBColor else None,
                alignment=PP_ALIGN.CENTER if PP_ALIGN else None,
            )

    elif slide_type == "chart":
        # チャートスライド
        _add_textbox(
            slide, title_text,
            left=Inches(0.8), top=Inches(0.4), width=Inches(11.733), height=Inches(0.8),
            font_name=title_font, font_size=Pt(title_size),
            color=primary, bold=True,
        )
        chart_data = slide_info.get("chart_data", {})
        chart_type_name = slide_info.get("chart_type", "COLUMN_CLUSTERED")
        _add_chart(slide, chart_type_name, chart_data,
                   left=Inches(1), top=Inches(1.5), width=Inches(11.333), height=Inches(5.2))

    elif slide_type == "table":
        # テーブルスライド
        _add_textbox(
            slide, title_text,
            left=Inches(0.8), top=Inches(0.4), width=Inches(11.733), height=Inches(0.8),
            font_name=title_font, font_size=Pt(title_size),
            color=primary, bold=True,
        )
        table_data = slide_info.get("table_data", [])
        if table_data:
            _add_table(slide, table_data, style,
                       left=Inches(0.8), top=Inches(1.5), width=Inches(11.733), height=Inches(5.2))

    elif slide_type == "two_column":
        # 2列レイアウト
        _add_textbox(
            slide, title_text,
            left=Inches(0.8), top=Inches(0.4), width=Inches(11.733), height=Inches(0.8),
            font_name=title_font, font_size=Pt(title_size),
            color=primary, bold=True,
        )
        mid = len(bullets) // 2
        left_bullets = bullets[:mid] if mid > 0 else bullets[:1]
        right_bullets = bullets[mid:] if mid > 0 else bullets[1:]

        bullet_text_l = "\n".join(f"  {b}" for b in left_bullets)
        bullet_text_r = "\n".join(f"  {b}" for b in right_bullets)

        _add_textbox(
            slide, bullet_text_l,
            left=Inches(0.8), top=Inches(1.5), width=Inches(5.5), height=Inches(5.2),
            font_name=body_font, font_size=Pt(body_size), color=text_color,
        )
        _add_textbox(
            slide, bullet_text_r,
            left=Inches(7), top=Inches(1.5), width=Inches(5.5), height=Inches(5.2),
            font_name=body_font, font_size=Pt(body_size), color=text_color,
        )

    else:
        # content / summary / default
        _add_textbox(
            slide, title_text,
            left=Inches(0.8), top=Inches(0.4), width=Inches(11.733), height=Inches(0.8),
            font_name=title_font, font_size=Pt(title_size),
            color=primary, bold=True,
        )
        if bullets:
            bullet_text = "\n".join(f"  {b}" for b in bullets)
            _add_textbox(
                slide, bullet_text,
                left=Inches(1), top=Inches(1.5), width=Inches(11.333), height=Inches(5.2),
                font_name=body_font, font_size=Pt(body_size), color=text_color,
            )

    # スピーカーノート
    if notes and slide.notes_slide:
        slide.notes_slide.notes_text_frame.text = notes


def _add_textbox(slide, text, left, top, width, height,
                 font_name="Meiryo", font_size=None, color=None,
                 alignment=None, bold=False):
    """テキストボックスを追加。"""
    from pptx.util import Pt

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    para = tf.paragraphs[0]
    para.text = text
    if alignment:
        para.alignment = alignment

    if font_size or font_name or color or bold:
        for run in para.runs:
            if font_name:
                run.font.name = font_name
            if font_size:
                run.font.size = font_size
            if color:
                run.font.color.rgb = color
            if bold:
                run.font.bold = True


def _add_bg_rect(slide, width, height, color):
    """背景矩形を追加。"""
    from pptx.util import Emu

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(0), Emu(0), width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_chart(slide, chart_type_name: str, chart_data: dict,
               left, top, width, height):
    """チャートを追加。"""
    try:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        type_map = {
            "COLUMN_CLUSTERED": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "LINE": XL_CHART_TYPE.LINE,
            "PIE": XL_CHART_TYPE.PIE,
            "BAR_CLUSTERED": XL_CHART_TYPE.BAR_CLUSTERED,
        }
        xl_type = type_map.get(chart_type_name.upper(), XL_CHART_TYPE.COLUMN_CLUSTERED)

        data = CategoryChartData()
        categories = chart_data.get("categories", [])
        if categories:
            data.categories = categories

        for s in chart_data.get("series", []):
            name = s.get("name", "")
            values = []
            for v in s.get("values", []):
                try:
                    values.append(float(v))
                except (ValueError, TypeError):
                    values.append(0)
            data.add_series(name, values)

        slide.shapes.add_chart(xl_type, left, top, width, height, data)

    except Exception as e:
        # チャート追加失敗時はテキストで代替
        _add_textbox(slide, f"[Chart: {chart_type_name}]",
                     left, top, width, height,
                     font_size=Pt(14) if Pt else None)


def _add_table(slide, table_data: list, style: dict,
               left, top, width, height):
    """テーブルを追加。"""
    if not table_data:
        return

    rows = len(table_data)
    cols = len(table_data[0]) if table_data else 0
    if rows == 0 or cols == 0:
        return

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    primary = style.get("primary", (0x25, 0x63, 0xEB))

    for row_idx, row in enumerate(table_data):
        for col_idx, val in enumerate(row):
            if col_idx >= cols:
                break
            cell = table.cell(row_idx, col_idx)
            cell.text = str(val) if val is not None else ""

            # ヘッダー行のスタイル
            if row_idx == 0 and RGBColor:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*primary)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                        run.font.bold = True


# =====================================================================
# メインエントリーポイント
# =====================================================================
def main():
    parser = argparse.ArgumentParser(
        description="PPTX Converter - テンプレート変換 & デッキ生成ツール",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
使用例:
  # テンプレート変換（1コマンド完結）
  %(prog)s convert source.pptx --topic "2026年Q1営業報告" -o output.pptx

  # マッピングYAML生成のみ
  %(prog)s extract source.pptx -o mapping.yaml

  # マッピング + データ → PPTX書き換え
  %(prog)s build source.pptx mapping.yaml --data data.yaml -o output.pptx

  # ゼロからデッキ生成
  %(prog)s deck --topic "AIエージェント活用提案" --type proposal -o output.pptx
        """,
    )
    subparsers = parser.add_subparsers(dest="command", help="サブコマンド")

    # convert
    p_convert = subparsers.add_parser("convert", help="テンプレート変換（1コマンド完結）")
    p_convert.add_argument("source", help="元PPTXファイル")
    p_convert.add_argument("--topic", "-t", required=True, help="新しいトピック")
    p_convert.add_argument("--output", "-o", required=True, help="出力PPTXファイル")

    # extract
    p_extract = subparsers.add_parser("extract", help="マッピングYAML生成")
    p_extract.add_argument("source", help="元PPTXファイル")
    p_extract.add_argument("--output", "-o", default="mapping.yaml", help="出力YAMLファイル")
    p_extract.add_argument("--no-gemini", action="store_true", help="Gemini解析を無効化")

    # build
    p_build = subparsers.add_parser("build", help="マッピング + データ → PPTX")
    p_build.add_argument("source", help="元PPTXファイル")
    p_build.add_argument("mapping", help="マッピングYAMLファイル")
    p_build.add_argument("--data", "-d", required=True, help="データYAMLファイル")
    p_build.add_argument("--output", "-o", required=True, help="出力PPTXファイル")

    # deck
    p_deck = subparsers.add_parser("deck", help="ゼロからデッキ生成")
    p_deck.add_argument("--topic", "-t", required=True, help="トピック")
    p_deck.add_argument(
        "--type", default="auto",
        choices=["auto", "presentation", "proposal", "report", "educational", "pitch"],
        help="デッキタイプ",
    )
    p_deck.add_argument(
        "--style", default="corporate",
        choices=["corporate", "creative", "minimal", "academic"],
        help="デザインスタイル",
    )
    p_deck.add_argument("--slides", type=int, default=0, help="スライド枚数（0=AI決定）")
    p_deck.add_argument("--audience", default="general", help="対象")
    p_deck.add_argument("--output", "-o", required=True, help="出力PPTXファイル")
    p_deck.add_argument("--outline-only", action="store_true", help="アウトラインのみ出力")

    args = parser.parse_args()

    if not args.command:
        parser.print_help()
        sys.exit(1)

    if args.command == "convert":
        cmd_convert(args)
    elif args.command == "extract":
        cmd_extract(args)
    elif args.command == "build":
        cmd_build(args)
    elif args.command == "deck":
        cmd_deck(args)
    else:
        parser.print_help()
        sys.exit(1)


if __name__ == "__main__":
    main()
