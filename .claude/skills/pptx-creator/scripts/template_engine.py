"""
template_engine.py — テンプレートコピー + PH注入 + リッチ要素ディスパッチ

テンプレート PPTX をコピーし、既存スライドを削除してからアウトラインの各スライドを追加。
レイアウトPH注入と rich_elements コード生成のハイブリッド。
"""

import shutil
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

import yaml
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from rich_elements import render as render_rich, available_patterns

# ---------------------------------------------------------------------------
# Catalog
# ---------------------------------------------------------------------------

def load_catalog(path=None):
    """catalog.yaml を読み込み"""
    if path is None:
        path = Path(__file__).resolve().parent.parent / "templates" / "catalog.yaml"
    else:
        path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"Catalog not found: {path}")
    with open(path, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def resolve_template(catalog, template_name):
    """テンプレート設定を解決、ファイル存在確認"""
    templates = catalog.get("templates", {})
    if template_name not in templates:
        raise ValueError(f"Template '{template_name}' not found. Available: {list(templates.keys())}")
    tpl = templates[template_name]
    tpl_dir = Path(__file__).resolve().parent.parent / "templates"
    tpl_file = tpl_dir / tpl["file"]
    tpl["_abs_path"] = tpl_file
    tpl["_name"] = template_name
    return tpl


def resolve_layout(catalog, template_name, slide_type):
    """slide_type → レイアウトキーを解決。"code_gen" or レイアウトキー"""
    mapping = catalog.get("slide_type_mapping", {})
    if slide_type in mapping:
        return mapping[slide_type].get(template_name, "content")
    print(f"[engine] Warning: no mapping for slide_type '{slide_type}', using 'content'", file=sys.stderr)
    return "content"


# ---------------------------------------------------------------------------
# Slide operations
# ---------------------------------------------------------------------------

def delete_all_slides(prs):
    """既存スライドを全削除（レイアウト/マスターは保持）"""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def _get_placeholder_by_idx(slide, idx):
    """スライドからプレースホルダーを idx で取得"""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def set_text_preserving_format(placeholder, text, bullets=None):
    """既存書式を完全保持してテキストを設定"""
    tf = placeholder.text_frame

    # テンプレートの既存書式を取得
    template_font_name = None
    template_font_size = None
    template_font_bold = None
    template_font_color = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        src_run = tf.paragraphs[0].runs[0]
        template_font_name = src_run.font.name
        template_font_size = src_run.font.size
        template_font_bold = src_run.font.bold
        try:
            template_font_color = src_run.font.color.rgb
        except (AttributeError, TypeError):
            pass

    if bullets:
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            # 段落書式をコピー
            if i > 0 and len(tf.paragraphs) > 1:
                try:
                    src_pPr = tf.paragraphs[0]._p.get_or_add_pPr()
                    dst_pPr = p._p.get_or_add_pPr()
                    for attr in src_pPr.attrib:
                        dst_pPr.set(attr, src_pPr.get(attr))
                except Exception:
                    pass
            run = p.add_run()
            run.text = bullet
            # テンプレートのフォント書式を継承
            if template_font_name:
                run.font.name = template_font_name
            if template_font_size:
                run.font.size = template_font_size
            if template_font_bold is not None:
                run.font.bold = template_font_bold
            if template_font_color:
                run.font.color.rgb = template_font_color
    else:
        if tf.paragraphs:
            p = tf.paragraphs[0]
            for run in p.runs:
                run.text = ""
            if p.runs:
                p.runs[0].text = text
            else:
                run = p.add_run()
                run.text = text
                if template_font_name:
                    run.font.name = template_font_name
                if template_font_size:
                    run.font.size = template_font_size
                if template_font_bold is not None:
                    run.font.bold = template_font_bold
                if template_font_color:
                    run.font.color.rgb = template_font_color
        else:
            p = tf.paragraphs[0]
            p.text = text


def inject_placeholder_text(slide, ph_map, slide_data):
    """PH idx でテキスト注入"""
    for field_name, ph_info in ph_map.items():
        idx = ph_info.get("idx")
        if idx is None:
            continue

        ph = _get_placeholder_by_idx(slide, idx)
        if ph is None:
            print(f"[engine] Warning: placeholder idx={idx} not found for field '{field_name}'", file=sys.stderr)
            continue

        if field_name in ("title", "subtitle"):
            text = slide_data.get(field_name, "")
            if text:
                set_text_preserving_format(ph, text)
        elif field_name == "key_message":
            km_text = slide_data.get("key_message", "")
            if km_text:
                set_text_preserving_format(ph, km_text)
        elif field_name == "body":
            bullets = slide_data.get("bullets")
            body_text = slide_data.get("body", "")
            if bullets:
                set_text_preserving_format(ph, None, bullets=bullets)
            elif body_text:
                set_text_preserving_format(ph, body_text)
        elif field_name in ("left_body", "right_body"):
            text = slide_data.get(field_name, "")
            if text:
                set_text_preserving_format(ph, text)
        elif field_name in ("left_subtitle", "right_subtitle"):
            text = slide_data.get(field_name, "")
            if text:
                set_text_preserving_format(ph, text)
        elif field_name == "extra":
            text = slide_data.get("extra", "")
            if text:
                set_text_preserving_format(ph, text)


def _add_fallback_content(slide, slide_data, w_emu, h_emu, theme):
    """リッチ要素未対応時のフォールバック: タイトル + 箇条書き"""
    from rich_elements import add_textbox, add_rect, hex_to_rgb
    from pptx.dml.color import RGBColor

    accent = theme.get("accent_color", "4472C4")
    body_font = theme.get("body_font", "Arial")
    heading_font = theme.get("heading_font", "Arial")

    # Title bar
    bar_h = int(0.10 * h_emu)
    add_rect(slide, 0, 0, w_emu, bar_h, fill_color=accent)
    title = slide_data.get("title", "")
    add_textbox(
        slide,
        int(0.03 * w_emu), int(0.02 * h_emu),
        int(0.94 * w_emu), bar_h,
        title,
        font_name=heading_font, font_size=24, bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.LEFT,
    )

    # Bullets or body
    bullets = slide_data.get("bullets", [])
    body = slide_data.get("body", "")
    content_text = "\n".join(f"• {b}" for b in bullets) if bullets else body

    if content_text:
        add_textbox(
            slide,
            int(0.06 * w_emu), int(0.16 * h_emu),
            int(0.88 * w_emu), int(0.75 * h_emu),
            content_text,
            font_name=body_font, font_size=16, bold=False,
            color=RGBColor(0x33, 0x33, 0x33),
            align=PP_ALIGN.LEFT,
        )


# ---------------------------------------------------------------------------
# Main pipeline
# ---------------------------------------------------------------------------

def build_presentation(outline, template_name, output_path, catalog_path=None):
    """
    メインパイプライン: アウトライン + テンプレート → PPTX 生成

    Args:
        outline: dict with 'title' and 'slides'
        template_name: "simple" or "standard"
        output_path: 出力 PPTX パス
        catalog_path: catalog.yaml パス (None=デフォルト)

    Returns:
        Path: 出力ファイルパス
    """
    catalog = load_catalog(catalog_path)
    tpl_config = resolve_template(catalog, template_name)
    tpl_path = tpl_config["_abs_path"]

    if not tpl_path.exists():
        raise FileNotFoundError(
            f"テンプレートファイルが見つかりません: {tpl_path}\n"
            f"テンプレートを配置してください:\n"
            f"  {tpl_path}"
        )

    # 1. Copy template to output
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(tpl_path, output)

    # 2. Open copied file
    prs = Presentation(str(output))

    # 3. Delete existing slides
    delete_all_slides(prs)

    # Slide dimensions in EMU
    w_emu = prs.slide_width
    h_emu = prs.slide_height
    theme = tpl_config.get("theme", {})
    layouts_config = tpl_config.get("layouts", {})

    slides = outline.get("slides", [])
    for slide_data in slides:
        slide_type = slide_data.get("slide_type", "content")
        layout_key = resolve_layout(catalog, template_name, slide_type)

        if layout_key == "code_gen":
            # Use blank layout + rich_elements
            blank_info = layouts_config.get("blank", {})
            blank_idx = blank_info.get("index", 6)
            try:
                layout = prs.slide_layouts[blank_idx]
            except IndexError:
                layout = prs.slide_layouts[-1]
            slide = prs.slides.add_slide(layout)

            # Try rich element rendering
            rendered = render_rich(slide_type, slide, slide_data, w_emu, h_emu, theme)
            if not rendered:
                # Fallback: title + bullets on blank
                _add_fallback_content(slide, slide_data, w_emu, h_emu, theme)
        else:
            # Use template layout + placeholder injection
            layout_info = layouts_config.get(layout_key)
            if layout_info is None:
                print(f"[engine] Warning: layout '{layout_key}' not in template config, using content", file=sys.stderr)
                layout_info = layouts_config.get("content", {"index": 0, "placeholders": {}})

            layout_idx = layout_info.get("index", 0)
            try:
                layout = prs.slide_layouts[layout_idx]
            except IndexError:
                print(f"[engine] Warning: layout index {layout_idx} out of range, using index 0", file=sys.stderr)
                layout = prs.slide_layouts[0]

            slide = prs.slides.add_slide(layout)
            ph_map = layout_info.get("placeholders", {})
            inject_placeholder_text(slide, ph_map, slide_data)

    # 4. Save
    prs.save(str(output))
    print(f"[engine] Generated: {output.resolve()} ({len(slides)} slides)")
    return output.resolve()
