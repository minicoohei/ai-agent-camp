"""
rich_elements.py — Blank レイアウト上にコード生成でリッチ要素を描画

全座標はスライドサイズに対する相対比率（0.0〜1.0）で計算し、
simple (13.333") / standard (20.0") 両方で動作する。

パターン:
  MVP:  kpi_dashboard, process_flow, comparison, table
  V2:   title_cover, section_header, content_bullets, closing_slide,
        numbered_agenda, two_column_content, key_message_hero
"""

from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Pattern registry
# ---------------------------------------------------------------------------
_PATTERNS = {}


def register_pattern(name):
    """Decorator to register a rich element pattern."""
    def decorator(func):
        _PATTERNS[name] = func
        return func
    return decorator


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def rel(fraction, total_emu):
    """相対座標(0.0〜1.0) → EMU 変換"""
    return int(fraction * total_emu)


def hex_to_rgb(hex_str):
    """'4472C4' → RGBColor"""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_textbox(slide, left, top, width, height, text,
                font_name="Arial", font_size=12, bold=False,
                color=None, align=PP_ALIGN.LEFT, word_wrap=True,
                anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """テキストボックスを追加"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    txbox.text_frame.word_wrap = word_wrap
    tf = txbox.text_frame
    tf.auto_size = None
    try:
        tf.anchor = anchor
    except (AttributeError, TypeError):
        pass
    try:
        tf.paragraphs[0].alignment = align
    except (AttributeError, IndexError):
        pass
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color if isinstance(color, RGBColor) else hex_to_rgb(color)
    try:
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    except Exception:
        pass
    if line_spacing is not None:
        try:
            tf.paragraphs[0].line_spacing = Pt(line_spacing)
        except Exception:
            pass
    return txbox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_name="Arial", font_size=12, bold=False,
                          color=None, align=PP_ALIGN.LEFT, word_wrap=True,
                          line_spacing=None, bullet_char=None):
    """複数行テキストボックス（箇条書き対応）"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    txbox.text_frame.word_wrap = word_wrap
    tf = txbox.text_frame
    tf.auto_size = None

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            try:
                p.line_spacing = Pt(line_spacing)
            except Exception:
                pass
        display_text = f"{bullet_char} {line}" if bullet_char else line
        run = p.add_run()
        run.text = display_text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color if isinstance(color, RGBColor) else hex_to_rgb(color)
    return txbox


def add_rounded_rect(slide, left, top, width, height,
                     fill_color=None, line_color=None, line_width=None):
    """角丸四角形を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = (
            fill_color if isinstance(fill_color, RGBColor) else hex_to_rgb(fill_color)
        )
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = (
            line_color if isinstance(line_color, RGBColor) else hex_to_rgb(line_color)
        )
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None):
    """四角形を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = (
            fill_color if isinstance(fill_color, RGBColor) else hex_to_rgb(fill_color)
        )
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = (
            line_color if isinstance(line_color, RGBColor) else hex_to_rgb(line_color)
        )
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, left, top, diameter, fill_color=None):
    """円を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, diameter, diameter
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = (
            fill_color if isinstance(fill_color, RGBColor) else hex_to_rgb(fill_color)
        )
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape


def add_arrow(slide, start_x, start_y, end_x, end_y, color=None, width=Pt(2)):
    """矢印コネクタを追加"""
    connector = slide.shapes.add_connector(
        1, start_x, start_y, end_x, end_y
    )
    if color:
        c = color if isinstance(color, RGBColor) else hex_to_rgb(color)
        connector.line.color.rgb = c
    connector.line.width = width
    ln = connector.line._ln
    tail = ln.makeelement(qn("a:tailEnd"), {})
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    ln.append(tail)
    return connector


def _resolve_fonts(theme):
    """テーマからフォント情報を取得"""
    return {
        "heading": theme.get("heading_font", "Arial"),
        "body": theme.get("body_font", "Arial"),
        "number": theme.get("number_font", theme.get("body_font", "Arial")),
        "accent": theme.get("accent_color", "4472C4"),
    }


def _title_bar(slide, title, w_emu, h_emu, theme):
    """スライド上部にタイトルバーを描画（共通）"""
    fonts = _resolve_fonts(theme)
    bar_h = rel(0.10, h_emu)
    add_rect(slide, 0, 0, w_emu, bar_h, fill_color=fonts["accent"])
    add_textbox(
        slide,
        rel(0.03, w_emu), rel(0.02, h_emu),
        rel(0.94, w_emu), bar_h,
        title,
        font_name=fonts["heading"], font_size=24, bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.LEFT,
    )
    return bar_h


def _accent_line(slide, w_emu, h_emu, y_frac, theme, width_frac=0.92):
    """アクセントラインを描画"""
    fonts = _resolve_fonts(theme)
    add_rect(
        slide,
        rel((1.0 - width_frac) / 2, w_emu), rel(y_frac, h_emu),
        rel(width_frac, w_emu), Pt(3),
        fill_color=fonts["accent"],
    )


# ===========================================================================
# V2 PATTERNS: テンプレート忠実度の高いパターン
# ===========================================================================

@register_pattern("title")
def _title_cover(slide, slide_data, w_emu, h_emu, theme):
    """表紙スライド: 中央タイトル + サブタイトル + アクセントライン"""
    title = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle", "")
    fonts = _resolve_fonts(theme)

    # 上部アクセントライン
    add_rect(
        slide, 0, 0, w_emu, rel(0.005, h_emu),
        fill_color=fonts["accent"],
    )

    # メインタイトル
    add_textbox(
        slide,
        rel(0.08, w_emu), rel(0.30, h_emu),
        rel(0.84, w_emu), rel(0.20, h_emu),
        title,
        font_name=fonts["heading"], font_size=40, bold=True,
        color=RGBColor(0x1A, 0x1A, 0x1A),
        align=PP_ALIGN.CENTER,
    )

    # アクセントライン（タイトル下）
    _accent_line(slide, w_emu, h_emu, 0.52, theme, width_frac=0.30)

    # サブタイトル
    if subtitle:
        add_textbox(
            slide,
            rel(0.10, w_emu), rel(0.56, h_emu),
            rel(0.80, w_emu), rel(0.12, h_emu),
            subtitle,
            font_name=fonts["body"], font_size=18, bold=False,
            color=RGBColor(0x66, 0x66, 0x66),
            align=PP_ALIGN.CENTER,
        )
    return True


@register_pattern("section")
def _section_header(slide, slide_data, w_emu, h_emu, theme):
    """セクション区切り: 大きなセクション番号 + タイトル"""
    title = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle", "")
    fonts = _resolve_fonts(theme)

    # 左側アクセントバー
    add_rect(
        slide,
        rel(0.04, w_emu), rel(0.20, h_emu),
        rel(0.005, w_emu), rel(0.60, h_emu),
        fill_color=fonts["accent"],
    )

    # セクション番号（slide_numberを使用）
    slide_num = slide_data.get("slide_number", "")
    if slide_num:
        add_textbox(
            slide,
            rel(0.06, w_emu), rel(0.25, h_emu),
            rel(0.15, w_emu), rel(0.20, h_emu),
            f"{slide_num:02d}" if isinstance(slide_num, int) else str(slide_num),
            font_name=fonts["number"], font_size=64, bold=True,
            color=hex_to_rgb(fonts["accent"]),
            align=PP_ALIGN.LEFT,
        )

    # セクションタイトル
    add_textbox(
        slide,
        rel(0.06, w_emu), rel(0.48, h_emu),
        rel(0.85, w_emu), rel(0.15, h_emu),
        title,
        font_name=fonts["heading"], font_size=36, bold=True,
        color=RGBColor(0x1A, 0x1A, 0x1A),
        align=PP_ALIGN.LEFT,
    )

    if subtitle:
        add_textbox(
            slide,
            rel(0.06, w_emu), rel(0.64, h_emu),
            rel(0.85, w_emu), rel(0.10, h_emu),
            subtitle,
            font_name=fonts["body"], font_size=16, bold=False,
            color=RGBColor(0x66, 0x66, 0x66),
            align=PP_ALIGN.LEFT,
        )
    return True


@register_pattern("content")
def _content_bullets(slide, slide_data, w_emu, h_emu, theme):
    """コンテンツスライド: タイトルバー + 箇条書き"""
    title = slide_data.get("title", "")
    bullets = slide_data.get("bullets", [])
    body = slide_data.get("body", "")
    fonts = _resolve_fonts(theme)

    _title_bar(slide, title, w_emu, h_emu, theme)

    if bullets:
        add_multiline_textbox(
            slide,
            rel(0.06, w_emu), rel(0.16, h_emu),
            rel(0.88, w_emu), rel(0.75, h_emu),
            bullets,
            font_name=fonts["body"], font_size=16, bold=False,
            color=RGBColor(0x33, 0x33, 0x33),
            align=PP_ALIGN.LEFT,
            line_spacing=28,
            bullet_char="●",
        )
    elif body:
        add_textbox(
            slide,
            rel(0.06, w_emu), rel(0.16, h_emu),
            rel(0.88, w_emu), rel(0.75, h_emu),
            body,
            font_name=fonts["body"], font_size=16, bold=False,
            color=RGBColor(0x33, 0x33, 0x33),
            align=PP_ALIGN.LEFT,
        )
    return True


@register_pattern("key_message")
def _key_message_hero(slide, slide_data, w_emu, h_emu, theme):
    """キーメッセージ: 大きな1つのメッセージを中央表示"""
    title = slide_data.get("title", "")
    body = slide_data.get("body", "")
    fonts = _resolve_fonts(theme)

    _title_bar(slide, title, w_emu, h_emu, theme)

    if body:
        # 中央に大きなメッセージ
        add_textbox(
            slide,
            rel(0.08, w_emu), rel(0.25, h_emu),
            rel(0.84, w_emu), rel(0.50, h_emu),
            body,
            font_name=fonts["heading"], font_size=28, bold=True,
            color=RGBColor(0x1A, 0x1A, 0x1A),
            align=PP_ALIGN.CENTER,
        )

        # 下部アクセントライン
        _accent_line(slide, w_emu, h_emu, 0.80, theme, width_frac=0.20)
    return True


@register_pattern("agenda")
def _numbered_agenda(slide, slide_data, w_emu, h_emu, theme):
    """アジェンダ: 番号付き項目リスト（Futura番号スタイル）"""
    title = slide_data.get("title", "アジェンダ")
    bullets = slide_data.get("bullets", [])
    fonts = _resolve_fonts(theme)

    _title_bar(slide, title, w_emu, h_emu, theme)

    n = min(len(bullets), 8)
    item_h_frac = min(0.08, 0.72 / max(n, 1))
    start_y = 0.16

    for i, item in enumerate(bullets[:n]):
        y = start_y + i * (item_h_frac + 0.01)

        # 番号円
        circle_d = rel(0.04, w_emu)
        add_circle(
            slide,
            rel(0.06, w_emu), rel(y, h_emu),
            circle_d,
            fill_color=fonts["accent"],
        )
        add_textbox(
            slide,
            rel(0.06, w_emu), rel(y + 0.002, h_emu),
            circle_d, circle_d,
            str(i + 1),
            font_name=fonts["number"], font_size=16, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
        )

        # 項目テキスト
        add_textbox(
            slide,
            rel(0.12, w_emu), rel(y + 0.005, h_emu),
            rel(0.80, w_emu), rel(item_h_frac, h_emu),
            item,
            font_name=fonts["body"], font_size=16, bold=False,
            color=RGBColor(0x33, 0x33, 0x33),
            align=PP_ALIGN.LEFT,
        )

        # 区切り線
        if i < n - 1:
            add_rect(
                slide,
                rel(0.06, w_emu), rel(y + item_h_frac + 0.005, h_emu),
                rel(0.88, w_emu), Pt(1),
                fill_color="E0E0E0",
            )
    return True


@register_pattern("closing")
def _closing_slide(slide, slide_data, w_emu, h_emu, theme):
    """クロージング: Thank you + 連絡先"""
    title = slide_data.get("title", "ご清聴ありがとうございました")
    subtitle = slide_data.get("subtitle", "")
    fonts = _resolve_fonts(theme)

    # 上部アクセントバー
    add_rect(slide, 0, 0, w_emu, rel(0.005, h_emu), fill_color=fonts["accent"])

    # メインメッセージ
    add_textbox(
        slide,
        rel(0.10, w_emu), rel(0.30, h_emu),
        rel(0.80, w_emu), rel(0.20, h_emu),
        title,
        font_name=fonts["heading"], font_size=36, bold=True,
        color=RGBColor(0x1A, 0x1A, 0x1A),
        align=PP_ALIGN.CENTER,
    )

    # アクセントライン
    _accent_line(slide, w_emu, h_emu, 0.52, theme, width_frac=0.20)

    # サブテキスト
    if subtitle:
        add_textbox(
            slide,
            rel(0.15, w_emu), rel(0.58, h_emu),
            rel(0.70, w_emu), rel(0.15, h_emu),
            subtitle,
            font_name=fonts["body"], font_size=16, bold=False,
            color=RGBColor(0x66, 0x66, 0x66),
            align=PP_ALIGN.CENTER,
        )
    return True


@register_pattern("two_column")
def _two_column_content(slide, slide_data, w_emu, h_emu, theme):
    """2カラムレイアウト: 左右にコンテンツ"""
    title = slide_data.get("title", "")
    left_subtitle = slide_data.get("left_subtitle", "")
    left_body = slide_data.get("left_body", "")
    right_subtitle = slide_data.get("right_subtitle", "")
    right_body = slide_data.get("right_body", "")
    fonts = _resolve_fonts(theme)

    _title_bar(slide, title, w_emu, h_emu, theme)

    col_w = 0.42
    left_x = 0.04
    right_x = 0.54
    subtitle_top = 0.14
    body_top = 0.20 if (left_subtitle or right_subtitle) else 0.15
    body_h = 0.73 if (left_subtitle or right_subtitle) else 0.78

    # 左サブタイトル
    if left_subtitle:
        add_textbox(
            slide,
            rel(left_x, w_emu), rel(subtitle_top, h_emu),
            rel(col_w, w_emu), rel(0.05, h_emu),
            left_subtitle,
            font_name=fonts["heading"], font_size=14, bold=True,
            color=hex_to_rgb(fonts["accent"]),
            align=PP_ALIGN.LEFT,
        )

    # 左カラム背景
    add_rounded_rect(
        slide,
        rel(left_x, w_emu), rel(body_top, h_emu),
        rel(col_w, w_emu), rel(body_h, h_emu),
        fill_color="F8F9FA", line_color="E0E0E0",
    )
    if left_body:
        add_textbox(
            slide,
            rel(left_x + 0.02, w_emu), rel(body_top + 0.03, h_emu),
            rel(col_w - 0.04, w_emu), rel(body_h - 0.06, h_emu),
            left_body,
            font_name=fonts["body"], font_size=14, bold=False,
            color=RGBColor(0x33, 0x33, 0x33),
            align=PP_ALIGN.LEFT,
        )

    # 右サブタイトル
    if right_subtitle:
        add_textbox(
            slide,
            rel(right_x, w_emu), rel(subtitle_top, h_emu),
            rel(col_w, w_emu), rel(0.05, h_emu),
            right_subtitle,
            font_name=fonts["heading"], font_size=14, bold=True,
            color=hex_to_rgb(fonts["accent"]),
            align=PP_ALIGN.LEFT,
        )

    # 右カラム背景
    add_rounded_rect(
        slide,
        rel(right_x, w_emu), rel(body_top, h_emu),
        rel(col_w, w_emu), rel(body_h, h_emu),
        fill_color="F8F9FA", line_color="E0E0E0",
    )
    if right_body:
        add_textbox(
            slide,
            rel(right_x + 0.02, w_emu), rel(body_top + 0.03, h_emu),
            rel(col_w - 0.04, w_emu), rel(body_h - 0.06, h_emu),
            right_body,
            font_name=fonts["body"], font_size=14, bold=False,
            color=RGBColor(0x33, 0x33, 0x33),
            align=PP_ALIGN.LEFT,
        )
    return True


# ===========================================================================
# MVP PATTERNS (improved with theme-aware fonts)
# ===========================================================================

@register_pattern("kpi_dashboard")
def _kpi_dashboard(slide, slide_data, w_emu, h_emu, theme):
    """3-4カード並列: 角丸背景 + 大数字 + ラベル + 変化率"""
    title = slide_data.get("title", "KPI Dashboard")
    kpis = slide_data.get("kpis", [])
    if not kpis:
        return False

    fonts = _resolve_fonts(theme)

    _title_bar(slide, title, w_emu, h_emu, theme)

    n = min(len(kpis), 4)
    margin_x = 0.04
    gap = 0.02
    card_w_frac = (1.0 - margin_x * 2 - gap * (n - 1)) / n
    card_h_frac = 0.55
    top_frac = 0.18

    for i, kpi in enumerate(kpis[:n]):
        cx = margin_x + i * (card_w_frac + gap)
        add_rounded_rect(
            slide,
            rel(cx, w_emu), rel(top_frac, h_emu),
            rel(card_w_frac, w_emu), rel(card_h_frac, h_emu),
            fill_color="F2F6FC", line_color="D0D8E8",
        )
        value_text = f"{kpi.get('value', '')}{kpi.get('unit', '')}"
        add_textbox(
            slide,
            rel(cx + 0.01, w_emu), rel(top_frac + 0.06, h_emu),
            rel(card_w_frac - 0.02, w_emu), rel(0.20, h_emu),
            value_text,
            font_name=fonts["number"], font_size=36, bold=True,
            color=hex_to_rgb(fonts["accent"]), align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            rel(cx + 0.01, w_emu), rel(top_frac + 0.28, h_emu),
            rel(card_w_frac - 0.02, w_emu), rel(0.10, h_emu),
            kpi.get("label", ""),
            font_name=fonts["body"], font_size=14, bold=False,
            color=RGBColor(0x33, 0x33, 0x33), align=PP_ALIGN.CENTER,
        )
        change = kpi.get("change", "")
        if change:
            is_positive = change.startswith("+")
            change_color = RGBColor(0x10, 0xB9, 0x81) if is_positive else RGBColor(0xEF, 0x44, 0x44)
            add_textbox(
                slide,
                rel(cx + 0.01, w_emu), rel(top_frac + 0.40, h_emu),
                rel(card_w_frac - 0.02, w_emu), rel(0.08, h_emu),
                change,
                font_name=fonts["body"], font_size=14, bold=True,
                color=change_color, align=PP_ALIGN.CENTER,
            )
    return True


@register_pattern("process_flow")
def _process_flow(slide, slide_data, w_emu, h_emu, theme):
    """番号付き円 + 矢印 + タイトル + 説明（3-5ステップ）"""
    title = slide_data.get("title", "Process Flow")
    steps = slide_data.get("steps", [])
    if not steps:
        return False

    fonts = _resolve_fonts(theme)

    _title_bar(slide, title, w_emu, h_emu, theme)

    n = min(len(steps), 5)
    margin_x = 0.06
    usable_w = 1.0 - margin_x * 2
    step_w = usable_w / n
    circle_d = 0.06
    circle_top = 0.22
    text_top = 0.38

    for i, step in enumerate(steps[:n]):
        # Handle both string and dict step formats
        if isinstance(step, str):
            step_title = step
            desc = ""
        else:
            step_title = step.get("title", f"Step {i+1}")
            desc = step.get("description", "")

        cx = margin_x + step_w * i + step_w / 2

        circle_left = cx - circle_d / 2
        d_emu = rel(circle_d, w_emu)
        add_circle(
            slide,
            rel(circle_left, w_emu), rel(circle_top, h_emu),
            d_emu, fill_color=fonts["accent"],
        )
        add_textbox(
            slide,
            rel(circle_left, w_emu), rel(circle_top + 0.005, h_emu),
            d_emu, d_emu,
            str(i + 1),
            font_name=fonts["number"], font_size=18, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
        )

        if i < n - 1:
            arrow_start_x = rel(cx + circle_d / 2 + 0.005, w_emu)
            next_cx = margin_x + step_w * (i + 1) + step_w / 2
            arrow_end_x = rel(next_cx - circle_d / 2 - 0.005, w_emu)
            arrow_y = rel(circle_top + circle_d / 2, h_emu)
            add_arrow(slide, arrow_start_x, arrow_y, arrow_end_x, arrow_y,
                      color=fonts["accent"])

        add_textbox(
            slide,
            rel(cx - step_w / 2 + 0.01, w_emu), rel(text_top, h_emu),
            rel(step_w - 0.02, w_emu), rel(0.08, h_emu),
            step_title,
            font_name=fonts["heading"], font_size=14, bold=True,
            color=RGBColor(0x33, 0x33, 0x33), align=PP_ALIGN.CENTER,
        )

        if desc:
            add_textbox(
                slide,
                rel(cx - step_w / 2 + 0.01, w_emu), rel(text_top + 0.10, h_emu),
                rel(step_w - 0.02, w_emu), rel(0.30, h_emu),
                desc,
                font_name=fonts["body"], font_size=11, bold=False,
                color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER,
            )
    return True


@register_pattern("comparison")
def _comparison(slide, slide_data, w_emu, h_emu, theme):
    """左右2ブロック + ヘッダーバー"""
    title = slide_data.get("title", "Comparison")
    left = slide_data.get("left", {})
    right = slide_data.get("right", {})

    # Support flat field format (left_subtitle/left_body/right_subtitle/right_body)
    if not left and (slide_data.get("left_subtitle") or slide_data.get("left_body")):
        left = {"title": slide_data.get("left_subtitle", ""), "items": [slide_data["left_body"]] if slide_data.get("left_body") else []}
    if not right and (slide_data.get("right_subtitle") or slide_data.get("right_body")):
        right = {"title": slide_data.get("right_subtitle", ""), "items": [slide_data["right_body"]] if slide_data.get("right_body") else []}

    if not left and not right:
        return False

    fonts = _resolve_fonts(theme)

    _title_bar(slide, title, w_emu, h_emu, theme)

    col_w = 0.42
    left_x = 0.04
    right_x = 0.54
    header_top = 0.14
    header_h = 0.08
    body_top = 0.24
    body_h = 0.65

    # Left header bar
    add_rect(
        slide,
        rel(left_x, w_emu), rel(header_top, h_emu),
        rel(col_w, w_emu), rel(header_h, h_emu),
        fill_color=fonts["accent"],
    )
    add_textbox(
        slide,
        rel(left_x + 0.01, w_emu), rel(header_top + 0.01, h_emu),
        rel(col_w - 0.02, w_emu), rel(header_h - 0.02, h_emu),
        left.get("title", "Option A"),
        font_name=fonts["heading"], font_size=16, bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
    )

    # Right header bar
    add_rect(
        slide,
        rel(right_x, w_emu), rel(header_top, h_emu),
        rel(col_w, w_emu), rel(header_h, h_emu),
        fill_color="6C757D",
    )
    add_textbox(
        slide,
        rel(right_x + 0.01, w_emu), rel(header_top + 0.01, h_emu),
        rel(col_w - 0.02, w_emu), rel(header_h - 0.02, h_emu),
        right.get("title", "Option B"),
        font_name=fonts["heading"], font_size=16, bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
    )

    # Left body items
    left_items = left.get("items", [])
    for j, item in enumerate(left_items[:8]):
        y = body_top + j * 0.07
        add_textbox(
            slide,
            rel(left_x + 0.02, w_emu), rel(y, h_emu),
            rel(col_w - 0.04, w_emu), rel(0.06, h_emu),
            f"● {item}",
            font_name=fonts["body"], font_size=12, bold=False,
            color=RGBColor(0x33, 0x33, 0x33), align=PP_ALIGN.LEFT,
        )

    # Right body items
    right_items = right.get("items", [])
    for j, item in enumerate(right_items[:8]):
        y = body_top + j * 0.07
        add_textbox(
            slide,
            rel(right_x + 0.02, w_emu), rel(y, h_emu),
            rel(col_w - 0.04, w_emu), rel(0.06, h_emu),
            f"● {item}",
            font_name=fonts["body"], font_size=12, bold=False,
            color=RGBColor(0x33, 0x33, 0x33), align=PP_ALIGN.LEFT,
        )
    return True


@register_pattern("table")
def _table(slide, slide_data, w_emu, h_emu, theme):
    """ヘッダー行（アクセント色）+ 交互背景色のテーブル"""
    title = slide_data.get("title", "Table")
    table_data = slide_data.get("table_data", {})
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    if not headers:
        return False

    fonts = _resolve_fonts(theme)

    _title_bar(slide, title, w_emu, h_emu, theme)

    margin_x = 0.06
    table_top = 0.16
    table_w = 1.0 - margin_x * 2
    n_cols = len(headers)
    n_rows = min(len(rows), 10) + 1
    row_h = min(0.07, (0.80 - table_top) / n_rows)

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        rel(margin_x, w_emu), rel(table_top, h_emu),
        rel(table_w, w_emu), rel(row_h * n_rows, h_emu),
    )
    tbl = tbl_shape.table

    for ci, header in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = str(header)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = fonts["heading"]
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn("a:solidFill"), {})
        srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": fonts["accent"]})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)

    for ri, row in enumerate(rows[:n_rows - 1]):
        for ci in range(n_cols):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(row[ci]) if ci < len(row) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = fonts["body"]
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            if ri % 2 == 1:
        
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn("a:solidFill"), {})
                srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": "F2F6FC"})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)
    return True


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def render(slide_type, slide, slide_data, width_emu, height_emu, theme):
    """
    リッチ要素を描画する。

    Args:
        slide_type: パターン名
        slide: python-pptx Slide object
        slide_data: dict with pattern-specific data
        width_emu: slide width in EMU
        height_emu: slide height in EMU
        theme: dict with heading_font, body_font, accent_color etc.

    Returns:
        True if rendered, False if pattern not found or data missing
    """
    pattern_fn = _PATTERNS.get(slide_type)
    if pattern_fn is None:
        return False
    try:
        return pattern_fn(slide, slide_data, width_emu, height_emu, theme)
    except Exception as e:
        import sys
        print(f"[rich_elements] Warning: failed to render '{slide_type}': {e}", file=sys.stderr)
        return False


def available_patterns():
    """登録済みパターン名一覧"""
    return list(_PATTERNS.keys())
