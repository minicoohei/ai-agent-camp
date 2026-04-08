#!/usr/bin/env python3
"""
PowerPoint Operations - PPTX ファイル操作ツール

python-pptx を使用して PowerPoint ファイルの読み取り・書き込み・分析を行います。

使用方法:
    uv run python tools/pptx_ops.py read <file.pptx>           # 読み取り
    uv run python tools/pptx_ops.py to-markdown <file.pptx>    # Markdown変換
    uv run python tools/pptx_ops.py analyze <file.pptx>        # 構造分析
    uv run python tools/pptx_ops.py extract-template <file.pptx> --output template.json  # テンプレート抽出
    uv run python tools/pptx_ops.py create <template.json> --output new.pptx  # 作成
"""

import argparse
import json
import os
import sys
from pathlib import Path
from typing import Optional, Dict, List, Any
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx library not found.")
    print("Please install with: pip install python-pptx")
    sys.exit(1)


class PowerPointOperator:
    """PowerPoint ファイル操作クラス"""
    
    def __init__(self, filepath: str = None):
        self.filepath = Path(filepath) if filepath else None
        self.presentation = None
    
    def load(self) -> bool:
        """プレゼンテーションを読み込み"""
        if not self.filepath or not self.filepath.exists():
            print(f"❌ File not found: {self.filepath}")
            return False
        
        try:
            self.presentation = Presentation(str(self.filepath))
            return True
        except Exception as e:
            print(f"❌ Failed to load presentation: {e}")
            return False
    
    def get_slide_count(self) -> int:
        """スライド数を取得"""
        if not self.presentation:
            return 0
        return len(self.presentation.slides)
    
    def read_slide(self, slide_index: int) -> Dict[str, Any]:
        """スライドの内容を読み取り"""
        if not self.presentation:
            return {"error": "Presentation not loaded"}
        
        if slide_index < 0 or slide_index >= len(self.presentation.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = self.presentation.slides[slide_index]
        
        slide_data = {
            "index": slide_index + 1,
            "slide_id": slide.slide_id,
            "shapes": [],
            "text_content": [],
            "notes": ""
        }
        
        # シェイプを取得
        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "text": ""
            }
            
            # テキストを抽出
            if shape.has_text_frame:
                texts = []
                for paragraph in shape.text_frame.paragraphs:
                    para_text = ""
                    for run in paragraph.runs:
                        para_text += run.text
                    if para_text.strip():
                        texts.append(para_text)
                
                if texts:
                    shape_info["text"] = "\n".join(texts)
                    slide_data["text_content"].extend(texts)
            
            # テーブルの場合
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text)
                    table_data.append(row_data)
                shape_info["table"] = table_data
            
            slide_data["shapes"].append(shape_info)
        
        # ノート（スピーカーノート）を取得
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            slide_data["notes"] = notes_text
        
        return slide_data
    
    def read_all_slides(self) -> List[Dict[str, Any]]:
        """全スライドの内容を読み取り"""
        if not self.presentation:
            return []
        
        slides = []
        for i in range(len(self.presentation.slides)):
            slides.append(self.read_slide(i))
        
        return slides
    
    def to_markdown(self) -> str:
        """プレゼンテーションをMarkdownに変換"""
        if not self.presentation:
            return "Error: Presentation not loaded"
        
        md_lines = []
        md_lines.append(f"# {self.filepath.name}")
        md_lines.append("")
        md_lines.append(f"**Slides**: {len(self.presentation.slides)}")
        md_lines.append("")
        md_lines.append("---")
        md_lines.append("")
        
        # 目次
        md_lines.append("## 目次")
        md_lines.append("")
        
        slides_data = self.read_all_slides()
        for slide in slides_data:
            # 最初のテキストをタイトルとして使用
            title = "Slide"
            if slide.get("text_content"):
                title = slide["text_content"][0][:50]
            md_lines.append(f"{slide['index']}. [{title}](#slide-{slide['index']})")
        
        md_lines.append("")
        
        # 各スライド
        for slide in slides_data:
            md_lines.append("---")
            md_lines.append("")
            md_lines.append(f"## Slide {slide['index']} {{#slide-{slide['index']}}}")
            md_lines.append("")
            
            # テキスト内容
            for i, text in enumerate(slide.get("text_content", [])):
                if i == 0:
                    md_lines.append(f"### {text}")
                else:
                    md_lines.append(text)
                md_lines.append("")
            
            # テーブル
            for shape in slide.get("shapes", []):
                if "table" in shape:
                    table = shape["table"]
                    if table:
                        # ヘッダー
                        md_lines.append("| " + " | ".join(table[0]) + " |")
                        md_lines.append("| " + " | ".join(["---"] * len(table[0])) + " |")
                        # データ
                        for row in table[1:]:
                            md_lines.append("| " + " | ".join(row) + " |")
                        md_lines.append("")
            
            # ノート
            if slide.get("notes"):
                md_lines.append("> **Speaker Notes:**")
                for line in slide["notes"].split("\n"):
                    md_lines.append(f"> {line}")
                md_lines.append("")
        
        return "\n".join(md_lines)
    
    def analyze(self) -> Dict[str, Any]:
        """プレゼンテーションの構造を分析"""
        if not self.presentation:
            return {"error": "Presentation not loaded"}
        
        report = {
            "filename": self.filepath.name,
            "slide_count": len(self.presentation.slides),
            "slide_width": self.presentation.slide_width,
            "slide_height": self.presentation.slide_height,
            "slides": [],
            "layouts_used": set(),
            "shape_types": {},
            "total_text_length": 0
        }
        
        for i, slide in enumerate(self.presentation.slides):
            slide_info = {
                "index": i + 1,
                "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
                "shape_count": len(slide.shapes),
                "has_notes": slide.has_notes_slide,
                "placeholders": []
            }
            
            report["layouts_used"].add(slide_info["layout_name"])
            
            for shape in slide.shapes:
                # シェイプタイプをカウント
                shape_type = str(shape.shape_type)
                report["shape_types"][shape_type] = report["shape_types"].get(shape_type, 0) + 1
                
                # テキスト長を累計
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            report["total_text_length"] += len(run.text)
                
                # プレースホルダー情報
                if shape.is_placeholder:
                    slide_info["placeholders"].append({
                        "name": shape.name,
                        "placeholder_format": str(shape.placeholder_format.type) if shape.placeholder_format else "N/A"
                    })
            
            report["slides"].append(slide_info)
        
        report["layouts_used"] = list(report["layouts_used"])
        
        return report
    
    def extract_template(self) -> Dict[str, Any]:
        """テンプレート情報を抽出"""
        if not self.presentation:
            return {"error": "Presentation not loaded"}
        
        template = {
            "source_file": self.filepath.name,
            "slide_width": self.presentation.slide_width,
            "slide_height": self.presentation.slide_height,
            "layouts": [],
            "slides": []
        }
        
        # レイアウト情報
        for layout in self.presentation.slide_master.slide_layouts:
            layout_info = {
                "name": layout.name,
                "placeholders": []
            }
            for shape in layout.placeholders:
                layout_info["placeholders"].append({
                    "idx": shape.placeholder_format.idx,
                    "type": str(shape.placeholder_format.type),
                    "name": shape.name
                })
            template["layouts"].append(layout_info)
        
        # スライド構造
        for i, slide in enumerate(self.presentation.slides):
            slide_template = {
                "index": i + 1,
                "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
                "content_structure": []
            }
            
            for shape in slide.shapes:
                shape_struct = {
                    "name": shape.name,
                    "type": str(shape.shape_type),
                    "is_placeholder": shape.is_placeholder,
                    "has_text": shape.has_text_frame,
                    "has_table": shape.has_table
                }
                
                if shape.has_text_frame:
                    # テキストの概要（最初の50文字）
                    text = ""
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            text += run.text
                    shape_struct["text_preview"] = text[:50] + "..." if len(text) > 50 else text
                
                slide_template["content_structure"].append(shape_struct)
            
            template["slides"].append(slide_template)
        
        return template
    
    def create_from_template(self, template_data: Dict, output_path: str) -> str:
        """テンプレートデータから新しいプレゼンテーションを作成"""
        prs = Presentation()
        
        # スライドサイズを設定（オプション）
        if "slide_width" in template_data:
            prs.slide_width = template_data["slide_width"]
        if "slide_height" in template_data:
            prs.slide_height = template_data["slide_height"]
        
        # レイアウトを取得
        blank_layout = prs.slide_layouts[6]  # 通常は空白レイアウト
        
        for slide_data in template_data.get("slides", []):
            slide = prs.slides.add_slide(blank_layout)
            
            # タイトル
            if "title" in slide_data:
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5), Inches(9), Inches(1)
                )
                title_frame = title_box.text_frame
                title_frame.text = slide_data["title"]
                title_frame.paragraphs[0].font.size = Pt(36)
                title_frame.paragraphs[0].font.bold = True
            
            # コンテンツ
            if "content" in slide_data:
                content_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.5), Inches(9), Inches(5)
                )
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                for item in slide_data["content"]:
                    p = content_frame.add_paragraph()
                    p.text = item
                    p.font.size = Pt(18)
            
            # ノート
            if "notes" in slide_data:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_data["notes"]
        
        prs.save(output_path)
        return output_path


def cmd_read(args):
    """読み取りコマンド"""
    op = PowerPointOperator(args.file)
    if not op.load():
        return
    
    if args.slide is not None:
        data = op.read_slide(args.slide - 1)  # 1-indexed to 0-indexed
        slides_data = [data]
    else:
        slides_data = op.read_all_slides()
    
    if args.format == "json":
        print(json.dumps(slides_data, ensure_ascii=False, indent=2, default=str))
    else:
        for slide in slides_data:
            print(f"\n=== Slide {slide['index']} ===")
            print(f"Shapes: {len(slide['shapes'])}")
            print("\nText content:")
            for text in slide.get("text_content", []):
                print(f"  - {text[:100]}{'...' if len(text) > 100 else ''}")
            if slide.get("notes"):
                print(f"\nNotes: {slide['notes'][:100]}...")


def cmd_to_markdown(args):
    """Markdown変換コマンド"""
    op = PowerPointOperator(args.file)
    if not op.load():
        return
    
    md = op.to_markdown()
    
    if args.output:
        Path(args.output).write_text(md, encoding="utf-8")
        print(f"✅ Saved to: {args.output}")
    else:
        print(md)


def cmd_analyze(args):
    """分析コマンド"""
    op = PowerPointOperator(args.file)
    if not op.load():
        return
    
    report = op.analyze()
    
    if args.format == "json":
        print(json.dumps(report, ensure_ascii=False, indent=2, default=str))
    else:
        print(f"📊 Analysis Report: {report['filename']}")
        print("=" * 50)
        print(f"Slides: {report['slide_count']}")
        print(f"Total text length: {report['total_text_length']} characters")
        print(f"Layouts used: {report['layouts_used']}")
        print()
        print("Shape types:")
        for shape_type, count in sorted(report["shape_types"].items(), key=lambda x: -x[1]):
            print(f"  {shape_type}: {count}")
        print()
        print("Slides overview:")
        for slide in report["slides"]:
            print(f"  {slide['index']}. {slide['layout_name']} ({slide['shape_count']} shapes)")


def cmd_extract_template(args):
    """テンプレート抽出コマンド"""
    op = PowerPointOperator(args.file)
    if not op.load():
        return
    
    template = op.extract_template()
    
    output = json.dumps(template, ensure_ascii=False, indent=2, default=str)
    
    if args.output:
        Path(args.output).write_text(output, encoding="utf-8")
        print(f"✅ Template saved to: {args.output}")
    else:
        print(output)


def cmd_create(args):
    """作成コマンド"""
    try:
        template_data = json.loads(Path(args.template).read_text(encoding="utf-8"))
    except Exception as e:
        print(f"❌ Failed to load template: {e}")
        return
    
    op = PowerPointOperator()
    output = op.create_from_template(template_data, args.output)
    print(f"✅ Created: {output}")


def main():
    parser = argparse.ArgumentParser(description="PowerPoint Operations Tool")
    subparsers = parser.add_subparsers(dest="command", help="Commands")
    
    # read コマンド
    read_parser = subparsers.add_parser("read", help="PPTXファイルを読み取り")
    read_parser.add_argument("file", help="PPTXファイルパス")
    read_parser.add_argument("--slide", "-s", type=int, help="特定スライド番号（1-indexed）")
    read_parser.add_argument("--format", "-f", choices=["text", "json"], default="text")
    
    # to-markdown コマンド
    md_parser = subparsers.add_parser("to-markdown", help="Markdownに変換")
    md_parser.add_argument("file", help="PPTXファイルパス")
    md_parser.add_argument("--output", "-o", help="出力ファイルパス")
    
    # analyze コマンド
    analyze_parser = subparsers.add_parser("analyze", help="構造を分析")
    analyze_parser.add_argument("file", help="PPTXファイルパス")
    analyze_parser.add_argument("--format", "-f", choices=["text", "json"], default="text")
    
    # extract-template コマンド
    template_parser = subparsers.add_parser("extract-template", help="テンプレート抽出")
    template_parser.add_argument("file", help="PPTXファイルパス")
    template_parser.add_argument("--output", "-o", help="出力JSONファイルパス")
    
    # create コマンド
    create_parser = subparsers.add_parser("create", help="テンプレートから作成")
    create_parser.add_argument("template", help="テンプレートJSONファイル")
    create_parser.add_argument("--output", "-o", required=True, help="出力PPTXファイルパス")
    
    args = parser.parse_args()
    
    if not args.command:
        parser.print_help()
        sys.exit(1)
    
    if args.command == "read":
        cmd_read(args)
    elif args.command == "to-markdown":
        cmd_to_markdown(args)
    elif args.command == "analyze":
        cmd_analyze(args)
    elif args.command == "extract-template":
        cmd_extract_template(args)
    elif args.command == "create":
        cmd_create(args)


if __name__ == "__main__":
    main()
