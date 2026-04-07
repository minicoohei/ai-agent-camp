#!/usr/bin/env python3
"""
PPTX操作スクリプト（Final Example）

PowerPointファイルの解析、生成、編集を行います。

必要条件:
- Python 3.9以上
- python-pptx

使用方法:
    python pptx_operations.py analyze --input presentation.pptx
    python pptx_operations.py generate --output new_presentation.pptx
    python pptx_operations.py extract-text --input presentation.pptx
"""

import os
import sys
import argparse
import json
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, List, Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx がインストールされていません")
    print("  インストール: pip install python-pptx")


class PPTXAnalyzer:
    """PPTX解析クラス"""
    
    def __init__(self, pptx_path: str):
        self.pptx_path = pptx_path
        self.prs = None
        
        if HAS_PPTX and os.path.exists(pptx_path):
            self.prs = Presentation(pptx_path)
    
    def get_basic_info(self) -> Dict[str, Any]:
        """基本情報を取得"""
        if not self.prs:
            return self._get_mock_info()
        
        return {
            "file_path": self.pptx_path,
            "slide_count": len(self.prs.slides),
            "slide_width": self.prs.slide_width,
            "slide_height": self.prs.slide_height,
            "slide_width_inches": round(self.prs.slide_width / 914400, 2),
            "slide_height_inches": round(self.prs.slide_height / 914400, 2)
        }
    
    def analyze_slides(self) -> List[Dict[str, Any]]:
        """全スライドを解析"""
        if not self.prs:
            return self._get_mock_slides()
        
        slides_info = []
        
        for i, slide in enumerate(self.prs.slides, 1):
            slide_info = {
                "number": i,
                "layout_name": slide.slide_layout.name,
                "shapes_count": len(slide.shapes),
                "shapes": []
            }
            
            for shape in slide.shapes:
                shape_info = {
                    "name": shape.name,
                    "type": self._get_shape_type(shape),
                    "position": {
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    }
                }
                
                # テキスト
                if shape.has_text_frame:
                    shape_info["has_text"] = True
                    shape_info["text_preview"] = shape.text[:100] if shape.text else ""
                
                # テーブル
                if shape.has_table:
                    shape_info["is_table"] = True
                    shape_info["table_size"] = {
                        "rows": len(shape.table.rows),
                        "columns": len(shape.table.columns)
                    }
                
                slide_info["shapes"].append(shape_info)
            
            slides_info.append(slide_info)
        
        return slides_info
    
    def _get_shape_type(self, shape) -> str:
        """シェイプタイプを文字列で取得"""
        try:
            if hasattr(shape, 'shape_type'):
                return str(shape.shape_type).split('.')[-1]
        except:
            pass
        return "UNKNOWN"
    
    def extract_all_text(self) -> List[Dict[str, Any]]:
        """全テキストを抽出"""
        if not self.prs:
            return self._get_mock_text()
        
        texts = []
        
        for i, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            texts.append({
                                "slide": i,
                                "shape": shape.name,
                                "level": paragraph.level,
                                "text": paragraph.text
                            })
        
        return texts
    
    def get_layouts(self) -> List[Dict[str, Any]]:
        """利用可能なレイアウトを取得"""
        if not self.prs:
            return self._get_mock_layouts()
        
        layouts = []
        
        for master in self.prs.slide_masters:
            for layout in master.slide_layouts:
                layout_info = {
                    "name": layout.name,
                    "placeholders": []
                }
                
                for placeholder in layout.placeholders:
                    layout_info["placeholders"].append({
                        "idx": placeholder.placeholder_format.idx,
                        "type": str(placeholder.placeholder_format.type),
                        "name": placeholder.name
                    })
                
                layouts.append(layout_info)
        
        return layouts
    
    def _get_mock_info(self) -> Dict[str, Any]:
        """モック基本情報"""
        return {
            "file_path": self.pptx_path,
            "slide_count": 5,
            "slide_width": 9144000,
            "slide_height": 5143500,
            "slide_width_inches": 10.0,
            "slide_height_inches": 5.63,
            "note": "モックデータ（ファイルが見つかりません）"
        }
    
    def _get_mock_slides(self) -> List[Dict[str, Any]]:
        """モックスライド情報"""
        return [
            {
                "number": 1,
                "layout_name": "Title Slide",
                "shapes_count": 2,
                "shapes": [
                    {"name": "Title 1", "type": "PLACEHOLDER", "has_text": True, "text_preview": "プレゼンテーションタイトル"},
                    {"name": "Subtitle 2", "type": "PLACEHOLDER", "has_text": True, "text_preview": "サブタイトル"}
                ]
            },
            {
                "number": 2,
                "layout_name": "Title and Content",
                "shapes_count": 3,
                "shapes": [
                    {"name": "Title 1", "type": "PLACEHOLDER", "has_text": True, "text_preview": "アジェンダ"},
                    {"name": "Content 2", "type": "PLACEHOLDER", "has_text": True, "text_preview": "項目1\n項目2\n項目3"}
                ]
            }
        ]
    
    def _get_mock_text(self) -> List[Dict[str, Any]]:
        """モックテキスト"""
        return [
            {"slide": 1, "shape": "Title 1", "level": 0, "text": "プレゼンテーションタイトル"},
            {"slide": 1, "shape": "Subtitle 2", "level": 0, "text": "サブタイトル"},
            {"slide": 2, "shape": "Title 1", "level": 0, "text": "アジェンダ"},
            {"slide": 2, "shape": "Content 2", "level": 0, "text": "項目1"},
            {"slide": 2, "shape": "Content 2", "level": 0, "text": "項目2"}
        ]
    
    def _get_mock_layouts(self) -> List[Dict[str, Any]]:
        """モックレイアウト"""
        return [
            {"name": "Title Slide", "placeholders": [{"idx": 0, "type": "TITLE", "name": "Title"}]},
            {"name": "Title and Content", "placeholders": [{"idx": 0, "type": "TITLE", "name": "Title"}, {"idx": 1, "type": "BODY", "name": "Content"}]},
            {"name": "Section Header", "placeholders": [{"idx": 0, "type": "TITLE", "name": "Title"}]},
            {"name": "Two Content", "placeholders": []},
            {"name": "Blank", "placeholders": []}
        ]


class PPTXGenerator:
    """PPTX生成クラス"""
    
    def __init__(self, template_path: str = None):
        self.prs = None
        
        if HAS_PPTX:
            if template_path and os.path.exists(template_path):
                self.prs = Presentation(template_path)
            else:
                self.prs = Presentation()
    
    def create_presentation(self, content: Dict[str, Any]) -> bool:
        """プレゼンテーションを生成"""
        if not self.prs or not HAS_PPTX:
            print("python-pptxが利用できないため、生成をスキップします")
            return False
        
        # タイトルスライド
        title_slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title_slide.shapes.title.text = content.get("title", "タイトル")
        
        if title_slide.placeholders[1]:
            title_slide.placeholders[1].text = content.get("subtitle", "")
        
        # コンテンツスライド
        content_layout = self.prs.slide_layouts[1]
        
        for section in content.get("sections", []):
            slide = self.prs.slides.add_slide(content_layout)
            slide.shapes.title.text = section.get("title", "")
            
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            
            for i, point in enumerate(section.get("points", [])):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = point
                p.level = 0
        
        return True
    
    def save(self, output_path: str):
        """保存"""
        if self.prs:
            self.prs.save(output_path)
            print(f"✅ 保存完了: {output_path}")
        else:
            print("❌ プレゼンテーションが作成されていません")


# サンプルコンテンツ
SAMPLE_CONTENT = {
    "title": "AIエージェント活用研修",
    "subtitle": "2025年版 - 非エンジニア向け実践ガイド",
    "sections": [
        {
            "title": "研修の目的",
            "points": [
                "AIエージェントの基本概念を理解する",
                "実務での活用方法を習得する",
                "自分で設定・カスタマイズできるようになる"
            ]
        },
        {
            "title": "カリキュラム概要",
            "points": [
                "Foundation: LLM基礎、トークン、エージェント概念",
                "Module 1-5: 画像生成、図表、データ分析",
                "Module 6-11: Slack連携、動画、GAS、エージェント開発"
            ]
        },
        {
            "title": "学習の進め方",
            "points": [
                "/start-X-Y コマンドで各レッスンを開始",
                "ハンズオン形式で実際に手を動かす",
                "成果物をexamples/final/と比較して確認"
            ]
        },
        {
            "title": "まとめ",
            "points": [
                "AIは道具 - 使いこなすことで業務効率化",
                "試行錯誤が大切 - 失敗を恐れずに",
                "継続学習 - 技術は日々進化"
            ]
        }
    ]
}


def main():
    parser = argparse.ArgumentParser(description="PPTX操作スクリプト")
    subparsers = parser.add_subparsers(dest="command", help="コマンド")
    
    # analyze
    analyze_parser = subparsers.add_parser("analyze", help="PPTXを解析")
    analyze_parser.add_argument("--input", "-i", required=True, help="入力ファイル")
    analyze_parser.add_argument("--output", "-o", help="出力JSONファイル")
    
    # generate
    generate_parser = subparsers.add_parser("generate", help="PPTXを生成")
    generate_parser.add_argument("--output", "-o", required=True, help="出力ファイル")
    generate_parser.add_argument("--template", "-t", help="テンプレートファイル")
    generate_parser.add_argument("--content", "-c", help="コンテンツJSONファイル")
    
    # extract-text
    extract_parser = subparsers.add_parser("extract-text", help="テキスト抽出")
    extract_parser.add_argument("--input", "-i", required=True, help="入力ファイル")
    extract_parser.add_argument("--output", "-o", help="出力ファイル")
    
    # layouts
    layouts_parser = subparsers.add_parser("layouts", help="レイアウト一覧")
    layouts_parser.add_argument("--input", "-i", required=True, help="入力ファイル")
    
    args = parser.parse_args()
    
    if args.command == "analyze":
        analyzer = PPTXAnalyzer(args.input)
        
        result = {
            "info": analyzer.get_basic_info(),
            "slides": analyzer.analyze_slides(),
            "layouts": analyzer.get_layouts()
        }
        
        if args.output:
            with open(args.output, 'w', encoding='utf-8') as f:
                json.dump(result, f, ensure_ascii=False, indent=2)
            print(f"✅ 解析結果を保存: {args.output}")
        else:
            print(json.dumps(result, ensure_ascii=False, indent=2))
    
    elif args.command == "generate":
        # コンテンツ読み込み
        if args.content and os.path.exists(args.content):
            with open(args.content, 'r', encoding='utf-8') as f:
                content = json.load(f)
        else:
            content = SAMPLE_CONTENT
            print("サンプルコンテンツを使用します")
        
        generator = PPTXGenerator(args.template)
        if generator.create_presentation(content):
            generator.save(args.output)
    
    elif args.command == "extract-text":
        analyzer = PPTXAnalyzer(args.input)
        texts = analyzer.extract_all_text()
        
        if args.output:
            with open(args.output, 'w', encoding='utf-8') as f:
                for t in texts:
                    f.write(f"[Slide {t['slide']}] {t['text']}\n")
            print(f"✅ テキストを保存: {args.output}")
        else:
            for t in texts:
                print(f"[Slide {t['slide']}] {t['text']}")
    
    elif args.command == "layouts":
        analyzer = PPTXAnalyzer(args.input)
        layouts = analyzer.get_layouts()
        
        print("利用可能なレイアウト:")
        for layout in layouts:
            print(f"\n  {layout['name']}")
            for ph in layout.get('placeholders', []):
                print(f"    - {ph['name']} (idx: {ph['idx']}, type: {ph['type']})")
    
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
